# -*- coding: utf-8 -*-
"""
ارائه جلسه توسعه سازمانی + منتورینگ گروهی

بخش اول (اسلاید ۱ تا ۱۴): چرایی توسعه فردی و سازمانی، امواج تمدنی،
مکاتب مدیریتی منطبق بر هر موج، جایگاه فعلی سازمان نفت و گاز و مسیر گذار.

بخش دوم (اسلاید ۱۵ تا ۳۳): محتوای منتورینگ گروهی که از ارائه
group-mentoring-event بازاستفاده می‌شود.

اجرا:  python3 build.py
"""
import argparse
import importlib.util
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
LIB = os.path.join(HERE, "..", "_lib")
GM = os.path.join(HERE, "..", "group-mentoring-event")
sys.path.insert(0, LIB)

import theme as T
from theme import (
    NAVY, NAVY_DEEP, BLUE_MID, BLUE_LIGHT, BLUE_PALE, ORANGE, ORANGE_DEEP,
    ORANGE_PALE, WHITE, INK, INK_SOFT, PAPER, GREEN, GREEN_PALE, RED, RED_PALE,
    PURPLE, AMBER, LINE, SW, SH, M, CONTENT_W,
    fa, rect, card, gradient, hline, textbox, write, picture_fill,
    round_picture, base_slide, blank, footer, notes, no_line, shadow_off,
    scrim, scrim_gradient, full_bleed,
)

IMG = os.path.join(HERE, "img")

INTRO_N = 14          # تعداد اسلایدهای بخش مقدمه
TOTAL = 33            # کل اسلایدهای ارائه ترکیبی
LABEL = "توسعه سازمانی"

# ═══════════════════════ وصله‌کردن پاورقی برای شماره‌گذاری پیوسته ═══════
_orig_footer = T.footer


def _shifted_footer(slide, n, total, label=None):
    """اسلایدهای بخش منتورینگ با ۱۰ واحد جابه‌جایی شماره می‌خورند."""
    return _orig_footer(slide, n + INTRO_N - 1, TOTAL,
                        label="منتورینگ گروهی" if label is None else label)


def img(name):
    p = os.path.join(IMG, name)
    return p, Image.open(p).size


# ═════════════════════════════════════════ کمک‌کارهای بخش مقدمه ═════════
def head(slide, title, kicker=None, on_dark=False):
    top = Inches(0.40)
    if kicker:
        kb = textbox(slide, M, Inches(0.34), CONTENT_W, Inches(0.34))
        write(kb.text_frame, kicker, size=15.5, bold=True, color=ORANGE,
              first=True, line=1.0)
        top = Inches(0.68)
    tb = textbox(slide, M, top, CONTENT_W, Inches(0.70))
    write(tb.text_frame, title, size=31, bold=True,
          color=WHITE if on_dark else NAVY, first=True, line=1.05)
    hline(slide, SW - M - Inches(1.2), Inches(1.36), Inches(1.2), ORANGE,
          Pt(4.5))
    return slide


def divider(prs, num, title, subtitle, image):
    s = blank(prs)
    p, sz = img(image)
    full_bleed(s, p, sz)
    scrim_gradient(s, 0, 0, SW, SH, angle=0, a_from=0.10, a_to=0.94)
    scrim(s, 0, 0, SW, SH, alpha=0.22)

    tw = Inches(6.3)
    tx = SW - M - tw
    nb = textbox(s, tx, Inches(2.18), tw, Inches(0.62))
    write(nb.text_frame, num, size=17, bold=True, color=ORANGE, first=True,
          line=1.0)
    hline(s, tx + tw - Inches(1.3), Inches(2.86), Inches(1.3), ORANGE, Pt(4.5))
    tb = textbox(s, tx, Inches(3.12), tw, Inches(1.9))
    tf = tb.text_frame
    for i, ln in enumerate(title.split("\n")):
        write(tf, ln, size=44, bold=True, color=WHITE, first=(i == 0),
              line=1.18)
    write(tf, subtitle, size=19, color=BLUE_PALE, space_before=14, line=1.3)
    return s


def band(slide, y, chunks, h=Inches(0.94), fill=NAVY, accent=ORANGE,
         size=21, align=PP_ALIGN.CENTER, x=M, w=None):
    w = w or CONTENT_W
    rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.10)
    rect(slide, x + w - Inches(0.09), y, Inches(0.09), h, fill=accent)
    tb = textbox(slide, x + Inches(0.34), y, w - Inches(0.68), h,
                 anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame, chunks, size=size, color=WHITE, first=True,
          align=align, line=1.2)
    return slide


def photo_card(slide, x, y, w, h, image, caption=None, cap_h=Inches(0.62)):
    p, sz = img(image)
    pic = picture_fill(slide, p, x, y, w, h, sz)
    round_picture(pic)
    if caption:
        cy = y + h - cap_h
        scrim(slide, x, cy, w, cap_h, alpha=0.72)
        tb = textbox(slide, x + Inches(0.20), cy, w - Inches(0.40), cap_h,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, caption, size=14, bold=True, color=WHITE,
              first=True, align=PP_ALIGN.CENTER, line=1.15)
    return pic


# ═══════════════════════════════════════════════ ۱ — جلد ════════════════
def i01(prs, cfg):
    s = blank(prs)
    p, sz = img("cover-wave.jpg")
    full_bleed(s, p, sz)
    scrim_gradient(s, 0, 0, SW, SH, angle=0, a_from=0.06, a_to=0.95)
    scrim(s, 0, 0, SW, SH, alpha=0.20)

    tw = Inches(6.9)
    tx = SW - M - tw
    hline(s, tx + tw - Inches(1.4), Inches(1.62), Inches(1.4), ORANGE, Pt(5))

    tb = textbox(s, tx, Inches(1.92), tw, Inches(3.5))
    tf = tb.text_frame
    write(tf, "جلسه توسعه فردی و سازمانی", size=18, bold=True,
          color=BLUE_LIGHT, first=True, line=1.0)
    write(tf, "موج بعدی را", size=45, bold=True, color=WHITE, space_before=14,
          line=1.22)
    write(tf, [("چه کسی ", True, WHITE, 45), ("می‌سازد؟", True, ORANGE, 45)],
          line=1.22)
    write(tf, "از موج دوم تا سازمان یادگیرنده — و نقطه شروع ما",
          size=19, color=BLUE_PALE, space_before=18, line=1.3)

    hline(s, tx + tw - Inches(3.2), Inches(5.52), Inches(3.2), BLUE_LIGHT,
          Pt(1.25))
    mb = textbox(s, tx, Inches(5.72), tw, Inches(0.62))
    write(mb.text_frame,
          [("تهیه و تنظیم:  ", False, BLUE_LIGHT, 15.5),
           (cfg.presenter, True, WHITE, 19)], first=True, line=1.15)
    db = textbox(s, tx, Inches(6.30), tw, Inches(0.40))
    write(db.text_frame,
          [("تاریخ:  ", False, BLUE_LIGHT, 14), (cfg.date, True, WHITE, 14),
           ("        ", False, BLUE_LIGHT, 14),
           (cfg.org, False, BLUE_LIGHT, 14)], first=True, line=1.15)

    site_w = Inches(2.62)
    rect(s, tx + tw - site_w, Inches(6.76), site_w, Inches(0.46), fill=ORANGE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.22)
    sb = textbox(s, tx + tw - site_w, Inches(6.76), site_w, Inches(0.46),
                 anchor=MSO_ANCHOR.MIDDLE)
    write(sb.text_frame, "www.coachroom.ir", size=15, bold=True, color=WHITE,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    notes(s, """
    جمله شروع: «امروز دو کار می‌کنیم. اول می‌فهمیم سازمان ما کجای تاریخ
    ایستاده، بعد یک ابزار مشخص برای حرکت به جلو تمرین می‌کنیم.»
    """)
    return s


# ═══════════════════════════════════ ۲ — نقشه جلسه ══════════════════════
def i02(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "مسیر امروز", kicker="دو بخش، یک هدف")

    parts = [
        ("بخش اول", "چرا باید تغییر کنیم؟", [
            "پنج موج تمدنی و مکتب مدیریتی هر موج",
            "جدول انطباق موج‌ها با سبک‌های مدیریت",
            "سازمان ما دقیقاً کجای این نقشه است",
            "مسیر گذار به موج بالاتر",
        ], NAVY, "۳۰ دقیقه"),
        ("بخش دوم", "از کجا شروع کنیم؟", [
            "منتورینگ گروهی: چیستی، چرایی، نحوه اجرا",
            "قواعد فضای امن و مهارت پرسیدن",
            "یک دور تمرین واقعی، همین امروز",
        ], ORANGE_DEEP, "۹۰ دقیقه"),
    ]
    gap = Inches(0.42)
    cw = (CONTENT_W - gap) / 2
    y = Inches(1.72)
    ch = Inches(3.52)
    for i, (tag, title, items, color, dur) in enumerate(parts):
        x = SW - M - cw - i * (cw + gap)
        card(s, x, y, cw, ch, fill=WHITE, line=LINE)
        rect(s, x, y, cw, Inches(0.86), fill=color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
        rect(s, x, y + Inches(0.70), cw, Inches(0.16), fill=color)
        tb = textbox(s, x + Inches(0.30), y + Inches(0.06), cw - Inches(0.60),
                     Inches(0.74), anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame,
              [(tag + "  ·  ", False, WHITE, 14), (title, True, WHITE, 21)],
              first=True, line=1.05)
        yy = y + Inches(1.06)
        for it in items:
            bx = textbox(s, x + Inches(0.30), yy, cw - Inches(0.60),
                         Inches(0.62))
            write(bx.text_frame, [("◂  ", True, color, 13),
                                  (it, False, INK, 16.5)], first=True,
                  line=1.25)
            yy += Inches(0.60)
        db = rect(s, x + Inches(0.30), y + ch - Inches(0.66), Inches(1.5),
                  Inches(0.44), fill=color,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.24)
        dtf = db.text_frame
        dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(dtf, dur, size=13.5, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)

    band(s, Inches(5.36),
         [("پیوند دو بخش:  ", True, ORANGE, 17),
          ("سازمان وقتی به موج بالاتر می‌رود که آدم‌هایش یاد بگیرند از هم یاد بگیرند.",
           False, WHITE, 20)],
         h=Inches(0.86), size=20)

    ab = textbox(s, M, Inches(6.36), CONTENT_W, Inches(0.5))
    write(ab.text_frame,
          "یک چالش کاری واقعی و حل‌نشده را از همین حالا در ذهنتان نگه دارید — در بخش دوم به آن برمی‌گردیم.",
          size=16, bold=True, color=NAVY, first=True, align=PP_ALIGN.CENTER,
          line=1.2)

    footer(s, 2, TOTAL, label=LABEL)
    notes(s, """
    تأکید کنید که بخش دوم عملی است، نه سخنرانی.
    درخواست «یک چالش واقعی در ذهن نگه دارید» را همین‌جا مطرح کنید.
    """)
    return s


# ══════════════════════════ ۳ — جداکننده: چرا تغییر؟ ════════════════════
def i03(prs, cfg):
    s = divider(prs, "بخش اول", "چرا باید\nتغییر کنیم؟",
                "سرعت تغییر بیرون، از سرعت یادگیری ما بیشتر شده است",
                "personal-growth.jpg")
    notes(s, "سه ثانیه سکوت، بعد اسلاید بعد.")
    return s


# ═══════════════════════ ۴ — قانون بقا: سرعت یادگیری ════════════════════
def i04(prs, cfg):
    s = blank(prs)
    gradient(s, 0, 0, SW, SH, NAVY_DEEP, NAVY, angle=315)

    kb = textbox(s, M, Inches(0.86), CONTENT_W, Inches(0.42))
    write(kb.text_frame, "یک قانون ساده", size=16, bold=True, color=ORANGE,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    qb = textbox(s, Inches(1.0), Inches(1.42), SW - Inches(2.0), Inches(1.7))
    tf = qb.text_frame
    write(tf, "هر سازمانی که سرعت یادگیری‌اش", size=31, bold=True, color=WHITE,
          first=True, align=PP_ALIGN.CENTER, line=1.28)
    write(tf, [("کمتر از سرعت تغییر محیط", True, ORANGE, 31),
               (" باشد،", True, WHITE, 31)],
          align=PP_ALIGN.CENTER, line=1.28)
    write(tf, "دیر یا زود حذف می‌شود.", size=31, bold=True, color=WHITE,
          align=PP_ALIGN.CENTER, line=1.28)

    cols = [
        ("تغییر محیط", ["فناوری", "انتظار ذی‌نفعان", "قواعد بازار انرژی",
                        "نسل جدید نیروی کار"], ORANGE, "سریع"),
        ("یادگیری ما", ["دوره‌های پراکنده", "تجربه‌های ثبت‌نشده",
                        "دانش در سیلوها", "آموزش به‌جای انتقال تجربه"],
         BLUE_LIGHT, "کند"),
    ]
    gap = Inches(0.42)
    cw = (CONTENT_W - gap) / 2
    y = Inches(3.44)
    ch = Inches(2.42)
    for i, (title, items, color, speed) in enumerate(cols):
        x = SW - M - cw - i * (cw + gap)
        rect(s, x, y, cw, ch, fill=WHITE, alpha=0.10,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10,
             line=color, line_w=Pt(1.5))
        tb = textbox(s, x + Inches(0.28), y + Inches(0.16), cw - Inches(0.56),
                     Inches(0.5))
        write(tb.text_frame,
              [(title, True, color, 21), ("     " + speed, True, WHITE, 15)],
              first=True, line=1.05)
        yy = y + Inches(0.78)
        for it in items:
            bx = textbox(s, x + Inches(0.30), yy, cw - Inches(0.60),
                         Inches(0.40))
            write(bx.text_frame, [("◂  ", True, color, 12),
                                  (it, False, WHITE, 16)], first=True,
                  line=1.15)
            yy += Inches(0.40)

    fb = textbox(s, M, Inches(6.12), CONTENT_W, Inches(0.62))
    write(fb.text_frame,
          [("این فاصله، ", False, BLUE_PALE, 19),
           ("همان شکاف توسعه", True, ORANGE, 20),
           (" است — و بستنش از توسعه فردی شروع می‌شود.", False, BLUE_PALE, 19)],
          first=True, align=PP_ALIGN.CENTER, line=1.2)

    notes(s, """
    این جمله از «ریوانز» است: وقتی نرخ یادگیری از نرخ تغییر کمتر شود، سازمان
    رو به زوال می‌رود. یک مثال واقعی از صنعت خودمان بزنید.
    """)
    return s


# ══════════════════════════════ ۵ — سه موج تمدنی ════════════════════════
def i05(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "پنج موجی که تمدن را ساخته‌اند",
         kicker="سه موج تافلر (۱۹۸۰) + دو موج صنعتی که بعد از او افزوده شد")

    waves = [
        ("۱", "کشاورزی", "۸۰۰۰ ق.م ـ ۱۷۵۰", "زمین",
         ["عضله انسان و دام", "خانواده و روستا", "دانش سینه‌به‌سینه"],
         GREEN, GREEN_PALE),
        ("۲", "صنعتی", "۱۷۵۰ ـ ۱۹۵۰", "ماشین",
         ["کارگر متخصص", "کارخانه و سلسله‌مراتب", "دانش در دستورالعمل"],
         BLUE_MID, BLUE_PALE),
        ("۳", "دانش", "۱۹۵۰ ـ ۲۰۱۰", "مغز",
         ["کارکن دانشی", "تیم و شبکه", "دانش، سرمایه اصلی"],
         PURPLE, T.RGBColor(0xF6, 0xEF, 0xFA)),
        ("۴", "دیجیتال", "۲۰۱۱ ـ امروز", "داده",
         ["هوش مصنوعی و اینترنت اشیا", "سیستم‌های خودتنظیم",
          "تصمیم مبتنی بر داده"], BLUE_LIGHT, BLUE_PALE),
        ("۵", "انسان‌محور", "۲۰۲۱ به بعد", "معنا",
         ["همکاری انسان و ماشین", "تاب‌آوری و پایداری",
          "خِرد جمعی"], ORANGE, ORANGE_PALE),
    ]
    gap = Inches(0.20)
    cw = (CONTENT_W - 4 * gap) / 5
    y = Inches(1.64)
    ch = Inches(3.46)
    for i, (num, name, era, driver, items, color, pale) in enumerate(waves):
        x = SW - M - cw - i * (cw + gap)
        card(s, x, y, cw, ch, fill=WHITE, line=LINE)
        rect(s, x, y, cw, Inches(1.10), fill=color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
        rect(s, x, y + Inches(0.94), cw, Inches(0.16), fill=color)
        tb = textbox(s, x + Inches(0.10), y + Inches(0.08), cw - Inches(0.20),
                     Inches(0.96), anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, "موج " + num, size=13, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        write(tb.text_frame, name, size=19, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, line=1.1)
        write(tb.text_frame, era, size=11.5, color=WHITE,
              align=PP_ALIGN.CENTER, line=1.1)

        db = rect(s, x + Inches(0.16), y + Inches(1.28), cw - Inches(0.32),
                  Inches(0.46), fill=pale,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.24)
        dtf = db.text_frame
        dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(dtf, driver, size=14, bold=True, color=color, first=True,
              align=PP_ALIGN.CENTER, line=1.0)

        yy = y + Inches(1.92)
        for it in items:
            bx = textbox(s, x + Inches(0.13), yy, cw - Inches(0.26),
                         Inches(0.48))
            write(bx.text_frame, it, size=13, color=INK, first=True,
                  align=PP_ALIGN.CENTER, line=1.18)
            yy += Inches(0.50)

    band(s, Inches(5.30),
         [("نکته کلیدی:  ", True, ORANGE, 17),
          ("در هر موج، ", False, WHITE, 19.5),
          ("شیوه مدیریت", True, ORANGE, 19.5),
          (" هم عوض شد — چون منبع اصلی ارزش عوض شده بود.", False, WHITE, 19.5)],
         h=Inches(0.84), size=19.5)

    fb = textbox(s, M, Inches(6.30), CONTENT_W, Inches(0.46))
    write(fb.text_frame,
          "امواج حذف نمی‌شوند؛ روی هم می‌افتند. سازمان ما هم‌زمان تکه‌هایی از چند موج را دارد.",
          size=15.5, color=INK_SOFT, first=True, align=PP_ALIGN.CENTER,
          line=1.2)

    footer(s, 5, TOTAL, label=LABEL)
    notes(s, """
    دقت در انتساب: آلوین تافلر در کتاب «موج سوم» (۱۹۸۰) سه موج اول را معرفی کرد.
    موج چهارم و پنجم افزوده‌های بعدی‌اند و از ادبیات صنعتی می‌آیند:
    «صنعت ۴٫۰» (نمایشگاه هانوفر، ۲۰۱۱) و «صنعت ۵٫۰» (کمیسیون اروپا، ۲۰۲۱).
    اگر کسی پرسید، همین را صادقانه بگویید؛ چارچوب پنج‌موجی یک تعمیم مفید است، نه نقل‌قول تافلر.
    تأکید: موج دوم بد نیست — برای زمان خودش درست بود. مسئله، ماندن در آن است.
    تاریخ‌ها تقریبی‌اند و مرز موج‌ها در منابع مختلف کمی جابه‌جا می‌شود.
    """)
    return s


# ════════════════════ ۶ — مکاتب مدیریتی روی خط زمان ═════════════════════
def i06(prs, cfg):
    s = base_slide(prs, WHITE)
    head(s, "هر موج، مکتب مدیریتی خودش را ساخت",
         kicker="نقشه تاریخی مکاتب مدیریت — پنج موج، پنج پاسخ متفاوت")

    cols = [
        (GREEN, GREEN_PALE, "موج ۱ · کشاورزی", "تا ۱۷۵۰", [
            ("مدیریت سنتی", "استادـشاگردی"),
            ("نظام صنفی", "مهارت در کارگاه"),
        ]),
        (BLUE_MID, BLUE_PALE, "موج ۲ · صنعتی", "۱۷۵۰ ـ ۱۹۵۰", [
            ("مدیریت علمی — تیلور ۱۹۱۱", "بهترین راه انجام کار"),
            ("مدیریت اداری — فایول ۱۹۱۶", "برنامه‌ریزی، سازمان‌دهی، کنترل"),
            ("بوروکراسی — وبر ۱۹۲۲", "سلسله‌مراتب، قاعده، غیرشخصی"),
            ("روابط انسانی — مِیو ۱۹۳۲", "انسان فقط ماشین نیست"),
            ("نظریه X و Y — مک‌گرگور ۱۹۶۰", "دو نگاه به انگیزه کارکنان"),
            ("سیستمی و اقتضایی ۱۹۷۰", "یک نسخه برای همه نیست"),
        ]),
        (PURPLE, T.RGBColor(0xF6, 0xEF, 0xFA), "موج ۳ · دانش", "۱۹۵۰ ـ ۲۰۱۰", [
            ("کیفیت جامع — دهه ۱۹۸۰", "بهبود مستمر، مشارکت همه"),
            ("سازمان یادگیرنده — سنگه ۱۹۹۰", "یادگیری، مزیت رقابتی پایدار"),
            ("مدیریت دانش — نوناکا ۱۹۹۵", "تجربه ضمنی، دانش سازمانی می‌شود"),
            ("سازمان چابک ۲۰۰۱", "تیم‌های کوچک، چرخه کوتاه"),
            ("رهبری مربی‌گرا", "مدیر به‌جای دستور، سؤال می‌پرسد"),
        ]),
        (BLUE_LIGHT, BLUE_PALE, "موج ۴ · دیجیتال", "۲۰۱۱ ـ امروز", [
            ("صنعت ۴٫۰ — هانوفر ۲۰۱۱", "اتصال ماشین، داده، خودکارسازی"),
            ("سازمان داده‌محور", "تصمیم بر پایه داده، نه حدس"),
            ("تیم‌های خودگردان", "اختیار توزیع‌شده"),
            ("سازمان شبکه‌ای", "همکاری برخط و فرامرزی"),
            ("مدیریت الگوریتمی", "داده، جای قضاوت را نمی‌گیرد"),
        ]),
        (ORANGE, ORANGE_PALE, "موج ۵ · انسان‌محور", "۲۰۲۱ به بعد", [
            ("صنعت ۵٫۰ — اروپا ۲۰۲۱", "انسان‌محور، تاب‌آور، پایدار"),
            ("رهبری تاب‌آور", "مدیریت در بحران و عدم‌قطعیت"),
            ("مدیریت انسان‌محور", "رفاه کارکنان کنار بهره‌وری"),
            ("پایداری و مسئولیت", "سود، کنار محیط زیست"),
            ("رهبری خدمتگزار", "مدیر، راه رشد را باز می‌کند"),
        ]),
    ]
    gap = Inches(0.16)
    widths = [Inches(1.52), Inches(2.86), Inches(2.62), Inches(2.30),
              Inches(2.19)]
    y = Inches(1.66)
    ch = Inches(4.52)
    hh = Inches(0.78)
    xcur = SW - M
    for (color, pale, wave, era, rows), cw in zip(cols, widths):
        x = xcur - cw
        card(s, x, y, cw, ch, fill=WHITE, line=LINE)
        rect(s, x, y, cw, hh, fill=color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
        rect(s, x, y + hh - Inches(0.16), cw, Inches(0.16), fill=color)
        tb = textbox(s, x + Inches(0.06), y, cw - Inches(0.12), hh,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, wave, size=13.5, bold=True, color=WHITE,
              first=True, align=PP_ALIGN.CENTER, line=1.05)
        write(tb.text_frame, fa(era), size=11, color=WHITE,
              align=PP_ALIGN.CENTER, line=1.1)

        yy = y + hh + Inches(0.10)
        rh = Inches(0.615)
        for j, (name, desc) in enumerate(rows):
            if j % 2 == 0:
                rect(s, x + Inches(0.08), yy, cw - Inches(0.16),
                     rh - Inches(0.05), fill=pale,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
            bx = textbox(s, x + Inches(0.16), yy, cw - Inches(0.32),
                         rh - Inches(0.05), anchor=MSO_ANCHOR.MIDDLE)
            write(bx.text_frame, name, size=12.5, bold=True, color=INK,
                  first=True, line=1.08)
            write(bx.text_frame, desc, size=10.5, color=INK_SOFT, line=1.12)
            yy += rh
        xcur -= cw + gap

    fb = textbox(s, M, Inches(6.30), CONTENT_W, Inches(0.46))
    write(fb.text_frame,
          [("پیام نقشه:  ", True, ORANGE, 14.5),
           ("هر ستون، پاسخ مدیریتیِ زمانِ خودش بود؛ مسئله وقتی شروع می‌شود که محیط "
            "به ستون چهارم رفته و ما هنوز با ستون دوم مدیریت می‌کنیم.",
            False, INK, 14.5)],
          first=True, align=PP_ALIGN.CENTER, line=1.2)

    footer(s, 6, TOTAL, label=LABEL)
    notes(s, """
    این اسلاید ستون فقرات بخش اول است — روی آن وقت بگذارید.
    تفکیک پنج ستون عمدی است: موج چهارم و پنجم دو منطق متفاوت‌اند،
    نه یک دوره واحد. موج چهارم درباره «داده و خودکارسازی» است،
    موج پنجم درباره «انسان، تاب‌آوری و پایداری».
    تاریخِ زیر عنوان هر ستون، بازه خودِ موج است و تاریخ کنار هر مکتب،
    سال انتشار اثر شاخص آن. مکاتب معمولاً با تأخیر نسبت به موج شکل می‌گیرند
    (مثلاً تیلور ۱۹۱۱، یعنی حدود ۱۵۰ سال پس از آغاز موج دوم).
    مکاتب هم حذف نمی‌شوند؛ روی هم انباشته می‌شوند.
    پیام: ستون آبی (موج دوم) هنوز شیوه غالب مدیریت در سازمان‌های ماست،
    در حالی که محیط کسب‌وکار وارد ستون‌های چهارم و پنجم شده است.
    از مخاطب بپرسید: «کدام ردیف‌ها را در سازمان خودمان می‌بینید؟»
    """)
    return s


# ═══════════ ۷ — جدول انطباق موج‌ها، مکاتب مدیریت و جایگاه ما ═══════════
def i07(prs, cfg):
    s = base_slide(prs, WHITE)
    head(s, "جدول انطباق: هر موج، کدام سبک مدیریت و کدام شیوه توسعه؟",
         kicker="نقشه یک‌نگاهی — وضعیت فعلی و مسیر پیش رو")

    # ستون‌ها از راست به چپ
    cols = [
        ("موج", Inches(1.30)),
        ("دوره", Inches(1.15)),
        ("مکتب و سبک مدیریت غالب", Inches(3.05)),
        ("نگاه به انسان", Inches(1.85)),
        ("شیوه توسعه افراد", Inches(2.55)),
        ("جایگاه سازمان ما", Inches(2.14)),
    ]
    rows = [
        (GREEN, "موج ۱", "کشاورزی", "تا ۱۷۵۰",
         "مدیریت سنتی، نظام صنفی، استادـشاگردی",
         "بازوی کار", "تقلید از استاد، سینه‌به‌سینه",
         "گذشته", INK_SOFT, "outline", None),
        (BLUE_MID, "موج ۲", "صنعتی", "۱۷۵۰ ـ ۱۹۵۰",
         "تیلور، فایول، وبر — مدیریت علمی و بوروکراسی",
         "جزئی از ماشین", "آموزش شغلی، دستورالعمل، کلاس",
         "بدنه اصلی ما", ORANGE, "solid", ORANGE_PALE),
        (PURPLE, "موج ۳", "دانش", "۱۹۵۰ ـ ۲۰۱۰",
         "سازمان یادگیرنده، مدیریت دانش، چابک",
         "صاحب دانش", "منتورینگ، یادگیری از هم",
         "جزیره‌ای و پراکنده", AMBER, "solid", None),
        (BLUE_LIGHT, "موج ۴", "دیجیتال", "۲۰۱۱ ـ امروز",
         "صنعت ۴٫۰، سازمان داده‌محور، تیم‌های خودگردان",
         "شریک ماشین", "یادگیری در جریان کار، بازخورد داده",
         "چند پایلوت", BLUE_MID, "solid", None),
        (T.RGBColor(0xE8, 0x59, 0x2B), "موج ۵", "انسان‌محور", "۲۰۲۱ به بعد",
         "صنعت ۵٫۰، مدیریت انسان‌محور و تاب‌آور",
         "سرمایه اصلی", "مربی‌گری، رشد معنامحور",
         "افق پیش رو", GREEN, "outline", GREEN_PALE),
    ]

    y0 = Inches(1.62)
    hh = Inches(0.48)
    rh = Inches(0.755)

    # سطر عنوان
    rect(s, M, y0, CONTENT_W, hh, fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
    rect(s, M, y0 + Inches(0.24), CONTENT_W, Inches(0.26), fill=NAVY)
    xc = SW - M
    for name, cw in cols:
        tb = textbox(s, xc - cw, y0, cw, hh, anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, name, size=13.5, bold=True, color=WHITE,
              first=True, align=PP_ALIGN.CENTER, line=1.0)
        xc -= cw

    # سطرهای داده
    y = y0 + hh + Inches(0.06)
    for (wc, wnum, wname, era, school, human, dev,
         status, scol, style, tint) in rows:
        if tint is not None:
            rect(s, M, y, CONTENT_W, rh - Inches(0.06), fill=tint,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
        else:
            rect(s, M, y, CONTENT_W, rh - Inches(0.06), fill=WHITE,
                 line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
        # نوار رنگی موج در لبه راست
        rect(s, SW - M - Inches(0.10), y, Inches(0.10), rh - Inches(0.06),
             fill=wc)

        xc = SW - M
        # ستون موج
        cw = cols[0][1]
        tb = textbox(s, xc - cw + Inches(0.14), y, cw - Inches(0.20),
                     rh - Inches(0.06), anchor=MSO_ANCHOR.MIDDLE)
        tf = tb.text_frame
        write(tf, wname, size=14.5, bold=True, color=wc, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        write(tf, wnum, size=10.5, color=INK_SOFT, align=PP_ALIGN.CENTER,
              line=1.05)
        xc -= cw
        # دوره
        cw = cols[1][1]
        tb = textbox(s, xc - cw, y, cw, rh - Inches(0.06),
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, fa(era), size=12.5, color=INK_SOFT, first=True,
              align=PP_ALIGN.CENTER, line=1.05)
        xc -= cw
        # مکتب
        cw = cols[2][1]
        tb = textbox(s, xc - cw + Inches(0.12), y, cw - Inches(0.24),
                     rh - Inches(0.06), anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, school, size=13, bold=True, color=INK,
              first=True, align=PP_ALIGN.RIGHT, line=1.16)
        xc -= cw
        # نگاه به انسان
        cw = cols[3][1]
        tb = textbox(s, xc - cw + Inches(0.10), y, cw - Inches(0.20),
                     rh - Inches(0.06), anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, human, size=13, color=INK, first=True,
              align=PP_ALIGN.CENTER, line=1.12)
        xc -= cw
        # شیوه توسعه
        cw = cols[4][1]
        tb = textbox(s, xc - cw + Inches(0.12), y, cw - Inches(0.24),
                     rh - Inches(0.06), anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, dev, size=12.5, color=INK_SOFT, first=True,
              align=PP_ALIGN.RIGHT, line=1.14)
        xc -= cw
        # جایگاه ما — قرص وضعیت
        cw = cols[5][1]
        pw = cw - Inches(0.42)
        ph = Inches(0.42)
        px = xc - cw + Inches(0.21)
        py = y + (rh - Inches(0.06) - ph) / 2
        if style == "solid":
            rect(s, px, py, pw, ph, fill=scol,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.42)
            tcol = WHITE
        else:
            rect(s, px, py, pw, ph, fill=WHITE, line=scol, line_w=Pt(1.5),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.42)
            tcol = scol
        tb = textbox(s, px, py, pw, ph, anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, status, size=12.5, bold=True, color=tcol,
              first=True, align=PP_ALIGN.CENTER, line=1.0)
        y += rh

    # جمع‌بندی
    by = y + Inches(0.10)
    bh = Inches(0.74)
    rect(s, M, by, CONTENT_W, bh, fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    rect(s, SW - M - Inches(0.09), by, Inches(0.09), bh, fill=ORANGE)
    tb = textbox(s, M + Inches(0.30), by, CONTENT_W - Inches(0.60),
                 bh, anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame,
          [("پیام جدول:  ", True, ORANGE, 15),
           ("ما با سبک مدیریت موج دوم، می‌خواهیم در محیط موج چهارم و پنجم رقابت کنیم. "
            "ستون «شیوه توسعه افراد» کوتاه‌ترین راه عبور است.",
            False, WHITE, 16)],
          first=True, line=1.22)

    footer(s, 7, TOTAL, label=LABEL)
    notes(s, """
    این اسلاید، جمع‌بندی بصری کل بخش اول است؛ اگر وقت کم بود، همین یکی را کامل بگویید.
    سه ستون سمت راست، «تاریخ» را می‌گوید و سه ستون سمت چپ، «تکلیف ما» را.
    روی ستون «نگاه به انسان» مکث کنید: از بازوی کار تا سرمایه اصلی.
    نکته کلیدی: سازمان نمی‌تواند از موج دوم مستقیم به موج پنجم بپرد،
    ولی ستون «شیوه توسعه افراد» را می‌تواند از همین امروز جلو ببرد —
    و این دقیقاً همان کاری است که در بخش دوم جلسه دنبال می‌کنیم.
    اگر پرسیدند «چرا موج ۴ فقط پایلوت است؟» بگویید: ابزار دیجیتال خریده‌ایم،
    ولی سبک تصمیم‌گیری هنوز موج دومی است.
    """)
    return s


# ════════════════ ۸ — موج چهارم: صنعت دیجیتال و داده ════════════════════
def i08(prs, cfg):
    s = blank(prs)
    p, sz = img("w4-digital.jpg")
    full_bleed(s, p, sz)
    scrim_gradient(s, 0, 0, SW, SH, angle=0, a_from=0.18, a_to=0.94)
    scrim(s, 0, 0, SW, SH, alpha=0.26)

    kb = textbox(s, M, Inches(0.50), CONTENT_W, Inches(0.40))
    write(kb.text_frame, "موج چهارم  ·  از ۲۰۱۱ تا امروز", size=15.5,
          bold=True, color=ORANGE, first=True, line=1.0)
    tb = textbox(s, M, Inches(0.90), CONTENT_W, Inches(0.66))
    write(tb.text_frame, "صنعت دیجیتال: وقتی ماشین‌ها حرف می‌زنند", size=31,
          bold=True, color=WHITE, first=True, line=1.05)
    hline(s, SW - M - Inches(1.2), Inches(1.64), Inches(1.2), ORANGE, Pt(4.5))

    tw = Inches(7.6)
    tx = SW - M - tw
    items = [
        ("موتور محرک", "داده و اتصال — اینترنت اشیا، هوش مصنوعی"),
        ("در نفت و گاز", "دوقلوی دیجیتال، نگهداشت پیش‌بین، پایش برخط چاه"),
        ("منطق تصمیم", "از «تجربه می‌گوید» به «داده نشان می‌دهد»"),
        ("نقش انسان", "کمتر اپراتور، بیشتر تحلیل‌گر"),
    ]
    y = Inches(2.02)
    rh = Inches(0.86)
    for i, (label, desc) in enumerate(items):
        rect(s, tx, y, tw, rh - Inches(0.08), fill=WHITE, alpha=0.12,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
        rect(s, tx + tw - Inches(0.08), y, Inches(0.08), rh - Inches(0.08),
             fill=BLUE_LIGHT)
        lb = textbox(s, tx + tw - Inches(2.35), y, Inches(2.10),
                     rh - Inches(0.08), anchor=MSO_ANCHOR.MIDDLE)
        write(lb.text_frame, label, size=16.5, bold=True, color=BLUE_LIGHT,
              first=True, align=PP_ALIGN.RIGHT, line=1.05)
        db = textbox(s, tx + Inches(0.28), y, tw - Inches(2.75),
                     rh - Inches(0.08), anchor=MSO_ANCHOR.MIDDLE)
        write(db.text_frame, desc, size=17, color=WHITE, first=True, line=1.15)
        y += rh

    rect(s, tx, Inches(5.66), tw, Inches(1.10), fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    rect(s, tx + tw - Inches(0.09), Inches(5.66), Inches(0.09), Inches(1.10),
         fill=ORANGE)
    fb = textbox(s, tx + Inches(0.30), Inches(5.66), tw - Inches(0.60),
                 Inches(1.10), anchor=MSO_ANCHOR.MIDDLE)
    write(fb.text_frame,
          [("هشدار مهم:  ", True, ORANGE, 16),
           ("دیجیتالی‌کردن یک سازمانِ موج دومی، فقط کاغذبازی را سریع‌تر می‌کند — نه هوشمندتر.",
            False, WHITE, 18)],
          first=True, line=1.22)

    notes(s, """
    موج چهارم همان «انقلاب صنعتی چهارم» یا Industry 4.0 است (از ۲۰۱۱).
    نکته کلیدی برای مدیران: خرید نرم‌افزار و سنسور، به‌تنهایی گذار نیست.
    اگر تصمیم‌گیری همچنان متمرکز و سلسله‌مراتبی بماند، داده فقط انبار می‌شود.
    """)
    return s


# ═════════════ ۸ — موج پنجم: انسان‌محوری، تاب‌آوری، پایداری ═════════════
def i09(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "موج پنجم: بازگشت انسان به مرکز",
         kicker="موج پنجم  ·  از ۲۰۲۱ — کمیسیون اروپا")

    pw = Inches(4.30)
    photo_card(s, M, Inches(1.66), pw, Inches(2.72), "w5-human.jpg",
               caption="انسان و ماشین، در کنار هم — نه به‌جای هم")

    card(s, M, Inches(4.54), pw, Inches(2.24), fill=NAVY, line=None)
    rect(s, M + pw - Inches(0.09), Inches(4.54), Inches(0.09), Inches(2.24),
         fill=ORANGE)
    tb = textbox(s, M + Inches(0.26), Inches(4.76), pw - Inches(0.55),
                 Inches(1.85))
    tf = tb.text_frame
    write(tf, "تفاوت با موج چهارم", size=14.5, bold=True, color=ORANGE,
          first=True, line=1.0)
    write(tf, "موج چهارم می‌پرسد: چطور کارخانه را کارآمدتر کنیم؟",
          size=15, color=BLUE_PALE, space_before=8, line=1.26)
    write(tf, "موج پنجم می‌پرسد: چطور کار را انسانی‌تر، تاب‌آورتر و پایدارتر کنیم؟",
          size=15.5, bold=True, color=WHITE, space_before=6, line=1.26)

    tx = M + pw + Inches(0.42)
    tw = SW - M - tx
    pillars = [
        ("۱", "انسان‌محوری", "فناوری در خدمت آدم‌ها، نه جایگزین آن‌ها",
         ORANGE, ORANGE_PALE),
        ("۲", "تاب‌آوری", "توان بازگشت سریع پس از بحران و تحریم",
         BLUE_MID, BLUE_PALE),
        ("۳", "پایداری", "مسئولیت زیست‌محیطی و کاهش انتشار", GREEN, GREEN_PALE),
    ]
    y = Inches(1.66)
    rh = Inches(1.16)
    for num, title, desc, color, pale in pillars:
        rect(s, tx, y, tw, rh - Inches(0.10), fill=pale,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11)
        rect(s, tx + tw - Inches(0.09), y, Inches(0.09), rh - Inches(0.10),
             fill=color)
        nb = rect(s, tx + tw - Inches(0.88), y + Inches(0.29), Inches(0.46),
                  Inches(0.46), fill=color, shape=MSO_SHAPE.OVAL)
        ntf = nb.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(ntf, num, size=16, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        bx = textbox(s, tx + Inches(0.28), y, tw - Inches(1.26),
                     rh - Inches(0.10), anchor=MSO_ANCHOR.MIDDLE)
        write(bx.text_frame, title, size=20.5, bold=True, color=color,
              first=True, line=1.05)
        write(bx.text_frame, desc, size=15, color=INK_SOFT, line=1.18)
        y += rh

    ky = Inches(5.20)
    rect(s, tx, ky, tw, Inches(1.58), fill=WHITE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11, line=ORANGE,
         line_w=Pt(1.75))
    kb = textbox(s, tx + Inches(0.28), ky + Inches(0.16), tw - Inches(0.56),
                 Inches(1.26))
    tf = kb.text_frame
    write(tf, "خبر خوب برای ما", size=14.5, bold=True, color=ORANGE_DEEP,
          first=True, line=1.0)
    write(tf, "موج پنجم بیش از آنکه به فناوری گران نیاز داشته باشد، به فرهنگ یادگیری و اعتماد نیاز دارد.",
          size=16.5, bold=True, color=NAVY, space_before=7, line=1.28)
    write(tf, "یعنی می‌شود از همین امروز شروع کرد — بدون بودجه ارزی.",
          size=15, color=INK_SOFT, space_before=5, line=1.22)

    footer(s, 9, TOTAL, label=LABEL)
    notes(s, """
    موج پنجم (Industry 5.0) را کمیسیون اروپا در ۲۰۲۱ رسمی کرد:
    صنعتِ انسان‌محور، تاب‌آور و پایدار.
    پیام کلیدی: موج پنجم جایگزین چهارم نیست؛ لایه ارزشی روی آن است.
    برای سازمان ما این خبر خوبی است — چون نقطه ورودش فرهنگ است، نه سرمایه.
    """)
    return s


# ═══════════════════ ۹ — ما کجای این نقشه ایستاده‌ایم؟ ═══════════════════
def i10(prs, cfg):
    s = blank(prs)
    p, sz = img("w1-refinery.jpg")
    full_bleed(s, p, sz)
    scrim_gradient(s, 0, 0, SW, SH, angle=0, a_from=0.14, a_to=0.95)
    scrim(s, 0, 0, SW, SH, alpha=0.30)

    kb = textbox(s, M, Inches(0.52), CONTENT_W, Inches(0.40))
    write(kb.text_frame, "تشخیص وضعیت", size=15.5, bold=True, color=ORANGE,
          first=True, line=1.0)
    tb = textbox(s, M, Inches(0.92), CONTENT_W, Inches(0.66))
    write(tb.text_frame, "صنعت نفت و گاز ایران کجای این نقشه است؟", size=31,
          bold=True, color=WHITE, first=True, line=1.05)
    hline(s, SW - M - Inches(1.2), Inches(1.66), Inches(1.2), ORANGE, Pt(4.5))

    tw = Inches(7.15)
    tx = SW - M - tw
    rows = [
        ("زیرساخت فنی", "موج دوم، با جزیره‌هایی از موج سوم", 0.62, BLUE_MID),
        ("ساختار سازمانی", "سلسله‌مراتبی، تصمیم متمرکز", 0.30, BLUE_MID),
        ("شیوه انتقال دانش", "شفاهی، وابسته به افراد کلیدی", 0.25, RED),
        ("سرمایه انسانی", "باتجربه، ولی تجربه‌اش ثبت‌نشده", 0.40, AMBER),
    ]
    y = Inches(2.00)
    rh = Inches(0.92)
    for name, desc, frac, color in rows:
        rect(s, tx, y, tw, rh - Inches(0.10), fill=WHITE, alpha=0.13,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
        lb = textbox(s, tx + tw - Inches(2.55), y, Inches(2.35),
                     rh - Inches(0.10), anchor=MSO_ANCHOR.MIDDLE)
        write(lb.text_frame, name, size=17, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.RIGHT, line=1.05)
        db = textbox(s, tx + Inches(1.75), y, Inches(2.90), rh - Inches(0.10),
                     anchor=MSO_ANCHOR.MIDDLE)
        write(db.text_frame, desc, size=14, color=BLUE_PALE, first=True,
              line=1.15)
        gw = Inches(1.55)
        gx = tx + Inches(0.16)
        gy = y + Inches(0.30)
        rect(s, gx, gy, gw, Inches(0.22), fill=WHITE, alpha=0.22,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.30)
        rect(s, gx + gw * (1 - frac), gy, gw * frac, Inches(0.22), fill=color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.30)
        y += rh

    rect(s, tx, Inches(5.76), tw, Inches(1.06), fill=ORANGE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    fb = textbox(s, tx + Inches(0.30), Inches(5.76), tw - Inches(0.60),
                 Inches(1.06), anchor=MSO_ANCHOR.MIDDLE)
    write(fb.text_frame,
          [("جمع‌بندی صادقانه:  ", True, WHITE, 16),
           ("بدنه سازمان در موج دوم است، با جزیره‌هایی از موج سوم و چهارم — و این ناهم‌ترازی، پرهزینه‌ترین وضعیت است.",
            False, WHITE, 17)],
          first=True, line=1.2)

    notes(s, """
    این اسلاید را با احتیاط و بدون سرزنش ارائه کنید.
    پیام: ما ضعیف نیستیم؛ ما در حال گذاریم و گذار، مدیریت می‌خواهد.
    اگر داده واقعی از سازمان دارید، نوارها را با آن تنظیم کنید.
    """)
    return s


# ═══════════════════════ ۸ — نشانه‌های ماندن در موج دوم ═════════════════
def i11(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "نشانه‌های ماندن در موج دوم", kicker="این‌ها را می‌شناسید؟")

    signs = [
        ("تصمیم فقط بالا گرفته می‌شود", "کارشناس منتظر دستور می‌ماند"),
        ("دانش با بازنشستگی می‌رود", "سی سال تجربه، بدون جانشین"),
        ("آموزش یعنی کلاس", "دوره برگزار می‌شود، رفتار عوض نمی‌شود"),
        ("واحدها با هم حرف نمی‌زنند", "هر واحد، جزیره‌ای مستقل"),
        ("خطا پنهان می‌شود", "چون به‌جای درس، دنبال مقصر می‌گردیم"),
        ("جلسه یعنی گزارش‌دهی", "نه حل مسئله مشترک"),
    ]
    gap = Inches(0.28)
    cw = (CONTENT_W - gap) / 2
    y0 = Inches(1.66)
    rh = Inches(0.94)
    for i, (title, desc) in enumerate(signs):
        col = i % 2
        row = i // 2
        x = SW - M - cw - col * (cw + gap)
        y = y0 + row * (rh + Inches(0.14))
        rect(s, x, y, cw, rh, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.11, line=LINE)
        rect(s, x + cw - Inches(0.08), y, Inches(0.08), rh, fill=RED)
        nb = rect(s, x + cw - Inches(0.80), y + Inches(0.25), Inches(0.42),
                  Inches(0.42), fill=RED_PALE, shape=MSO_SHAPE.OVAL)
        ntf = nb.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(ntf, fa(i + 1), size=15, bold=True, color=RED, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, x + Inches(0.24), y, cw - Inches(1.16), rh,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=18, bold=True, color=NAVY, first=True,
              line=1.05)
        write(tb.text_frame, desc, size=14, color=INK_SOFT, line=1.18)

    band(s, Inches(5.06),
         [("هیچ‌کدام از این‌ها نشانه بی‌کفایتی نیست — ", False, WHITE, 19),
          ("نشانه یک سیستم قدیمی است.", True, ORANGE, 19.5)],
         h=Inches(0.80), size=19)

    fb = textbox(s, M, Inches(6.04), CONTENT_W, Inches(0.72))
    write(fb.text_frame,
          "و سیستم را نمی‌شود با بخشنامه عوض کرد؛ با تغییر رفتار روزمره آدم‌ها عوض می‌شود.",
          size=18, bold=True, color=NAVY, first=True, align=PP_ALIGN.CENTER,
          line=1.25)

    footer(s, 11, TOTAL, label=LABEL)
    notes(s, """
    از مخاطب بپرسید کدام مورد را بیشتر می‌بینند و رأی‌گیری دستی کنید.
    این مشارکت، مقاومت را کم می‌کند چون خودشان تشخیص را تأیید کرده‌اند.
    """)
    return s


# ══════════════════════════ ۹ — مسیر گذار به موج سوم ════════════════════
def i12(prs, cfg):
    s = base_slide(prs, WHITE)
    head(s, "مسیر گذار به موج بالاتر",
         kicker="چهار گام، به ترتیب — از موج دوم به سوم، و آماده‌سازی برای چهارم و پنجم")

    pw = Inches(4.55)
    photo_card(s, M, Inches(1.66), pw, Inches(2.62), "w2-control.jpg",
               caption="موج سوم یعنی تصمیم، نزدیک به محل کار گرفته شود")

    card(s, M, Inches(4.44), pw, Inches(2.34), fill=NAVY, line=None)
    rect(s, M + pw - Inches(0.09), Inches(4.44), Inches(0.09), Inches(2.34),
         fill=ORANGE)
    tb = textbox(s, M + Inches(0.26), Inches(4.66), pw - Inches(0.55),
                 Inches(1.95))
    tf = tb.text_frame
    write(tf, "چرا از گام ۱ شروع می‌کنیم؟", size=14.5, bold=True, color=ORANGE,
          first=True, line=1.0)
    write(tf, "چون سه گام بعدی بدون تغییر رفتار روزمره، فقط سند و نمودار می‌ماند.",
          size=16, color=WHITE, space_before=9, line=1.3)
    write(tf, "ارزان‌ترین و سریع‌ترین گام هم همین اولی است.",
          size=15, color=BLUE_PALE, space_before=7, line=1.28)

    tx = M + pw + Inches(0.42)
    tw = SW - M - tx
    steps = [
        ("۱", "تغییر رفتار یادگیری", "از «کلاس رفتن» به «از هم یاد گرفتن»",
         "منتورینگ گروهی، بازخورد همتا", ORANGE, True),
        ("۲", "ثبت و گردش دانش", "تجربه ضمنی افراد، دارایی سازمان شود",
         "مستندسازی درس‌آموخته‌ها", BLUE_MID, False),
        ("۳", "توزیع اختیار تصمیم", "تصمیم نزدیک به محل کار گرفته شود",
         "تیم‌های حل مسئله", BLUE_MID, False),
        ("۴", "ساختار چابک", "چرخه‌های کوتاه، بازنگری منظم",
         "بازطراحی فرایندها", PURPLE, False),
    ]
    y = Inches(1.66)
    rh = Inches(1.22)
    for num, title, desc, how, color, hi in steps:
        rect(s, tx, y, tw, rh - Inches(0.10),
             fill=ORANGE_PALE if hi else PAPER,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11,
             line=ORANGE if hi else None, line_w=Pt(2))
        rect(s, tx + tw - Inches(0.08), y, Inches(0.08), rh - Inches(0.10),
             fill=color)
        nb = rect(s, tx + tw - Inches(0.86), y + Inches(0.32), Inches(0.46),
                  Inches(0.46), fill=color, shape=MSO_SHAPE.OVAL)
        ntf = nb.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(ntf, num, size=16, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, tx + Inches(0.26), y + Inches(0.06),
                     tw - Inches(1.24), rh - Inches(0.22))
        tf = tb.text_frame
        ch = [(title, True, color if hi else NAVY, 19)]
        if hi:
            ch.append(("     ← از اینجا شروع می‌کنیم", True, ORANGE, 13.5))
        write(tf, ch, first=True, line=1.05)
        write(tf, desc, size=14.5, color=INK_SOFT, line=1.18)
        write(tf, "ابزار: " + how, size=13, color=color, line=1.15)
        y += rh

    footer(s, 12, TOTAL, label=LABEL)
    notes(s, """
    ترتیب گام‌ها مهم است. اگر مدیری پرسید چرا از ساختار شروع نمی‌کنیم،
    بگویید: ساختار جدید با رفتار قدیمی، همان سازمان قبلی است با نمودار نو.
    """)
    return s


# ═════════════════════ ۱۰ — از توسعه فردی تا سازمانی ════════════════════
def i13(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "چرا توسعه فردی، مقدمه توسعه سازمانی است",
         kicker="سه حلقه به‌هم‌پیوسته")

    rings = [
        ("فرد", "من", ["مهارت پرسیدن", "تأمل بر تجربه", "پذیرش بازخورد"],
         ORANGE, ORANGE_PALE),
        ("تیم", "ما", ["اعتماد و فضای امن", "حل مسئله مشترک",
                       "یادگیری از خطا"], BLUE_MID, BLUE_PALE),
        ("سازمان", "همه", ["گردش دانش", "تصمیم توزیع‌شده", "حافظه سازمانی"],
         GREEN, GREEN_PALE),
    ]
    gap = Inches(0.30)
    cw = (CONTENT_W - 2 * gap) / 3
    y = Inches(1.68)
    ch = Inches(2.72)
    for i, (name, pron, items, color, pale) in enumerate(rings):
        x = SW - M - cw - i * (cw + gap)
        card(s, x, y, cw, ch, fill=WHITE, line=LINE)
        rect(s, x, y, cw, Inches(0.07), fill=color)
        cb = rect(s, x + (cw - Inches(0.92)) / 2, y + Inches(0.30),
                  Inches(0.92), Inches(0.92), fill=pale, shape=MSO_SHAPE.OVAL,
                  line=color, line_w=Pt(2.5))
        ctf = cb.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(ctf, pron, size=21, bold=True, color=color, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, x + Inches(0.16), y + Inches(1.34), cw - Inches(0.32),
                     Inches(0.44))
        write(tb.text_frame, name, size=20, bold=True, color=color, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        yy = y + Inches(1.82)
        for it in items:
            bx = textbox(s, x + Inches(0.18), yy, cw - Inches(0.36),
                         Inches(0.30))
            write(bx.text_frame, it, size=14, color=INK, first=True,
                  align=PP_ALIGN.CENTER, line=1.1)
            yy += Inches(0.30)

    ay = Inches(4.56)
    for i in range(2):
        ax = SW - M - cw - i * (cw + gap) - gap - Inches(0.11)
        ar = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, ax, ay, Inches(0.22),
                                Inches(0.20))
        ar.fill.solid()
        ar.fill.fore_color.rgb = ORANGE
        no_line(ar)
        shadow_off(ar)

    band(s, Inches(5.02),
         [("سازمان یادگیرنده از بالا ابلاغ نمی‌شود — ", False, WHITE, 19.5),
          ("از یک گفتگوی درست بین دو نفر شروع می‌شود.", True, ORANGE, 20)],
         h=Inches(0.88), size=19.5)

    gap2 = Inches(0.34)
    hw = (CONTENT_W - gap2) / 2
    by = Inches(6.06)
    for x, label, value, color in [
        (SW - M - hw, "اگر فقط فرد رشد کند", "دانش با رفتنش می‌رود", RED),
        (M, "اگر فرد در گروه رشد کند", "دانش در سازمان می‌ماند", GREEN),
    ]:
        rect(s, x, by, hw, Inches(0.72), fill=WHITE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11, line=color,
             line_w=Pt(1.5))
        tb = textbox(s, x + Inches(0.22), by, hw - Inches(0.44), Inches(0.72),
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, [(label + ":  ", False, INK_SOFT, 14.5),
                              (value, True, color, 16)], first=True,
              align=PP_ALIGN.CENTER, line=1.1)

    footer(s, 13, TOTAL, label=LABEL)
    notes(s, """
    این اسلاید پل بخش اول به بخش دوم است.
    جمله انتقال: «حالا برویم سراغ ابزاری که این حلقه اول را می‌سازد.»
    """)
    return s


# ════════════════════ ۱۱ — پل به بخش دوم: منتورینگ گروهی ════════════════
def i14(prs, cfg):
    s = blank(prs)
    p, sz = img("01-hands-together.jpg")
    full_bleed(s, p, sz)
    gradient(s, 0, 0, SW, SH, NAVY_DEEP, NAVY, angle=45,
             alpha1=0.93, alpha2=0.88)
    gradient(s, 0, 0, SW, SH, ORANGE_DEEP, NAVY_DEEP, angle=45,
             alpha1=0.30, alpha2=0.0)

    kb = textbox(s, M, Inches(1.06), CONTENT_W, Inches(0.44))
    write(kb.text_frame, "پایان بخش اول", size=16, bold=True, color=ORANGE,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    tb = textbox(s, Inches(1.1), Inches(1.62), SW - Inches(2.2), Inches(2.1))
    tf = tb.text_frame
    write(tf, "تشخیص روشن است.", size=30, bold=True, color=WHITE, first=True,
          align=PP_ALIGN.CENTER, line=1.28)
    write(tf, [("حالا ", True, WHITE, 38), ("اولین گام", True, ORANGE, 38),
               (" را با هم برمی‌داریم", True, WHITE, 38)],
          align=PP_ALIGN.CENTER, line=1.28)

    items = [("چیستی", "این روش دقیقاً چیست"),
             ("چرایی", "چرا جواب می‌دهد"),
             ("نحوه اجرا", "یک جلسه واقعی، دقیقه به دقیقه")]
    gap = Inches(0.32)
    cw = (CONTENT_W - 2 * gap) / 3
    y = Inches(4.06)
    ch = Inches(1.30)
    for i, (title, desc) in enumerate(items):
        x = SW - M - cw - i * (cw + gap)
        rect(s, x, y, cw, ch, fill=WHITE, alpha=0.14,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10, line=WHITE,
             line_w=Pt(1))
        tb = textbox(s, x + Inches(0.16), y, cw - Inches(0.32), ch,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=21, bold=True, color=WHITE,
              first=True, align=PP_ALIGN.CENTER, line=1.1)
        write(tb.text_frame, desc, size=14, color=BLUE_PALE,
              align=PP_ALIGN.CENTER, line=1.2)

    hline(s, (SW - Inches(3.4)) / 2, Inches(5.78), Inches(3.4), ORANGE, Pt(2.5))
    fb = textbox(s, M, Inches(6.00), CONTENT_W, Inches(0.9))
    tf = fb.text_frame
    write(tf, "منتورینگ گروهی", size=32, bold=True, color=ORANGE, first=True,
          align=PP_ALIGN.CENTER, line=1.15)
    write(tf, "جایی که با هم بلد می‌شویم", size=18, color=WHITE,
          align=PP_ALIGN.CENTER, line=1.2)

    notes(s, """
    اینجا یک استراحت کوتاه بدهید و بعد بخش دوم را شروع کنید.
    اگر جلسه فشرده است، مستقیم رد شوید ولی لحن را عوض کنید:
    از تحلیل به عمل.
    """)
    return s


# ═════════════════════════════════════ ساخت و اجرا ══════════════════════
INTRO_BUILDERS = [i01, i02, i03, i04, i05, i06, i07, i08, i09, i10,
                  i11, i12, i13, i14]


def _load_mentoring():
    """ماژول ارائه منتورینگ را بارگذاری می‌کند (بدون اجرای main)."""
    T.footer = _shifted_footer          # پاورقی با شماره‌گذاری پیوسته
    spec = importlib.util.spec_from_file_location(
        "gm_build", os.path.join(GM, "build.py"))
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


def build(cfg):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    T.set_font(cfg.font)

    for fn in INTRO_BUILDERS:
        fn(prs, cfg)

    gm = _load_mentoring()
    # اسلاید جلد منتورینگ (s01) حذف می‌شود چون ارائه جلد خودش را دارد
    for fn in gm.BUILDERS[1:]:
        fn(prs, cfg)

    prs.save(cfg.out)
    return cfg.out


def main():
    ap = argparse.ArgumentParser(
        description="ساخت ارائه ترکیبی توسعه سازمانی و منتورینگ گروهی")
    ap.add_argument("--out", default=os.path.join(
        HERE, "توسعه-سازمانی-و-منتورینگ-گروهی.pptx"))
    ap.add_argument("--font", default="IRANSans")
    ap.add_argument("--date", default="[تاریخ]")
    ap.add_argument("--presenter", default="وحید مجیدی")
    ap.add_argument("--org", default="[نام واحد / سازمان]")
    ap.add_argument("--first-session", dest="first_session", default="[تاریخ]")
    ap.add_argument("--time", default="[ساعت]")
    ap.add_argument("--place", default="[مکان]")
    ap.add_argument("--contact", default="[داخلی]")
    cfg = ap.parse_args()
    print("ساخته شد:", build(cfg))


if __name__ == "__main__":
    main()
