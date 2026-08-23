# -*- coding: utf-8 -*-
"""
رندر فایل PPTX به تصاویر PNG برای بازبینی چیدمان (بدون نیاز به LibreOffice).
شکل‌ها، گرادیانت‌ها، تصاویر و متن راست‌چین فارسی را تقریبی بازسازی می‌کند.

اجرا: python3 render.py <file.pptx> -o shots/ [--font-dir ~/.fonts]
"""
import argparse
import io
import os
import sys

from PIL import Image, ImageDraw, ImageFont, ImageFilter
from pptx import Presentation
from pptx.oxml.ns import qn

import arabic_reshaper
from bidi.algorithm import get_display

EMU_IN = 914400.0
DPI = 110.0


def px(v):
    return int(round((v or 0) / EMU_IN * DPI))


def pxf(v):
    return (v or 0) / EMU_IN * DPI


def shape_text(t):
    """آماده‌سازی متن فارسی برای رندر (اتصال حروف + ترتیب دوسویه)."""
    try:
        return get_display(arabic_reshaper.reshape(t))
    except Exception:
        return t


class Fonts:
    def __init__(self, font_dir):
        self.dir = os.path.expanduser(font_dir)
        self.cache = {}
        self.reg = self._find(["Vazirmatn-Regular.ttf", "DejaVuSans.ttf"])
        self.bold = self._find(["Vazirmatn-Bold.ttf", "DejaVuSans-Bold.ttf"])

    def _find(self, names):
        for n in names:
            p = os.path.join(self.dir, n)
            if os.path.exists(p):
                return p
        for root, _, files in os.walk("/usr/share/fonts"):
            for n in names:
                if n in files:
                    return os.path.join(root, n)
        return None

    def get(self, size_pt, bold=False):
        size = max(6, int(round(size_pt * DPI / 72.0)))
        key = (size, bold)
        if key not in self.cache:
            path = self.bold if bold else self.reg
            self.cache[key] = (ImageFont.truetype(path, size) if path
                               else ImageFont.load_default())
        return self.cache[key]


def color_of(el, default=None):
    if el is None:
        return default, 1.0
    srgb = el.find(qn("a:srgbClr"))
    if srgb is None:
        return default, 1.0
    hexv = srgb.get("val")
    a = srgb.find(qn("a:alpha"))
    alpha = int(a.get("val")) / 100000.0 if a is not None else 1.0
    return tuple(int(hexv[i:i + 2], 16) for i in (0, 2, 4)), alpha


def rounded_mask(size, radius):
    m = Image.new("L", size, 0)
    d = ImageDraw.Draw(m)
    d.rounded_rectangle([0, 0, size[0] - 1, size[1] - 1], radius=radius, fill=255)
    return m


def ellipse_mask(size):
    m = Image.new("L", size, 0)
    ImageDraw.Draw(m).ellipse([0, 0, size[0] - 1, size[1] - 1], fill=255)
    return m


def grad_image(size, stops, angle_deg):
    """گرادیانت خطی ساده (فقط دو رنگ)."""
    w, h = size
    (c1, a1), (c2, a2) = stops
    import math
    ang = math.radians((angle_deg) % 360)
    dx, dy = math.cos(ang), math.sin(ang)
    img = Image.new("RGBA", (w, h))
    px_ = img.load()
    # نمونه‌برداری کم‌هزینه: ساخت در ابعاد کوچک و بزرگ‌نمایی
    sw, sh = max(2, w // 6), max(2, h // 6)
    small = Image.new("RGBA", (sw, sh))
    sp = small.load()
    proj = [[0] * sw for _ in range(sh)]
    mn, mx = 1e18, -1e18
    for y in range(sh):
        for x in range(sw):
            t = (x / sw) * dx + (y / sh) * dy
            proj[y][x] = t
            mn = min(mn, t); mx = max(mx, t)
    rng = (mx - mn) or 1
    for y in range(sh):
        for x in range(sw):
            t = (proj[y][x] - mn) / rng
            sp[x, y] = (
                int(c1[0] + (c2[0] - c1[0]) * t),
                int(c1[1] + (c2[1] - c1[1]) * t),
                int(c1[2] + (c2[2] - c1[2]) * t),
                int(255 * (a1 + (a2 - a1) * t)),
            )
    return small.resize((w, h), Image.BILINEAR)


def geom(sp):
    spPr = getattr(sp._element, "spPr", None)
    if spPr is None:
        return None
    g = spPr.find(qn("a:prstGeom"))
    return g.get("prst") if g is not None else None


def adj_val(sp, idx, default):
    try:
        return sp.adjustments[idx]
    except Exception:
        return default


ALIGN = {None: "r", "r": "r", "l": "l", "ctr": "c", "just": "r"}
ANCHOR = {None: "t", "t": "t", "ctr": "c", "b": "b"}


def draw_shape(sp, canvas, fonts):
    x, y = px(sp.left), px(sp.top)
    w, h = max(1, px(sp.width)), max(1, px(sp.height))
    prst = geom(sp)

    # ---------- تصویر ----------
    if sp.shape_type is not None and int(sp.shape_type) == 13:
        try:
            im = Image.open(io.BytesIO(sp.image.blob)).convert("RGBA")
        except Exception:
            return
        cl = sp.crop_left or 0; cr = sp.crop_right or 0
        ct = sp.crop_top or 0; cb = sp.crop_bottom or 0
        iw, ih = im.size
        box = (int(iw * cl), int(ih * ct), int(iw * (1 - cr)), int(ih * (1 - cb)))
        if box[2] > box[0] and box[3] > box[1]:
            im = im.crop(box)
        im = im.resize((w, h), Image.LANCZOS)
        if prst == "roundRect":
            im.putalpha(rounded_mask((w, h), int(min(w, h) * 0.05)))
        canvas.alpha_composite(im, (x, y))
        return

    spPr = getattr(sp._element, "spPr", None)
    if spPr is None:
        return

    layer = Image.new("RGBA", (w, h), (0, 0, 0, 0))

    # ---------- پرکردن ----------
    solid = spPr.find(qn("a:solidFill"))
    grad = spPr.find(qn("a:gradFill"))
    if solid is not None:
        c, a = color_of(solid)
        if c:
            layer.paste((c[0], c[1], c[2], int(255 * a)), [0, 0, w, h])
    elif grad is not None:
        gsLst = grad.find(qn("a:gsLst"))
        stops = []
        for gs in gsLst:
            c, a = color_of(gs, (0, 0, 0))
            stops.append((c, a))
        lin = grad.find(qn("a:lin"))
        ang = int(lin.get("ang", "0")) / 60000.0 if lin is not None else 0
        if len(stops) >= 2:
            layer = grad_image((w, h), (stops[0], stops[-1]), ang)

    # ---------- ماسک شکل ----------
    if prst == "roundRect":
        r = int(min(w, h) * adj_val(sp, 0, 0.16))
        layer.putalpha(Image.composite(rounded_mask((w, h), r),
                                       Image.new("L", (w, h), 0),
                                       layer.split()[3].point(lambda v: 255 if v else 0))
                       if False else _mul_alpha(layer, rounded_mask((w, h), r)))
    elif prst in ("ellipse", "donut"):
        m = ellipse_mask((w, h))
        if prst == "donut":
            inner = adj_val(sp, 0, 0.25)
            d = ImageDraw.Draw(m)
            iw2 = int(w * inner); ih2 = int(h * inner)
            d.ellipse([iw2, ih2, w - iw2, h - ih2], fill=0)
        layer.putalpha(_mul_alpha(layer, m))
    elif prst == "blockArc":
        layer.putalpha(_mul_alpha(layer, _blockarc_mask(sp, w, h)))
    elif prst is None and spPr.find(qn("a:custGeom")) is not None:
        layer.putalpha(_mul_alpha(layer, _custgeom_mask(sp, spPr, w, h)))
    elif prst == "upArrow":
        m = Image.new("L", (w, h), 0)
        ImageDraw.Draw(m).polygon(
            [(w // 2, 0), (w, int(h * .55)), (int(w * .72), int(h * .55)),
             (int(w * .72), h), (int(w * .28), h), (int(w * .28), int(h * .55)),
             (0, int(h * .55))], fill=255)
        layer.putalpha(_mul_alpha(layer, m))

    canvas.alpha_composite(layer, (x, y))

    # ---------- خط دور ----------
    ln = spPr.find(qn("a:ln"))
    if ln is not None and ln.find(qn("a:noFill")) is None:
        sf = ln.find(qn("a:solidFill"))
        if sf is not None:
            c, a = color_of(sf)
            lw = max(1, int(round(int(ln.get("w", "12700")) / 12700.0 * DPI / 72.0)))
            d = ImageDraw.Draw(canvas)
            if prst == "roundRect":
                d.rounded_rectangle([x, y, x + w - 1, y + h - 1],
                                    radius=int(min(w, h) * adj_val(sp, 0, 0.16)),
                                    outline=(*c, int(255 * a)), width=lw)
            elif prst == "ellipse":
                d.ellipse([x, y, x + w - 1, y + h - 1],
                          outline=(*c, int(255 * a)), width=lw)
            else:
                d.rectangle([x, y, x + w - 1, y + h - 1],
                            outline=(*c, int(255 * a)), width=lw)

    draw_text(sp, canvas, fonts, x, y, w, h)


def _mul_alpha(layer, mask):
    a = layer.split()[3]
    return Image.composite(a, Image.new("L", a.size, 0), mask)


def _custgeom_mask(sp, spPr, w, h):
    """ماسک شکل آزاد (custGeom) از روی مسیرهای moveTo/lnTo."""
    cust = spPr.find(qn("a:custGeom"))
    m = Image.new("L", (w, h), 0)
    d = ImageDraw.Draw(m)
    pathLst = cust.find(qn("a:pathLst"))
    if pathLst is None:
        return m
    for path in pathLst:
        pw = int(path.get("w", "0")) or 1
        ph = int(path.get("h", "0")) or 1
        pts = []
        for node in path:
            tag = node.tag.split("}")[1]
            if tag in ("moveTo", "lnTo"):
                pt = node.find(qn("a:pt"))
                if pt is None:
                    continue
                pts.append((int(pt.get("x")) / pw * w, int(pt.get("y")) / ph * h))
        if len(pts) >= 3:
            d.polygon(pts, fill=255)
    return m


def _blockarc_mask(sp, w, h):
    import math
    m = Image.new("L", (w, h), 0)
    d = ImageDraw.Draw(m)
    a0 = adj_val(sp, 0, 0.0) * 360
    a1 = adj_val(sp, 1, 0.5) * 360
    inner = adj_val(sp, 2, 0.25)
    d.pieslice([0, 0, w - 1, h - 1], a0, a1 if a1 > a0 else a1 + 360, fill=255)
    iw = int(w * inner); ih = int(h * inner)
    d.ellipse([iw, ih, w - iw, h - ih], fill=0)
    return m


def draw_text(sp, canvas, fonts, bx, by, bw, bh):
    """رسم تقریبی متن راست‌چین چندپاراگرافی داخل کادر شکل."""
    if not sp.has_text_frame:
        return
    tf = sp.text_frame
    txBody = sp._element.find(qn("p:txBody"))
    bodyPr = txBody.find(qn("a:bodyPr")) if txBody is not None else None
    anchor = ANCHOR.get(bodyPr.get("anchor") if bodyPr is not None else None)
    ml, mr = px(tf.margin_left), px(tf.margin_right)
    mt, mb = px(tf.margin_top), px(tf.margin_bottom)
    aw = bw - ml - mr
    if aw <= 4:
        return

    d = ImageDraw.Draw(canvas)
    lines = []          # dict(segs, height, align, pad_top)

    for p in tf.paragraphs:
        runs = [(r.text, r.font) for r in p.runs if r.text]
        if not runs:
            continue
        pPr = p._p.find(qn("a:pPr"))
        algn = ALIGN.get(pPr.get("algn") if pPr is not None else None)
        sb = (p.space_before.pt if p.space_before else 0) * DPI / 72.0
        sa = (p.space_after.pt if p.space_after else 0) * DPI / 72.0
        ls = p.line_spacing or 1.2

        words = []
        for text, f in runs:
            size = f.size.pt if f.size else 18
            fnt = fonts.get(size, bool(f.bold))
            col = (31, 45, 56)
            try:
                if f.color and f.color.type is not None and f.color.rgb:
                    rgb = str(f.color.rgb)
                    col = tuple(int(rgb[i:i + 2], 16) for i in (0, 2, 4))
            except Exception:
                pass
            parts = text.split(" ")
            for j, wd in enumerate(parts):
                if wd:
                    words.append((wd, fnt, col, size))
                if j < len(parts) - 1:
                    words.append((" ", fnt, col, size))

        par_lines = []
        cur, curw, maxsz = [], 0.0, 0.0
        for wd, fnt, col, size in words:
            wpx = d.textlength(shape_text(wd), font=fnt)
            if cur and wd != " " and curw + wpx > aw:
                while cur and cur[-1][0] == " ":
                    curw -= d.textlength(" ", font=cur[-1][1])
                    cur.pop()
                par_lines.append((cur, maxsz, curw))
                cur, curw, maxsz = [], 0.0, 0.0
                if wd == " ":
                    continue
            cur.append((wd, fnt, col))
            curw += wpx
            maxsz = max(maxsz, size)
        if cur:
            par_lines.append((cur, maxsz, curw))

        for k, (segs, msz, tw) in enumerate(par_lines):
            lines.append({
                "segs": segs,
                "w": tw,
                "h": msz * DPI / 72.0 * ls,
                "align": algn,
                "pad_top": sb if k == 0 else 0.0,
                "pad_bot": sa if k == len(par_lines) - 1 else 0.0,
            })

    if not lines:
        return

    total = sum(l["h"] + l["pad_top"] + l["pad_bot"] for l in lines)
    if anchor == "c":
        y = by + mt + max(0.0, (bh - mt - mb - total) / 2)
    elif anchor == "b":
        y = by + bh - mb - total
    else:
        y = by + mt

    for l in lines:
        y += l["pad_top"]
        if l["align"] == "c":
            x = bx + ml + (aw - l["w"]) / 2
        elif l["align"] == "l":
            x = bx + ml
        else:
            x = bx + bw - mr - l["w"]
        cx = x
        for wd, fnt, col in reversed(l["segs"]):
            wpx = d.textlength(shape_text(wd), font=fnt)
            d.text((cx, y + l["h"] * 0.10), shape_text(wd), font=fnt, fill=col)
            cx += wpx
        y += l["h"] + l["pad_bot"]


def render(path, outdir, font_dir, scale=1.0, only=None):
    prs = Presentation(path)
    W, H = px(prs.slide_width), px(prs.slide_height)
    fonts = Fonts(font_dir)
    os.makedirs(outdir, exist_ok=True)
    paths = []
    for i, slide in enumerate(prs.slides, 1):
        if only and i not in only:
            continue
        canvas = Image.new("RGBA", (W, H), (255, 255, 255, 255))
        for sp in slide.shapes:
            try:
                draw_shape(sp, canvas, fonts)
            except Exception as e:
                sys.stderr.write("slide %d shape skipped: %s\n" % (i, e))
        out = os.path.join(outdir, "slide-%02d.png" % i)
        canvas.convert("RGB").save(out, quality=92)
        paths.append(out)
    return paths


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("-o", "--out", default="shots")
    ap.add_argument("--font-dir", default="~/.fonts")
    ap.add_argument("--only", default="")
    a = ap.parse_args()
    only = set(int(x) for x in a.only.split(",") if x.strip())
    for p in render(a.pptx, a.out, a.font_dir, only=only):
        print(p)
