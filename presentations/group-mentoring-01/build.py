# -*- coding: utf-8 -*-
"""
ساخت ارائه «از گفتگو تا سازمان یادگیرنده» — جلسه اول منتورینگ گروهی
اجرا:  python3 build.py [--out out.pptx] [--font IRANSans] [--date ...] [--presenter ...]
"""
import argparse
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "_lib"))
import theme as T
from theme import (
    NAVY, NAVY_DEEP, BLUE_MID, BLUE_LIGHT, BLUE_PALE, ORANGE, ORANGE_DEEP,
    ORANGE_PALE, WHITE, INK, INK_SOFT, PAPER, GREEN, GREEN_PALE, RED, RED_PALE,
    PURPLE, AMBER, LINE, SW, SH, M, CONTENT_W, CONTENT_TOP,
    fa, rect, card, gradient, hline, textbox, write, bullet, picture_fill,
    round_picture, base_slide, blank, header, footer, notes, no_line, shadow_off,
)

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")
TOTAL = 12


def img(name):
    p = os.path.join(IMG, name)
    return p, Image.open(p).size


# ═════════════════════════════════════════════════════ ۱ — جلد ═══════════
def slide_01(prs, cfg):
    s = blank(prs)
    gradient(s, 0, 0, SW, SH, NAVY_DEEP, BLUE_MID, angle=315)

    # تصویر سمت چپ با محو شدن تدریجی به سمت راست
    pw = Inches(6.1)
    p, sz = img("01-hands-together.jpg")
    picture_fill(s, p, SW - pw, 0, pw, SH, sz)
    gradient(s, SW - pw, 0, pw, SH, NAVY_DEEP, NAVY_DEEP,
             angle=0, alpha1=0.97, alpha2=0.10)

    # ستون متن (راست‌چین، سمت راست کادر متن)
    tx = Inches(0.85)
    tw = Inches(6.6)
    hline(s, tx + tw - Inches(1.5), Inches(1.62), Inches(1.5), ORANGE, Pt(5))

    tb = textbox(s, tx, Inches(1.88), tw, Inches(3.6))
    tf = tb.text_frame
    write(tf, "جلسه اول منتورینگ گروهی", size=19, bold=True,
          color=BLUE_LIGHT, first=True, line=1.0)
    write(tf, "از گفتگو", size=50, bold=True, color=WHITE,
          space_before=16, line=1.22)
    write(tf, [("تا ", True, WHITE, 50), ("سازمان یادگیرنده", True, ORANGE, 50)],
          line=1.22)
    write(tf, "نقشه راه توسعه فردی و سازمانی", size=20, color=BLUE_PALE,
          space_before=22, line=1.3)
    write(tf, "با رویکرد منتورینگ گروهی", size=20, color=BLUE_PALE, line=1.3)

    hline(s, tx + tw - Inches(3.4), Inches(5.62), Inches(3.4), BLUE_LIGHT, Pt(1.25))

    mb = textbox(s, tx, Inches(5.85), tw, Inches(1.0))
    write(mb.text_frame,
          [("تاریخ:  ", False, BLUE_LIGHT, 17), (cfg.date, True, WHITE, 17),
           ("        ارائه‌کننده:  ", False, BLUE_LIGHT, 17),
           (cfg.presenter, True, WHITE, 17)],
          first=True, line=1.2)
    write(mb.text_frame, cfg.org, size=14, color=BLUE_LIGHT, space_before=6, line=1.2)

    notes(s, """
    شروع بدون مقدمه طولانی. لبخند، تماس چشمی، و یک جمله ساده:
    «امروز نیامده‌ام چیزی به شما یاد بدهم؛ آمده‌ام یک پیشنهاد بدهم.»
    زمان: حدود ۱ دقیقه.
    """)
    return s


# ═══════════════════════════════════ ۲ — یک سؤال برای شروع ═══════════════
def slide_02(prs, cfg):
    s = base_slide(prs, WHITE)
    pw = Inches(5.25)
    p, sz = img("02-thinking.jpg")
    pic = picture_fill(s, p, 0, 0, pw, SH, sz)
    gradient(s, 0, 0, pw, SH, NAVY_DEEP, NAVY_DEEP, angle=0,
             alpha1=0.05, alpha2=0.55)

    tx = pw + Inches(0.7)
    tw = SW - tx - M

    kb = textbox(s, tx, Inches(0.95), tw, Inches(0.4))
    write(kb.text_frame, "یک سؤال ساده…", size=19, bold=True, color=ORANGE,
          first=True, line=1.0)

    qb = textbox(s, tx, Inches(1.55), tw, Inches(2.6))
    tf = qb.text_frame
    write(tf, [("آخرین بار ", False, INK, 34), ("کِی", True, ORANGE, 34),
               (" با همکارانتان یک ", False, INK, 34),
               ("گفتگوی واقعی", True, NAVY, 34),
               (" درباره یک ", False, INK, 34),
               ("چالش واقعی", True, NAVY, 34), (" داشتید؟", False, INK, 34)],
          first=True, line=1.35)

    yb = textbox(s, tx, Inches(3.95), tw, Inches(1.5))
    tf = yb.text_frame
    write(tf, "نه گزارش دادن.", size=24, color=INK_SOFT, first=True, line=1.35)
    write(tf, "نه دستور گرفتن.", size=24, color=INK_SOFT, line=1.35)
    write(tf, "گفتگوی واقعی.", size=27, bold=True, color=NAVY, line=1.35)

    hline(s, SW - M - Inches(2.6), Inches(5.62), Inches(2.6), ORANGE, Pt(3))

    fb = textbox(s, tx, Inches(5.85), tw, Inches(1.0))
    tf = fb.text_frame
    write(tf, "اگر یادتان نمی‌آید، شما تنها نیستید.", size=20, color=INK, first=True)
    write(tf, "و همین‌جاست که مشکل شروع می‌شود.", size=20, bold=True, color=ORANGE_DEEP)

    footer(s, 2, TOTAL)
    notes(s, """
    سؤال را بپرسید و ۵ ثانیه سکوت کنید. اجازه دهید سنگینی سؤال حس شود.
    اگر کسی خواست جواب دهد، بشنوید — همین شروع گفتگوست.
    """)
    return s


# ═══════════════════════════════════════ ۳ — واقعیت امروز ما ═════════════
def slide_03(prs, cfg):
    s = base_slide(prs, PAPER)
    header(s, "ما کجای مسیر هستیم؟", kicker="واقعیت امروز ما")

    gap = Inches(0.42)
    cw = (CONTENT_W - gap) / 2
    top = Inches(1.78)
    ch = Inches(0.62)          # ارتفاع سرستون
    rowh = Inches(0.62)
    rows = 4
    boxh = ch + rows * rowh + Inches(0.26)

    cols = [
        # (x, عنوان, رنگ سر, رنگ پس‌زمینه, نشانه, آیتم‌ها)
        (M + cw + gap, "آنچه داریم", GREEN, GREEN_PALE, "✓",
         ["تخصص فنی بالا", "سال‌ها تجربه", "نیروی متعهد", "جلسات فراوان"]),
        (M, "آنچه نداریم", RED, RED_PALE, "✕",
         ["فضای گفتگوی آزاد", "فرصت یادگیری از هم", "احساس شنیده شدن",
          "جلسات مفید و صمیمی"]),
    ]

    for x, title, cmain, cpale, mark, items in cols:
        card(s, x, top, cw, boxh, fill=WHITE, line=LINE)
        head = rect(s, x, top, cw, ch, fill=cmain,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
        rect(s, x, top + ch - Inches(0.16), cw, Inches(0.16), fill=cmain)
        tf = head.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(tf, [(mark + "   ", True, WHITE, 20), (title, True, WHITE, 23)],
              first=True, align=PP_ALIGN.CENTER, line=1.0)

        for i, it in enumerate(items):
            y = top + ch + Inches(0.13) + i * rowh
            if i % 2 == 0:
                rect(s, x + Inches(0.14), y, cw - Inches(0.28), rowh - Inches(0.06),
                     fill=cpale, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
            tb = textbox(s, x + Inches(0.34), y, cw - Inches(0.68),
                         rowh - Inches(0.06), anchor=MSO_ANCHOR.MIDDLE)
            write(tb.text_frame, [(mark + "  ", True, cmain, 16), (it, False, INK, 21)],
                  first=True, line=1.0)

    # نوار نتیجه
    ry = top + boxh + Inches(0.30)
    band = rect(s, M, ry, CONTENT_W, Inches(0.98), fill=NAVY,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    rect(s, SW - M - Inches(0.09), ry, Inches(0.09), Inches(0.98), fill=ORANGE)
    tb = textbox(s, M + Inches(0.3), ry, CONTENT_W - Inches(0.72), Inches(0.98),
                 anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame,
          [("نتیجه:   ", True, ORANGE, 20),
           ("دانش در سر افراد ", False, WHITE, 23), ("حبس", True, ORANGE, 23),
           (" شده، اشتباهات ", False, WHITE, 23), ("تکرار", True, ORANGE, 23),
           (" می‌شوند و نوآوری‌ها ", False, WHITE, 23), ("خفه", True, ORANGE, 23),
           (" می‌شوند.", False, WHITE, 23)],
          first=True, line=1.15)

    footer(s, 3, TOTAL)
    notes(s, """
    ستون سمت راست را با افتخار بخوانید و ستون چپ را بدون سرزنش.
    پیام: مشکل از آدم‌ها نیست، از نبودِ فضاست.
    """)
    return s


# ═══════════════════════════════════════════════ ۴ — نقشه راه ════════════
def slide_04(prs, cfg):
    s = base_slide(prs, PAPER)
    header(s, "ما در کجای مسیر تحول ایستاده‌ایم؟", kicker="نقشه راه")

    # ستون پله‌ها (سمت راست)
    colx = Inches(5.55)
    colw = SW - M - colx
    waves = [
        ("موج ۵", "انسان‌محور", BLUE_LIGHT, None),
        ("موج ۴", "دیجیتال", BLUE_MID, None),
        ("موج ۳", "دانش‌محور", NAVY, "مقصد ۲ ساله ما"),
        ("موج ۲", "صنعتی", ORANGE, "اینجا هستیم"),
        ("موج ۱", "پیشامدرن", INK_SOFT, None),
    ]
    top = Inches(1.80)
    rh = Inches(0.88)
    step = Inches(0.30)          # پله‌ای شدن به سمت بالا

    for i, (num, name, color, tag) in enumerate(waves):
        x = colx + step * i
        w = colw - step * i
        y = top + i * (rh + Inches(0.15))
        rect(s, x, y, w, rh, fill=color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
        tb = textbox(s, x + Inches(0.24), y, w - Inches(0.48), rh,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame,
              [(num + " — ", True, WHITE, 18), (name, True, WHITE, 24)],
              first=True, line=1.0)
        if tag:
            # برچسب داخل خود نوار، سمت چپ (تا با ستون کناری برخورد نکند)
            lw = Inches(2.05)
            lx = x + Inches(0.16)
            lbl = rect(s, lx, y + Inches(0.16), lw, rh - Inches(0.32),
                       fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.30)
            ltf = lbl.text_frame
            ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
            write(ltf, tag, size=15, bold=True, color=color, first=True,
                  align=PP_ALIGN.CENTER, line=1.0)
        # فلش صعود بین موج‌ها
        if i < len(waves) - 1:
            ax = x + w - Inches(1.15)
            ar = s.shapes.add_shape(MSO_SHAPE.UP_ARROW, ax,
                                    y - Inches(0.16), Inches(0.24), Inches(0.19))
            ar.fill.solid(); ar.fill.fore_color.rgb = LINE
            no_line(ar); shadow_off(ar)

    # ستون پیام (سمت چپ) + تصویر
    px = M
    pwid = colx - M - Inches(0.45)
    p, sz = img("04-climb.jpg")
    pic = picture_fill(s, p, px, Inches(1.80), pwid, Inches(2.45), sz)
    round_picture(pic)

    card(s, px, Inches(4.45), pwid, Inches(2.30), fill=NAVY, line=None)
    rect(s, px + pwid - Inches(0.09), Inches(4.45), Inches(0.09), Inches(2.30),
         fill=ORANGE)
    tb = textbox(s, px + Inches(0.28), Inches(4.68), pwid - Inches(0.62),
                 Inches(1.90))
    tf = tb.text_frame
    write(tf, "پیام کلیدی", size=16, bold=True, color=ORANGE, first=True, line=1.0)
    write(tf, [("پرش از موج ۲ به موج ۳ نیازمند ", False, WHITE, 21),
               ("یک چیز", True, ORANGE, 21), (" است:", False, WHITE, 21)],
          space_before=10, line=1.25)
    write(tf, "گفتگوی واقعی بین آدم‌ها.", size=25, bold=True, color=WHITE,
          space_before=6, line=1.2)
    write(tf, "نه بخشنامه، نه نرم‌افزار، نه ساختار جدید.",
          size=16, color=BLUE_LIGHT, space_before=8, line=1.2)

    footer(s, 4, TOTAL)
    notes(s, """
    از پایین به بالا بخوانید. روی «اینجا هستیم» مکث کنید.
    تأکید: فاصله موج ۲ تا ۳ فاصله فناوری نیست، فاصله فرهنگ گفتگوست.
    """)
    return s


# ═══════════════════════════════ ۵ — چیستی منتورینگ گروهی ════════════════
def slide_05(prs, cfg):
    s = base_slide(prs, WHITE)
    header(s, "منتورینگ گروهی چیست؟", kicker="تعریف")

    pw = Inches(5.15)
    p, sz = img("05-circle.jpg")
    pic = picture_fill(s, p, M, Inches(1.80), pw, Inches(4.05), sz)
    round_picture(pic)

    # کادر تعریف ساده روی تصویر
    qb = card(s, M + Inches(0.25), Inches(4.62), pw - Inches(0.5), Inches(1.05),
              fill=NAVY, line=None)
    tf = qb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(tf, "«فضایی امن برای فکر کردن با هم —", size=18, bold=True,
          color=WHITE, first=True, align=PP_ALIGN.CENTER, line=1.2)
    write(tf, "به‌جای فکر کردن تنها»", size=18, bold=True, color=ORANGE,
          align=PP_ALIGN.CENTER, line=1.2)

    # ستون راست
    tx = M + pw + Inches(0.55)
    tw = SW - M - tx

    head = card(s, tx, Inches(1.80), tw, Inches(1.02), fill=BLUE_PALE, line=None)
    rect(s, tx + tw - Inches(0.09), Inches(1.80), Inches(0.09), Inches(1.02),
         fill=NAVY)
    tf = head.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tb = textbox(s, tx + Inches(0.3), Inches(1.80), tw - Inches(0.62), Inches(1.02),
                 anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame,
          [("جمع شدن ", False, INK, 21), ("۸ تا ۱۲ نفر", True, NAVY, 21),
           (" از همکاران در یک جلسه ", False, INK, 21),
           ("۹۰ دقیقه‌ای", True, NAVY, 21), ("، برای اینکه:", False, INK, 21)],
          first=True, line=1.3)

    items = [
        ("تجربه‌های واقعی خود را ", "به اشتراک", " بگذارند"),
        ("به‌جای نصیحت، ", "سؤال", " بپرسند"),
        ("به‌جای قضاوت، ", "گوش", " بدهند"),
        ("به‌جای رقابت، ", "همکاری", " کنند"),
    ]
    y0 = Inches(3.05)
    rh = Inches(0.72)
    for i, (a, b, c) in enumerate(items):
        y = y0 + i * (rh + Inches(0.12))
        rect(s, tx, y, tw, rh, fill=PAPER,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
        dot = rect(s, tx + tw - Inches(0.62), y + Inches(0.19), Inches(0.34),
                   Inches(0.34), fill=ORANGE, shape=MSO_SHAPE.OVAL)
        dtf = dot.text_frame
        dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(dtf, fa(i + 1), size=14, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, tx + Inches(0.28), y, tw - Inches(1.0), rh,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame,
              [(a, False, INK, 21), (b, True, ORANGE_DEEP, 21), (c, False, INK, 21)],
              first=True, line=1.1)

    footer(s, 5, TOTAL)
    notes(s, """
    تأکید کنید که این «کلاس» نیست و «جلسه گزارش‌دهی» هم نیست.
    مثال بزنید: یک چالش واقعی که خودتان تنها با آن دست‌وپنجه نرم کرده‌اید.
    """)
    return s


# ═══════════════════════════ ۶ — چرا کار می‌کند؟ (۷۰-۲۰-۱۰) ══════════════
def slide_06(prs, cfg):
    s = base_slide(prs, PAPER)
    header(s, "چگونه واقعاً یاد می‌گیریم؟", kicker="علم پشت ماجرا")

    # نمودار حلقه‌ای سمت راست
    cx = Inches(9.85)
    cy = Inches(3.90)
    R = Inches(1.70)
    R_IN = Inches(1.02)
    GAP = 2.0                     # فاصله زاویه‌ای بین قطاع‌ها (درجه)

    segs = [(0.0, 252.0, GREEN, 0.0),          # ۷۰٪
            (252.0, 324.0, ORANGE, 0.13),      # ۲۰٪ (کمی بیرون‌زده)
            (324.0, 360.0, BLUE_MID, 0.0)]     # ۱۰٪
    for a0, a1, color, pop in segs:
        r_out = R + Inches(pop)
        r_in = R_IN + Inches(pop * 0.5)
        T.donut_segment(s, cx, cy, r_out, r_in, a0 + GAP, a1 - GAP, color)

    ctb = textbox(s, cx - Inches(1.05), cy - Inches(0.55), Inches(2.1),
                  Inches(1.1), anchor=MSO_ANCHOR.MIDDLE)
    write(ctb.text_frame, "۷۰-۲۰-۱۰", size=25, bold=True, color=NAVY,
          first=True, align=PP_ALIGN.CENTER, line=1.0)
    write(ctb.text_frame, "مدل یادگیری", size=13, color=INK_SOFT,
          align=PP_ALIGN.CENTER, line=1.1)

    src = textbox(s, cx - Inches(2.1), cy + R + Inches(0.34), Inches(4.2),
                  Inches(0.4))
    write(src.text_frame, "مرجع: مؤسسه CCL آمریکا", size=13, color=INK_SOFT,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    # ردیف‌های توضیح سمت راست‌ترِ متن (چپ اسلاید)
    tx = M
    tw = Inches(7.35)
    rows = [
        ("۷۰٪", "یادگیری از انجام کار واقعی", GREEN, GREEN_PALE, False),
        ("۲۰٪", "یادگیری از تعامل با دیگران", ORANGE, ORANGE_PALE, True),
        ("۱۰٪", "یادگیری از کلاس و کتاب", BLUE_MID, BLUE_PALE, False),
    ]
    y = Inches(1.82)
    for pct, label, color, pale, hi in rows:
        h = Inches(1.02) if hi else Inches(0.86)
        rect(s, tx, y, tw, h, fill=pale, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.12, line=color if hi else None, line_w=Pt(2))
        rect(s, tx + tw - Inches(0.09), y, Inches(0.09), h, fill=color)
        pb = textbox(s, tx + tw - Inches(1.65), y, Inches(1.45), h,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(pb.text_frame, pct, size=31 if hi else 27, bold=True, color=color,
              first=True, align=PP_ALIGN.CENTER, line=1.0)
        lb = textbox(s, tx + Inches(0.3), y, tw - Inches(2.1), h,
                     anchor=MSO_ANCHOR.MIDDLE)
        chunks = [(label, hi, INK if not hi else ORANGE_DEEP, 23 if hi else 21)]
        if hi:
            chunks.append(("      ← اینجا!", True, ORANGE, 17))
        write(lb.text_frame, chunks, first=True, line=1.15)
        y += h + Inches(0.16)

    # واقعیت تلخ
    ry = Inches(5.32)
    rect(s, M, ry, Inches(7.35), Inches(1.42), fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    rect(s, M + Inches(7.35) - Inches(0.09), ry, Inches(0.09), Inches(1.42),
         fill=ORANGE)
    tb = textbox(s, M + Inches(0.3), ry + Inches(0.16), Inches(6.75), Inches(1.1))
    tf = tb.text_frame
    write(tf, "واقعیت تلخ", size=15, bold=True, color=ORANGE, first=True, line=1.0)
    write(tf, [("سازمان ما ", False, WHITE, 20), ("۹۰٪", True, ORANGE, 20),
               (" بودجه آموزشی را روی همان ", False, WHITE, 20),
               ("۱۰٪", True, ORANGE, 20), (" خرج می‌کند.", False, WHITE, 20)],
          space_before=7, line=1.2)
    write(tf, "منتورینگ گروهی، همان ۲۰٪ گمشده است.", size=20, bold=True,
          color=WHITE, space_before=5, line=1.2)

    footer(s, 6, TOTAL)
    notes(s, """
    این اسلاید «مجوز علمی» طرح است. عدد ۷۰-۲۰-۱۰ را ساده توضیح دهید.
    جمله پایانی را آرام و با مکث بگویید.
    """)
    return s


# ═══════════════════════════════════ ۷ — سه‌گانه رشد فردی (KSA) ══════════
def slide_07(prs, cfg):
    s = base_slide(prs, WHITE)
    header(s, "هر یادگیری واقعی، سه بُعد دارد", kicker="سه‌گانه رشد — K S A")

    # نمودار ون سمت راست
    cx = Inches(9.95)
    cy = Inches(4.10)
    r = Inches(1.42)
    off = Inches(0.80)
    circles = [
        (cx, cy - off, BLUE_MID, "K", "دانش", "می‌دانم"),
        (cx + off * 1.05, cy + off * 0.72, AMBER, "S", "مهارت", "می‌توانم"),
        (cx - off * 1.05, cy + off * 0.72, RED, "A", "نگرش", "می‌خواهم"),
    ]
    for ccx, ccy, color, letter, name, verb in circles:
        sp = rect(s, ccx - r, ccy - r, 2 * r, 2 * r, fill=color,
                  shape=MSO_SHAPE.OVAL, alpha=0.62)
    for ccx, ccy, color, letter, name, verb in circles:
        lx = ccx - Inches(0.9)
        ly = ccy - r + Inches(0.16) if ccy < cy else ccy + r - Inches(1.05)
        tb = textbox(s, lx, ly, Inches(1.8), Inches(0.9))
        tf = tb.text_frame
        write(tf, letter, size=24, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        write(tf, name, size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
              line=1.05)
        write(tf, "«%s»" % verb, size=13, color=WHITE, align=PP_ALIGN.CENTER,
              line=1.05)

    core = textbox(s, cx - Inches(0.85), cy + Inches(0.14), Inches(1.7),
                   Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
    write(core.text_frame, "رشد واقعی", size=19, bold=True, color=NAVY,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    # جدول سمت چپ
    tx = M
    tw = Inches(7.15)
    hd = rect(s, tx, Inches(1.82), tw, Inches(0.68), fill=NAVY,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
    tb = textbox(s, tx + Inches(0.25), Inches(1.82), tw - Inches(0.5),
                 Inches(0.68), anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame, "در جلسات منتورینگ گروهی، هر بُعد چطور رشد می‌کند؟",
          size=19, bold=True, color=WHITE, first=True, line=1.0)

    rows = [
        ("K", "دانش", BLUE_MID, "از حرف‌ها و تجربیات واقعی دیگران"),
        ("S", "مهارت", AMBER, "از تمرین گفتگو و پرسشگری"),
        ("A", "نگرش", RED, "از حس تعلق و شنیده شدن"),
    ]
    y = Inches(2.72)
    rh = Inches(1.04)
    for letter, name, color, desc in rows:
        rect(s, tx, y, tw, rh, fill=PAPER, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.12)
        chip = rect(s, tx + tw - Inches(1.24), y + Inches(0.15), Inches(1.02),
                    rh - Inches(0.30), fill=color,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.22)
        ctf = chip.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(ctf, letter, size=25, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, tx + Inches(0.3), y, tw - Inches(1.7), rh,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, name, size=15, bold=True, color=color, first=True,
              line=1.0)
        write(tb.text_frame, desc, size=21, color=INK, line=1.15)
        y += rh + Inches(0.14)

    hline(s, tx, Inches(6.16), tw, LINE, Pt(1.25))
    nb = textbox(s, tx, Inches(6.32), tw, Inches(0.5))
    write(nb.text_frame,
          [("فقط ", False, INK_SOFT, 17), ("K", True, NAVY, 17),
           ("؟ می‌شود آموزش. هر سه با هم؟ می‌شود ", False, INK_SOFT, 17),
           ("توسعه", True, ORANGE_DEEP, 17), (".", False, INK_SOFT, 17)],
          first=True, line=1.1)

    footer(s, 7, TOTAL)
    notes(s, """
    یک مثال شخصی برای هر بُعد بگویید: چیزی که دانستید، مهارتی که تمرین کردید،
    و نگرشی که در شما عوض شد.
    """)
    return s


# ═════════════════════════════ ۸ — چه چیزی به دست می‌آوریم؟ ══════════════
def slide_08(prs, cfg):
    s = base_slide(prs, PAPER)
    header(s, "دستاوردهای این طرح", kicker="چه چیزی به دست می‌آوریم؟")

    gap = Inches(0.42)
    cw = (CONTENT_W - gap) / 2
    top = Inches(1.78)
    ch = Inches(0.66)
    rowh = Inches(0.60)
    rows = 5
    boxh = ch + rows * rowh + Inches(0.24)

    cols = [
        (M + cw + gap, "برای من (فردی)", NAVY, BLUE_PALE,
         ["مهارت حل مسئله", "اعتمادبه‌نفس بیشتر", "کاهش استرس شغلی",
          "گسترش شبکه ارتباطی", "رشد شغلی سریع‌تر"]),
        (M, "برای سازمان", ORANGE_DEEP, ORANGE_PALE,
         ["انتقال دانش بین‌نسلی", "کاهش تکرار اشتباهات", "افزایش چابکی",
          "نگهداشت نیروها", "تغییر فرهنگ سازمانی"]),
    ]

    for x, title, cmain, cpale, items in cols:
        card(s, x, top, cw, boxh, fill=WHITE, line=LINE)
        rect(s, x, top, cw, ch, fill=cmain, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.11)
        rect(s, x, top + ch - Inches(0.16), cw, Inches(0.16), fill=cmain)
        tb = textbox(s, x, top, cw, ch, anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=23, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        for i, it in enumerate(items):
            y = top + ch + Inches(0.12) + i * rowh
            if i % 2 == 0:
                rect(s, x + Inches(0.14), y, cw - Inches(0.28), rowh - Inches(0.05),
                     fill=cpale, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
            tb = textbox(s, x + Inches(0.34), y, cw - Inches(0.68),
                         rowh - Inches(0.05), anchor=MSO_ANCHOR.MIDDLE)
            write(tb.text_frame, [("✓  ", True, cmain, 17), (it, False, INK, 21)],
                  first=True, line=1.0)

    by = top + boxh + Inches(0.26)
    rect(s, M, by, CONTENT_W, Inches(0.84), fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11)
    tb = textbox(s, M, by, CONTENT_W, Inches(0.84), anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame,
          [("همه این‌ها، با فقط  ", False, WHITE, 23),
           ("۹۰ دقیقه در هفته", True, ORANGE, 25)],
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    footer(s, 8, TOTAL)
    notes(s, """
    ستون سازمان را برای مدیران و ستون فردی را برای کارشناسان پررنگ کنید.
    جمله «۹۰ دقیقه در هفته» را به‌عنوان کوچک‌ترین هزینه ممکن قاب بگیرید.
    """)
    return s


# ═════════════════════════════ ۹ — ساختار جلسه ۹۰ دقیقه‌ای ═══════════════
def slide_09(prs, cfg):
    s = base_slide(prs, WHITE)
    header(s, "ساختار هر جلسه — ۹۰ دقیقه", kicker="چگونه اجرا می‌شود؟")

    steps = [
        (10, "خوش‌آمد و گرم‌کردن", "یک سؤال ساده", BLUE_LIGHT),
        (10, "مرور جلسه قبل", "تعهدات هفته گذشته", BLUE_LIGHT),
        (15, "ارائه چالش", "یک نفر داوطلب", BLUE_MID),
        (15, "طوفان سؤالات", "فقط سؤال، بدون راه‌حل", ORANGE),
        (15, "تبادل تجربه", "همه ایده بدهند", ORANGE),
        (15, "تعهد فردی", "۳ اقدام مشخص", NAVY),
        (10, "جمع‌بندی", "یک کلمه از هر نفر", NAVY),
    ]

    # نوار پیشرفت افقی (راست به چپ)
    bar_y = Inches(1.80)
    bar_h = Inches(0.42)
    total_min = sum(x[0] for x in steps)
    xcur = SW - M
    for mins, title, desc, color in steps:
        w = CONTENT_W * mins / total_min
        rect(s, xcur - w, bar_y, w - Inches(0.035), bar_h, fill=color)
        tb = textbox(s, xcur - w, bar_y, w, bar_h, anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, fa(mins), size=14, bold=True, color=WHITE,
              first=True, align=PP_ALIGN.CENTER, line=1.0)
        xcur -= w

    # ردیف‌های جدول
    y = Inches(2.40)
    rh = Inches(0.505)
    for i, (mins, title, desc, color) in enumerate(steps):
        if i % 2 == 0:
            rect(s, M, y, CONTENT_W, rh - Inches(0.04), fill=PAPER,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.20)
        chip = rect(s, SW - M - Inches(1.30), y + Inches(0.045), Inches(1.14),
                    rh - Inches(0.14), fill=color,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.26)
        ctf = chip.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(ctf, fa("%d دقیقه" % mins), size=14, bold=True, color=WHITE,
              first=True, align=PP_ALIGN.CENTER, line=1.0)

        tb = textbox(s, SW - M - Inches(5.35), y, Inches(3.9), rh - Inches(0.04),
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=21, bold=True, color=NAVY, first=True,
              line=1.0)

        db = textbox(s, M + Inches(0.2), y, Inches(7.5), rh - Inches(0.04),
                     anchor=MSO_ANCHOR.MIDDLE)
        write(db.text_frame, desc, size=19, color=INK_SOFT, first=True, line=1.0)
        y += rh

    gy = Inches(6.06)
    rect(s, M, gy, CONTENT_W, Inches(0.70), fill=ORANGE_PALE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.13, line=ORANGE, line_w=Pt(1.75))
    tb = textbox(s, M, gy, CONTENT_W, Inches(0.70), anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame,
          [("قانون طلایی:   ", True, ORANGE_DEEP, 18),
           ("۷۰٪ گروه صحبت می‌کند، ۳۰٪ تسهیل‌گر.", True, NAVY, 21)],
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    footer(s, 9, TOTAL)
    notes(s, """
    تأکید بر «طوفان سؤالات»: در آن ۱۵ دقیقه هیچ‌کس اجازه راه‌حل دادن ندارد.
    این عجیب‌ترین و مؤثرترین بخش جلسه است.
    """)
    return s


# ══════════════════════════════════════ ۱۰ — قوانین طلایی ════════════════
def slide_10(prs, cfg):
    s = base_slide(prs, PAPER)
    header(s, "قواعد بازی — تا فضا امن باشد", kicker="قوانین طلایی جلسه")

    rules = [
        ("۱", "سلسله‌مراتب تعطیل است", "همه هم‌فکریم، نه رئیس و زیردست", ORANGE, ORANGE_PALE),
        ("۲", "محرمانگی کامل", "هر چه اینجا گفته شود، همین‌جا می‌ماند", NAVY, BLUE_PALE),
        ("۳", "سؤال بپرس، جواب نده", "به‌جای نصیحت، کنجکاوی", GREEN, GREEN_PALE),
        ("۴", "قضاوت ممنوع", "هیچ سؤالی احمقانه نیست", PURPLE, T.RGBColor(0xF6, 0xEF, 0xFA)),
        ("۵", "گوش بده تا بفهمی", "نه اینکه گوش بدهی تا جواب بدهی", RED, RED_PALE),
    ]

    tw = Inches(8.25)
    tx = SW - M - tw
    y = Inches(1.78)
    rh = Inches(0.92)
    for num, title, desc, color, pale in rules:
        rect(s, tx, y, tw, rh, fill=pale, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.11)
        rect(s, tx + tw - Inches(0.10), y, Inches(0.10), rh, fill=color)
        badge = rect(s, tx + tw - Inches(0.96), y + Inches(0.22), Inches(0.50),
                     Inches(0.50), fill=color, shape=MSO_SHAPE.OVAL)
        btf = badge.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(btf, num, size=18, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, tx + Inches(0.3), y, tw - Inches(1.35), rh,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=23, bold=True, color=color, first=True,
              line=1.05)
        write(tb.text_frame, desc, size=17, color=INK_SOFT, line=1.15)
        y += rh + Inches(0.10)

    # ستون تصویر سمت چپ
    px = M
    pw = tx - M - Inches(0.42)
    p, sz = img("10-safe.jpg")
    pic = picture_fill(s, p, px, Inches(1.78), pw, Inches(3.28), sz)
    round_picture(pic)

    card(s, px, Inches(5.24), pw, Inches(1.55), fill=NAVY, line=None)
    rect(s, px + pw - Inches(0.09), Inches(5.24), Inches(0.09), Inches(1.55),
         fill=ORANGE)
    tb = textbox(s, px + Inches(0.26), Inches(5.46), pw - Inches(0.55), Inches(1.15))
    tf = tb.text_frame
    write(tf, "چرا این قوانین؟", size=15, bold=True, color=ORANGE, first=True,
          line=1.0)
    write(tf, "بدون امنیت روانی، هیچ‌کس از شکست‌هایش حرف نمی‌زند — و یادگیری آنجاست.",
          size=16, color=WHITE, space_before=7, line=1.25)

    footer(s, 10, TOTAL)
    notes(s, """
    این ۵ قانون را در جلسه اول با صدای بلند بخوانید و از گروه تأیید بگیرید.
    بهتر است روی دیوار اتاق نصب شود.
    """)
    return s


# ══════════════════════════════════ ۱۱ — از یادگیری به اقدام ═════════════
def slide_11(prs, cfg):
    s = base_slide(prs, WHITE)
    header(s, "پایان هر جلسه: سه تعهد شخصی", kicker="از یادگیری به اقدام")

    cards = [
        ("K", "دانش", BLUE_MID, BLUE_PALE, "یک چیز جدید یاد می‌گیرم",
         "مثلاً: از تجربه فلانی درباره آن پروژه می‌پرسم"),
        ("S", "مهارت", AMBER, T.RGBColor(0xFD, 0xF3, 0xE2), "یک رفتار جدید تمرین می‌کنم",
         "مثلاً: به‌جای دستور دادن، سؤال می‌پرسم"),
        ("A", "نگرش", RED, RED_PALE, "یک نگرش را تغییر می‌دهم",
         "مثلاً: از یک همکار صادقانه تشکر می‌کنم"),
    ]

    gap = Inches(0.34)
    cw = (CONTENT_W - 2 * gap) / 3
    top = Inches(1.80)
    chh = Inches(3.30)

    for i, (letter, name, color, pale, title, example) in enumerate(cards):
        x = SW - M - cw - i * (cw + gap)     # چیدمان راست به چپ
        card(s, x, top, cw, chh, fill=WHITE, line=LINE)
        rect(s, x, top, cw, Inches(0.92), fill=color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
        rect(s, x, top + Inches(0.76), cw, Inches(0.16), fill=color)
        hb = textbox(s, x, top, cw, Inches(0.92), anchor=MSO_ANCHOR.MIDDLE)
        write(hb.text_frame,
              [(letter, True, WHITE, 30), ("   ", False, WHITE, 20),
               (name, True, WHITE, 21)],
              first=True, align=PP_ALIGN.CENTER, line=1.0)

        tb = textbox(s, x + Inches(0.26), top + Inches(1.14), cw - Inches(0.52),
                     Inches(1.0))
        write(tb.text_frame, title, size=22, bold=True, color=NAVY, first=True,
              align=PP_ALIGN.CENTER, line=1.25)

        rect(s, x + Inches(0.22), top + Inches(2.20), cw - Inches(0.44),
             Inches(0.88), fill=pale, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.12)
        eb = textbox(s, x + Inches(0.36), top + Inches(2.20), cw - Inches(0.72),
                     Inches(0.88), anchor=MSO_ANCHOR.MIDDLE)
        write(eb.text_frame, example, size=15, color=INK, first=True,
              align=PP_ALIGN.CENTER, line=1.25)

    # نوار پایین + تصویر کوچک
    by = Inches(5.42)
    bh = Inches(1.42)
    pw = Inches(3.05)
    p, sz = img("11-commit.jpg")
    pic = picture_fill(s, p, M, by, pw, bh, sz)
    round_picture(pic)

    bx = M + pw + Inches(0.3)
    bw = SW - M - bx
    rect(s, bx, by, bw, bh, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.10)
    rect(s, bx + bw - Inches(0.09), by, Inches(0.09), bh, fill=ORANGE)
    tb = textbox(s, bx + Inches(0.3), by + Inches(0.18), bw - Inches(0.62),
                 Inches(1.1))
    tf = tb.text_frame
    write(tf, "هفته بعد", size=15, bold=True, color=ORANGE, first=True, line=1.0)
    write(tf, "۵ دقیقه اول جلسه مرور می‌کنیم که تعهدات را چطور اجرا کردیم.",
          size=19, color=WHITE, space_before=7, line=1.2)
    write(tf, "نه برای بازخواست — برای یادگیری از تجربه شما.", size=19,
          bold=True, color=ORANGE, space_before=4, line=1.2)

    footer(s, 11, TOTAL)
    notes(s, """
    اینجا کارت تعهد کاغذی پخش کنید. تعهدها باید کوچک، مشخص و هفتگی باشند.
    """)
    return s


# ═══════════════════════════════════════ ۱۲ — دعوت به سفر ════════════════
def slide_12(prs, cfg):
    s = blank(prs)
    p, sz = img("12-road.jpg")
    picture_fill(s, p, 0, 0, SW, SH, sz)
    gradient(s, 0, 0, SW, SH, NAVY_DEEP, ORANGE_DEEP, angle=45,
             alpha1=0.88, alpha2=0.34)
    rect(s, 0, 0, SW, SH, fill=NAVY_DEEP, alpha=0.28)

    tb = textbox(s, M, Inches(0.62), CONTENT_W, Inches(0.5))
    write(tb.text_frame, "دعوت به سفر", size=18, bold=True, color=ORANGE,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    hb = textbox(s, M, Inches(1.08), CONTENT_W, Inches(0.8))
    write(hb.text_frame, "آماده‌اید؟", size=44, bold=True, color=WHITE,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    qb = textbox(s, Inches(1.9), Inches(2.06), SW - Inches(3.8), Inches(1.6))
    tf = qb.text_frame
    write(tf, "«ما سال‌ها هزینه کردیم تا افراد متخصص جذب کنیم.", size=24,
          color=WHITE, first=True, align=PP_ALIGN.CENTER, line=1.35)
    write(tf, "حالا وقت آن است که اجازه دهیم", size=24, color=WHITE,
          align=PP_ALIGN.CENTER, line=1.35)
    write(tf, "این متخصصان از هم یاد بگیرند.»", size=28, bold=True, color=ORANGE,
          align=PP_ALIGN.CENTER, line=1.35)

    # چهار کارت اطلاعات (راست به چپ)
    infos = [
        ("اولین جلسه رسمی", cfg.first_session),
        ("ساعت", cfg.time),
        ("مکان", cfg.place),
        ("ظرفیت", "۱۲ نفر داوطلب"),
    ]
    gap = Inches(0.28)
    cw = (CONTENT_W - 3 * gap) / 4
    y = Inches(4.10)
    ch = Inches(1.22)
    for i, (label, value) in enumerate(infos):
        x = SW - M - cw - i * (cw + gap)
        box = rect(s, x, y, cw, ch, fill=WHITE, alpha=0.13,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10,
                   line=WHITE, line_w=Pt(1))
        tb = textbox(s, x + Inches(0.16), y, cw - Inches(0.32), ch,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, label, size=14, bold=True, color=ORANGE, first=True,
              align=PP_ALIGN.CENTER, line=1.1)
        write(tb.text_frame, value, size=21, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, line=1.2)

    hline(s, (SW - Inches(3.2)) / 2, Inches(5.72), Inches(3.2), ORANGE, Pt(2.5))

    fb = textbox(s, M, Inches(5.98), CONTENT_W, Inches(1.0))
    tf = fb.text_frame
    write(tf, "بهترین زمان برای شروع، الان است.", size=25, bold=True, color=WHITE,
          first=True, align=PP_ALIGN.CENTER, line=1.25)
    write(tf, "بهترین مکان، همین‌جاست.", size=25, bold=True, color=ORANGE,
          align=PP_ALIGN.CENTER, line=1.25)

    notes(s, """
    پایان: نگویید «سؤالی نیست؟». بگویید «چه کسی می‌خواهد نفر اول باشد؟»
    و دستتان را بالا ببرید تا اولین داوطلب راحت‌تر بلند شود.
    """)
    return s


# ═══════════════════════════════════════════════════════ اجرا ════════════
BUILDERS = [slide_01, slide_02, slide_03, slide_04, slide_05, slide_06,
            slide_07, slide_08, slide_09, slide_10, slide_11, slide_12]


def build(cfg):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    T.set_font(cfg.font)
    for fn in BUILDERS:
        fn(prs, cfg)
    prs.save(cfg.out)
    return cfg.out


def main():
    ap = argparse.ArgumentParser(description="ساخت ارائه منتورینگ گروهی")
    ap.add_argument("--out", default=os.path.join(
        HERE, "از-گفتگو-تا-سازمان-یادگیرنده.pptx"))
    ap.add_argument("--font", default="IRANSans",
                    help="نام فونت فارسی (مثلاً IRANSans یا B Nazanin یا Vazirmatn)")
    ap.add_argument("--date", default="[تاریخ جلسه]")
    ap.add_argument("--presenter", default="[نام ارائه‌کننده]")
    ap.add_argument("--org", default="[نام واحد / سازمان]")
    ap.add_argument("--first-session", dest="first_session", default="[تاریخ]")
    ap.add_argument("--time", default="[ساعت]")
    ap.add_argument("--place", default="[مکان]")
    cfg = ap.parse_args()
    path = build(cfg)
    print("ساخته شد:", path)


if __name__ == "__main__":
    main()
