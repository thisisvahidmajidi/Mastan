# -*- coding: utf-8 -*-
"""
توکن‌های طراحی و توابع کمکی مشترک برای ساخت ارائه‌های راست‌چین (RTL).
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from copy import deepcopy

# ---------------------------------------------------------------- رنگ‌ها ----
NAVY        = RGBColor(0x1B, 0x49, 0x65)   # آبی تیره اصلی
NAVY_DEEP   = RGBColor(0x0E, 0x2C, 0x40)
BLUE_MID    = RGBColor(0x2A, 0x6F, 0x97)
BLUE_LIGHT  = RGBColor(0x62, 0xB6, 0xCB)
BLUE_PALE   = RGBColor(0xEC, 0xF3, 0xF7)
ORANGE      = RGBColor(0xFF, 0x6B, 0x35)
ORANGE_DEEP = RGBColor(0xD9, 0x4E, 0x1F)
ORANGE_PALE = RGBColor(0xFF, 0xF3, 0xEC)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
INK         = RGBColor(0x1F, 0x2D, 0x38)   # متن اصلی
INK_SOFT    = RGBColor(0x5A, 0x6B, 0x77)   # متن کم‌رنگ‌تر
PAPER       = RGBColor(0xF7, 0xFA, 0xFC)
GREEN       = RGBColor(0x1E, 0x8E, 0x5A)
GREEN_PALE  = RGBColor(0xE8, 0xF6, 0xEF)
RED         = RGBColor(0xC0, 0x39, 0x2B)
RED_PALE    = RGBColor(0xFD, 0xEE, 0xEC)
PURPLE      = RGBColor(0x8E, 0x44, 0xAD)
AMBER       = RGBColor(0xF3, 0x9C, 0x12)
LINE        = RGBColor(0xD8, 0xE2, 0xE9)

# ------------------------------------------------------------- اندازه‌ها ----
SW = Inches(13.333)          # عرض اسلاید (16:9)
SH = Inches(7.5)             # ارتفاع اسلاید
M  = Inches(0.62)            # حاشیه کناری
CONTENT_TOP = Inches(1.72)   # شروع ناحیه محتوا در اسلایدهای داخلی
CONTENT_W   = SW - 2 * M

FONT = "IRANSans"            # با سوییچ --font قابل تغییر است
FONT_FALLBACK = "Tahoma"     # فونت جایگزین برای دستگاه‌هایی که IRANSans ندارند

_DIGITS = str.maketrans("0123456789", "۰۱۲۳۴۵۶۷۸۹")


def fa(text):
    """تبدیل ارقام لاتین به ارقام فارسی."""
    return str(text).translate(_DIGITS)


def set_font(cfg_font):
    global FONT
    FONT = cfg_font


# ------------------------------------------------------- کمک‌کارهای XML ----
def _pPr(p):
    return p._p.get_or_add_pPr()


def rtl(p, align=PP_ALIGN.RIGHT):
    """راست‌چین کردن پاراگراف در سطح XML (rtl='1')."""
    pr = _pPr(p)
    pr.set("rtl", "1")
    p.alignment = align
    return p


def set_alpha(fill, alpha):
    """شفافیت رنگ توپر یک شکل (alpha بین 0 و 1)."""
    xPr = fill._xPr
    solid = xPr.find(qn("a:solidFill"))
    if solid is None:
        return
    clr = solid[0]
    for old in clr.findall(qn("a:alpha")):
        clr.remove(old)
    node = parse_xml(
        '<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="%d"/>'
        % int(alpha * 100000)
    )
    clr.append(node)


def no_line(shape):
    shape.line.fill.background()
    return shape


def shadow_off(shape):
    """حذف سایه پیش‌فرض تم."""
    spPr = shape._element.spPr
    for tag in ("a:effectLst", "a:effectRef"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    spPr.append(parse_xml(
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
    ))
    return shape


# ------------------------------------------------------------- شکل‌سازی ----
def rect(slide, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE,
         radius=None, alpha=None, line=None, line_w=Pt(1)):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sp.adjustments[0] = radius
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
        if alpha is not None:
            set_alpha(sp.fill, alpha)
    if line is None:
        no_line(sp)
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w
    shadow_off(sp)
    sp.text_frame.word_wrap = True
    return sp


def card(slide, x, y, w, h, fill=WHITE, radius=0.045, line=LINE):
    return rect(slide, x, y, w, h, fill=fill,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius, line=line)


def gradient(slide, x, y, w, h, c1, c2, angle=45, alpha1=None, alpha2=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    no_line(sp)
    shadow_off(sp)
    sp.fill.gradient()
    stops = sp.fill.gradient_stops
    stops[0].color.rgb = c1
    stops[0].position = 0.0
    stops[1].color.rgb = c2
    stops[1].position = 1.0
    while len(stops._gsLst) > 2:
        stops._gsLst.remove(stops._gsLst[-1])
    sp.fill.gradient_angle = angle
    if alpha1 is not None:
        _stop_alpha(stops._gsLst[0], alpha1)
    if alpha2 is not None:
        _stop_alpha(stops._gsLst[1], alpha2)
    return sp


def _stop_alpha(gs, alpha):
    clr = gs[0]
    node = parse_xml(
        '<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="%d"/>'
        % int(alpha * 100000)
    )
    clr.append(node)


def hline(slide, x, y, w, color=LINE, thickness=Pt(1.25)):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, thickness)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    no_line(sp)
    shadow_off(sp)
    return sp


# --------------------------------------------------------------- متن‌ها ----
def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tb


def para(tf, first=False):
    return tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()


def write(tf, chunks, size=24, color=INK, bold=False, align=PP_ALIGN.RIGHT,
          space_before=0, space_after=0, line=1.25, first=False, font=None):
    """
    نوشتن یک پاراگراف راست‌چین.
    chunks: رشته، یا لیستی از (متن، bold) / (متن، bold، رنگ) / (متن، bold، رنگ، سایز)
    """
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs
                             and len(tf.paragraphs) == 1) else tf.add_paragraph()
    rtl(p, align)
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line
    if isinstance(chunks, str):
        chunks = [(chunks, bold, color, size)]
    for ch in chunks:
        txt = ch[0]
        b = ch[1] if len(ch) > 1 else bold
        c = ch[2] if len(ch) > 2 and ch[2] is not None else color
        s = ch[3] if len(ch) > 3 and ch[3] is not None else size
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(s)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = font or FONT
        # اعمال فونت روی اسکریپت‌های عربی/لاتین
        rPr = r._r.get_or_add_rPr()
        for tag in ("a:cs", "a:ea"):
            el = rPr.find(qn(tag))
            if el is None:
                el = parse_xml(
                    '<%s xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="%s"/>'
                    % (tag, font or FONT))
                rPr.append(el)
            else:
                el.set("typeface", font or FONT)
    return p


def bullet(tf, text, size=24, color=INK, bold=False, marker="◆",
           marker_color=ORANGE, space_before=8, line=1.2):
    p = write(tf, [(marker + "  ", False, marker_color, size * 0.72)],
              space_before=space_before, line=line)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:cs", "a:ea"):
        el = parse_xml(
            '<%s xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="%s"/>'
            % (tag, FONT))
        rPr.append(el)
    return p


# --------------------------------------------------------------- تصاویر ----
def picture_fill(slide, path, x, y, w, h, img_size):
    """درج تصویر با برش (crop) طوری که کادر داده‌شده را کامل پر کند."""
    iw, ih = img_size
    box_ar = w / h
    img_ar = iw / ih
    pic = slide.shapes.add_picture(path, x, y, w, h)
    if img_ar > box_ar:                       # تصویر پهن‌تر → برش افقی
        keep = box_ar / img_ar
        c = (1 - keep) / 2
        pic.crop_left = c
        pic.crop_right = c
    else:                                     # تصویر بلندتر → برش عمودی
        keep = img_ar / box_ar
        c = (1 - keep) / 2
        pic.crop_top = c
        pic.crop_bottom = c
    return pic


def donut_segment(slide, cx, cy, r_out, r_in, a0, a1, color, steps=48):
    """
    یک قطاع حلقه‌ای (بخشی از نمودار دونات) به‌صورت شکل آزاد (freeform).
    زاویه‌ها بر حسب درجه و در جهت ساعتگرد از بالای دایره اندازه‌گیری می‌شوند.
    """
    import math

    def pt(r, ang_deg):
        a = math.radians(ang_deg - 90)     # ۰ درجه = بالای دایره
        return (int(cx + r * math.cos(a)), int(cy + r * math.sin(a)))

    pts = [pt(r_out, a0 + (a1 - a0) * i / steps) for i in range(steps + 1)]
    pts += [pt(r_in, a1 - (a1 - a0) * i / steps) for i in range(steps + 1)]

    builder = slide.shapes.build_freeform(pts[0][0], pts[0][1])
    builder.add_line_segments(pts[1:], close=True)
    sp = builder.convert_to_shape()
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    no_line(sp)
    shadow_off(sp)
    return sp


def round_picture(pic):
    """گوشه‌گرد کردن تصویر."""
    pic._element.spPr.find(qn("a:prstGeom")).set("prst", "roundRect")
    return pic


# ---------------------------------------------------- اسکلت اسلاید داخلی ----
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def base_slide(prs, bg=PAPER):
    s = blank(prs)
    rect(s, 0, 0, SW, SH, fill=bg)
    return s


def header(slide, title, kicker=None, color=NAVY, accent=ORANGE, on_dark=False):
    """تیتر بالای اسلاید + نوار تأکید نارنجی (راست‌چین)."""
    top = Inches(0.44)
    if kicker:
        kb = textbox(slide, M, Inches(0.34), CONTENT_W, Inches(0.36))
        write(kb.text_frame, kicker, size=17, bold=True,
              color=ORANGE if not on_dark else BLUE_LIGHT, first=True, line=1.0)
        top = Inches(0.74)
    tb = textbox(slide, M, top, CONTENT_W, Inches(0.72))
    write(tb.text_frame, title, size=33, bold=True,
          color=WHITE if on_dark else color, first=True, line=1.05)
    bar_w = Inches(1.35)
    hline(slide, SW - M - bar_w, Inches(1.44), bar_w, color=accent, thickness=Pt(4.5))
    return slide


def footer(slide, number, total=12, dark=False, label="از گفتگو تا سازمان یادگیرنده"):
    tb = textbox(slide, M, SH - Inches(0.52), Inches(4.0), Inches(0.3))
    write(tb.text_frame, label, size=11.5,
          color=WHITE if dark else INK_SOFT, first=True, line=1.0,
          align=PP_ALIGN.RIGHT)
    nb = textbox(slide, SW - M - Inches(1.6), SH - Inches(0.52), Inches(1.6), Inches(0.3))
    write(nb.text_frame, fa("%d / %d" % (number, total)), size=11.5,
          color=WHITE if dark else INK_SOFT, first=True, line=1.0,
          align=PP_ALIGN.LEFT)
    return slide


def notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = ""
    for i, line in enumerate(text.strip().split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        rtl(p)
        r = p.add_run()
        r.text = line.strip()
        r.font.size = Pt(13)
        r.font.name = FONT
    return slide
