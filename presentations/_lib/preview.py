# -*- coding: utf-8 -*-
"""
رندر تقریبی فایل PPTX به HTML برای بازبینی چیدمان در مرورگر.
اجرا: python3 preview.py <file.pptx> -o preview.html
"""
import argparse
import base64
import html
import os
import sys

from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

EMU_IN = 914400.0
PX = 96.0  # پیکسل بر اینچ در خروجی HTML


def emu2px(v):
    return (v or 0) / EMU_IN * PX


def _color_of(el):
    """استخراج (hex, alpha) از یک نود رنگ DrawingML."""
    if el is None:
        return None, 1.0
    srgb = el.find(qn("a:srgbClr"))
    if srgb is None:
        return None, 1.0
    hexv = srgb.get("val")
    alpha = 1.0
    a = srgb.find(qn("a:alpha"))
    if a is not None:
        alpha = int(a.get("val")) / 100000.0
    return "#" + hexv, alpha


def shape_fill_css(sp):
    spPr = getattr(sp._element, "spPr", None)
    if spPr is None:
        return "transparent"
    if spPr.find(qn("a:noFill")) is not None:
        return "transparent"
    solid = spPr.find(qn("a:solidFill"))
    if solid is not None:
        c, a = _color_of(solid)
        if c:
            return _rgba(c, a)
    grad = spPr.find(qn("a:gradFill"))
    if grad is not None:
        stops = []
        gsLst = grad.find(qn("a:gsLst"))
        for gs in gsLst:
            c, a = _color_of(gs)
            pos = int(gs.get("pos", "0")) / 100000.0
            stops.append((pos, _rgba(c or "#000000", a)))
        lin = grad.find(qn("a:lin"))
        ang = 0
        if lin is not None:
            ang = int(lin.get("ang", "0")) / 60000.0
        css_ang = (ang + 90) % 360   # تبدیل زاویه OOXML به CSS
        return "linear-gradient(%ddeg, %s)" % (
            css_ang, ", ".join("%s %d%%" % (c, p * 100) for p, c in stops))
    return "transparent"


def _rgba(hexv, alpha):
    hexv = hexv.lstrip("#")
    r, g, b = int(hexv[0:2], 16), int(hexv[2:4], 16), int(hexv[4:6], 16)
    return "rgba(%d,%d,%d,%.3f)" % (r, g, b, alpha)


def line_css(sp):
    spPr = getattr(sp._element, "spPr", None)
    if spPr is None:
        return ""
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        return ""
    if ln.find(qn("a:noFill")) is not None:
        return ""
    solid = ln.find(qn("a:solidFill"))
    if solid is None:
        return ""
    c, a = _color_of(solid)
    w = int(ln.get("w", "12700")) / 12700.0
    return "border:%.2fpx solid %s;" % (w * PX / 72.0, _rgba(c or "#000", a))


def geom_of(sp):
    spPr = getattr(sp._element, "spPr", None)
    if spPr is None:
        return None
    g = spPr.find(qn("a:prstGeom"))
    return g.get("prst") if g is not None else None


ALIGN = {None: "right", "r": "right", "l": "left", "ctr": "center",
         "just": "justify"}
ANCHOR = {None: "flex-start", "t": "flex-start", "ctr": "center",
          "b": "flex-end"}


def render_text(sp, out, x, y, w, h):
    if not sp.has_text_frame:
        return
    tf = sp.text_frame
    txBody = sp._element.find(qn("p:txBody"))
    if txBody is None:
        txBody = sp._element.find(qn("a:txBody"))
    bodyPr = txBody.find(qn("a:bodyPr")) if txBody is not None else None
    anchor = ANCHOR.get(bodyPr.get("anchor") if bodyPr is not None else None)
    ml = emu2px(tf.margin_left); mr = emu2px(tf.margin_right)
    mt = emu2px(tf.margin_top); mb = emu2px(tf.margin_bottom)
    paras = []
    any_text = False
    for p in tf.paragraphs:
        pPr = p._p.find(qn("a:pPr"))
        algn = pPr.get("algn") if pPr is not None else None
        sb = p.space_before.pt if p.space_before else 0
        sa = p.space_after.pt if p.space_after else 0
        ls = p.line_spacing if p.line_spacing else 1.2
        runs = []
        for r in p.runs:
            any_text = True
            f = r.font
            size = f.size.pt if f.size else 18
            color = "#000000"
            try:
                if f.color and f.color.type is not None and f.color.rgb:
                    color = "#" + str(f.color.rgb)
            except Exception:
                pass
            runs.append(
                '<span style="font-size:%.2fpx;font-weight:%s;color:%s">%s</span>'
                % (size * PX / 72.0, "700" if f.bold else "400", color,
                   html.escape(r.text).replace("\n", "<br>")))
        paras.append(
            '<p style="margin:%.1fpx 0 %.1fpx 0;line-height:%.2f;text-align:%s">%s</p>'
            % (sb * PX / 72.0, sa * PX / 72.0, ls, ALIGN.get(algn), "".join(runs) or "&nbsp;"))
    if not any_text:
        return
    out.append(
        '<div class="tf" style="left:%.1fpx;top:%.1fpx;width:%.1fpx;height:%.1fpx;'
        'padding:%.1fpx %.1fpx %.1fpx %.1fpx;justify-content:%s">%s</div>'
        % (x, y, w, h, mt, mr, mb, ml, anchor, "".join(paras)))


def render_shape(sp, out, tmpdir, idx):
    x, y = emu2px(sp.left), emu2px(sp.top)
    w, h = emu2px(sp.width), emu2px(sp.height)

    if sp.shape_type is not None and sp.shape_type == 13:  # PICTURE
        blob = sp.image.blob
        b64 = base64.b64encode(blob).decode()
        radius = "18px" if geom_of(sp) == "roundRect" else "0"
        out.append(
            '<div style="position:absolute;left:%.1fpx;top:%.1fpx;width:%.1fpx;'
            'height:%.1fpx;overflow:hidden;border-radius:%s">'
            '<img src="data:image/jpeg;base64,%s" style="width:100%%;height:100%%;'
            'object-fit:cover"></div>' % (x, y, w, h, radius, b64))
        return

    prst = geom_of(sp)
    fill = shape_fill_css(sp)
    border = line_css(sp)
    style = ("position:absolute;left:%.1fpx;top:%.1fpx;width:%.1fpx;height:%.1fpx;"
             "background:%s;%s" % (x, y, w, h, fill, border))
    if prst == "roundRect":
        try:
            adj = sp.adjustments[0]
        except Exception:
            adj = 0.16
        style += "border-radius:%.1fpx;" % (min(w, h) * adj)
    elif prst == "ellipse":
        style += "border-radius:50%;"
    elif prst == "donut":
        style += "border-radius:50%;"
    elif prst == "blockArc":
        style = ("position:absolute;left:%.1fpx;top:%.1fpx;width:%.1fpx;height:%.1fpx;"
                 "border-radius:50%%;background:%s;opacity:.9;"
                 "clip-path:polygon(50%% 50%%, 100%% 0, 100%% 100%%, 0 100%%);"
                 % (x, y, w, h, fill))
    elif prst == "upArrow":
        style += "clip-path:polygon(50% 0,100% 55%,72% 55%,72% 100%,28% 100%,28% 55%,0 55%);"
    out.append('<div style="%s"></div>' % style)
    render_text(sp, out, x, y, w, h)


def render(path, out_path, font_stack):
    prs = Presentation(path)
    W = emu2px(prs.slide_width)
    H = emu2px(prs.slide_height)
    parts = ["""<!doctype html><html lang="fa" dir="rtl"><head><meta charset="utf-8">
<title>پیش‌نمایش ارائه</title>
<style>
 body{background:#20262b;margin:0;padding:26px;font-family:%s;}
 .slide{position:relative;width:%.0fpx;height:%.0fpx;background:#fff;margin:0 auto 26px;
        box-shadow:0 10px 34px rgba(0,0,0,.5);overflow:hidden;}
 .tf{position:absolute;display:flex;flex-direction:column;box-sizing:border-box;
     direction:rtl;}
 .tf p{margin:0;}
 .num{color:#8fa3b0;text-align:center;font-size:13px;margin:0 0 8px;}
</style></head><body>""" % (font_stack, W, H)]
    for i, slide in enumerate(prs.slides, 1):
        parts.append('<div class="num">اسلاید %d</div><div class="slide">' % i)
        buf = []
        for sp in slide.shapes:
            try:
                render_shape(sp, buf, None, i)
            except Exception as e:
                sys.stderr.write("skip shape on slide %d: %s\n" % (i, e))
        parts.append("".join(buf))
        parts.append("</div>")
    parts.append("</body></html>")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("".join(parts))
    return out_path


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("-o", "--out", default="preview.html")
    ap.add_argument("--font", default="Vazirmatn, Tahoma, 'DejaVu Sans', sans-serif")
    a = ap.parse_args()
    print(render(a.pptx, a.out, a.font))
