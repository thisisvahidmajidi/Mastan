# -*- coding: utf-8 -*-
"""
ساخت ارائه جلسه افتتاحیه منتورingگ گروهی: «با هم بلد می‌شویم»
ساختار: چیستی → چرایی → نحوه اجرا → تجربه زنده

اجرا:  python3 build.py [--presenter ...] [--font IRANSans] [--out out.pptx]
"""
import argparse
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                "..", "_lib"))
import theme as T
from theme import (
    NAVY, NAVY_DEEP, BLUE_MID, BLUE_LIGHT, BLUE_PALE, ORANGE, ORANGE_DEEP,
    ORANGE_PALE, WHITE, INK, INK_SOFT, PAPER, GREEN, GREEN_PALE, RED, RED_PALE,
    PURPLE, AMBER, LINE, SW, SH, M, CONTENT_W,
    fa, rect, card, gradient, hline, textbox, write, picture_fill,
    round_picture, base_slide, blank, footer, notes, no_line, shadow_off,
)

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")
TOTAL = 16

# سه فصل ارائه — در سربرگ هر اسلاید نمایش داده می‌شود
CHAPTERS = ["چیستی", "چرایی", "نحوه اجرا"]


def img(name):
    p = os.path.join(IMG, name)
    return p, Image.open(p).size


# ═══════════════════════════════════════════════ سربرگ اختصاصی ═══════════
def head(slide, title, chapter=None, kicker=None):
    """
    سربرگ با نشانگر فصل: سه چیپ کوچک که فصل جاری در آن‌ها پررنگ است.
    این نشانگر جای سه اسلاید جداکننده فصل را می‌گیرد.
    """
    if chapter is not None:
        cx = M
        for i, name in enumerate(CHAPTERS):
            active = (i == chapter)
            w = Inches(1.16)
            x = cx + i * (w + Inches(0.10))
            rect(slide, x, Inches(0.46), w, Inches(0.34),
                 fill=NAVY if active else T.RGBColor(0xE4, 0xEC, 0xF2),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.42)
            tb = textbox(slide, x, Inches(0.46), w, Inches(0.34),
                         anchor=MSO_ANCHOR.MIDDLE)
            write(tb.text_frame, name, size=11.5, bold=True,
                  color=WHITE if active else INK_SOFT, first=True,
                  align=PP_ALIGN.CENTER, line=1.0)

    top = Inches(0.40)
    if kicker:
        kb = textbox(slide, M, Inches(0.34), CONTENT_W, Inches(0.34))
        write(kb.text_frame, kicker, size=16, bold=True, color=ORANGE,
              first=True, line=1.0)
        top = Inches(0.70)

    tb = textbox(slide, M, top, CONTENT_W, Inches(0.70))
    write(tb.text_frame, title, size=32, bold=True, color=NAVY, first=True,
          line=1.05)
    bar_w = Inches(1.25)
    hline(slide, SW - M - bar_w, Inches(1.38), bar_w, ORANGE, Pt(4.5))
    return slide


def quote_band(slide, y, chunks, h=Inches(0.92), fill=NAVY, accent=ORANGE,
               size=22, align=PP_ALIGN.CENTER, x=M, w=None):
    """نوار جمع‌بندی تیره با لبه نارنجی در سمت راست."""
    w = w or CONTENT_W
    rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.10)
    rect(slide, x + w - Inches(0.09), y, Inches(0.09), h, fill=accent)
    tb = textbox(slide, x + Inches(0.34), y, w - Inches(0.68), h,
                 anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame, chunks, size=size, color=WHITE, first=True,
          align=align, line=1.2)
    return slide


# ══════════════════════════════════════════════════ ۱ — جلد ══════════════
def s01(prs, cfg):
    s = blank(prs)
    gradient(s, 0, 0, SW, SH, NAVY_DEEP, BLUE_MID, angle=315)

    pw = Inches(6.35)
    p, sz = img("01-kickoff.jpg")
    picture_fill(s, p, SW - pw, 0, pw, SH, sz)
    gradient(s, SW - pw, 0, pw, SH, NAVY_DEEP, NAVY_DEEP, angle=0,
             alpha1=0.96, alpha2=0.08)

    tx = Inches(0.80)
    tw = Inches(6.35)
    hline(s, tx + tw - Inches(1.4), Inches(1.70), Inches(1.4), ORANGE, Pt(5))

    tb = textbox(s, tx, Inches(2.02), tw, Inches(3.3))
    tf = tb.text_frame
    write(tf, "جلسه افتتاحیه منتورینگ گروهی", size=18, bold=True,
          color=BLUE_LIGHT, first=True, line=1.0)
    write(tf, "با هم", size=52, bold=True, color=WHITE, space_before=16,
          line=1.22)
    write(tf, [("بلد ", True, ORANGE, 52), ("می‌شویم", True, WHITE, 52)],
          line=1.22)
    write(tf, "چیستی، چرایی و نحوه اجرای یادگیری همتامحور",
          size=20, color=BLUE_PALE, space_before=20, line=1.3)

    hline(s, tx + tw - Inches(3.2), Inches(5.72), Inches(3.2), BLUE_LIGHT,
          Pt(1.25))
    mb = textbox(s, tx, Inches(5.95), tw, Inches(0.9))
    write(mb.text_frame,
          [("تاریخ:  ", False, BLUE_LIGHT, 16), (cfg.date, True, WHITE, 16),
           ("        تسهیل‌گر:  ", False, BLUE_LIGHT, 16),
           (cfg.presenter, True, WHITE, 16)], first=True, line=1.2)
    write(mb.text_frame, cfg.org, size=13.5, color=BLUE_LIGHT, space_before=6,
          line=1.2)

    notes(s, """
    قبل از شروع، صندلی‌ها را دایره‌ای بچینید. پشت میز یا جایگاه نایستید.
    جمله شروع: «این آخرین جلسه‌ای است که من جلوی شما می‌ایستم و حرف می‌زنم.»
    زمان: ۱ دقیقه.
    """)
    return s


# ═════════════════════════════════════════ ۲ — قرار امروز ════════════════
def s02(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "قرار امروز چیست؟", kicker="۹۰ دقیقه پیشِ رو")

    items = [
        ("۱۵", "چیستی", "این روش دقیقاً چیست و چه چیزی نیست", NAVY),
        ("۲۰", "چرایی", "چرا این روش کار می‌کند و به چه درد ما می‌خورد", BLUE_MID),
        ("۲۰", "نحوه اجرا", "یک جلسه واقعی چطور پیش می‌رود", BLUE_LIGHT),
        ("۳۰", "تجربه زنده", "یک دور کوتاه، همین امروز، با هم", ORANGE),
        ("۰۵", "تصمیم", "چه کسی می‌خواهد ادامه دهد", GREEN),
    ]

    y = Inches(1.75)
    rh = Inches(0.80)
    tw = Inches(8.15)
    tx = SW - M - tw
    for mins, title, desc, color in items:
        rect(s, tx, y, tw, rh, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.13, line=LINE)
        chip = rect(s, tx + tw - Inches(1.20), y + Inches(0.13), Inches(1.02),
                    rh - Inches(0.26), fill=color,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.26)
        ctf = chip.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(ctf, fa(mins) + "′", size=17, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, tx + Inches(0.30), y, tw - Inches(1.62), rh,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=21, bold=True, color=color, first=True,
              line=1.05)
        write(tb.text_frame, desc, size=15, color=INK_SOFT, line=1.15)
        y += rh + Inches(0.11)

    # ستون تعهد سمت چپ
    px = M
    pw = tx - M - Inches(0.42)
    box = card(s, px, Inches(1.75), pw, Inches(2.55), fill=NAVY, line=None)
    rect(s, px + pw - Inches(0.09), Inches(1.75), Inches(0.09), Inches(2.55),
         fill=ORANGE)
    tb = textbox(s, px + Inches(0.26), Inches(1.98), pw - Inches(0.55),
                 Inches(2.1))
    tf = tb.text_frame
    write(tf, "قولِ امروز", size=15, bold=True, color=ORANGE, first=True,
          line=1.0)
    write(tf, "بیش از یک‌سومِ این جلسه، خودِ تجربه است — نه توضیحِ تجربه.",
          size=18, color=WHITE, space_before=10, line=1.3)
    write(tf, "امروز درباره منتورینگ گروهی حرف نمی‌زنیم؛ یک بار انجامش می‌دهیم.",
          size=15, color=BLUE_LIGHT, space_before=8, line=1.3)

    box2 = card(s, px, Inches(4.48), pw, Inches(2.28), fill=ORANGE_PALE,
                line=ORANGE)
    tb = textbox(s, px + Inches(0.26), Inches(4.70), pw - Inches(0.52),
                 Inches(1.9))
    tf = tb.text_frame
    write(tf, "از شما یک چیز می‌خواهم", size=15, bold=True, color=ORANGE_DEEP,
          first=True, line=1.0)
    write(tf, "یک چالش کاری واقعی و حل‌نشده را همین حالا در ذهنتان نگه دارید.",
          size=18, bold=True, color=NAVY, space_before=10, line=1.3)
    write(tf, "در پایان جلسه به آن برمی‌گردیم.", size=15, color=INK_SOFT,
          space_before=7, line=1.2)

    footer(s, 2, TOTAL, label="با هم بلد می‌شویم")
    notes(s, """
    روی «تجربه زنده» مکث کنید. مخاطب باید بداند این جلسه سخنرانی نیست.
    درخواستِ «یک چالش واقعی در ذهن نگه دارید» را جدی و آرام بگویید.
    """)
    return s


# ══════════════════════════════════ ۳ — سؤال شروع (قلاب) ═════════════════
def s03(prs, cfg):
    s = base_slide(prs, WHITE)
    gradient(s, 0, 0, SW, SH, NAVY_DEEP, NAVY, angle=315)

    tb = textbox(s, Inches(1.4), Inches(1.15), SW - Inches(2.8), Inches(0.5))
    write(tb.text_frame, "۶۰ ثانیه فکر کنید", size=17, bold=True, color=ORANGE,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    qb = textbox(s, Inches(1.25), Inches(1.78), SW - Inches(2.5), Inches(2.0))
    tf = qb.text_frame
    write(tf, "آخرین چیزی که در کارتان یاد گرفتید،", size=36, bold=True,
          color=WHITE, first=True, align=PP_ALIGN.CENTER, line=1.3)
    write(tf, [("از کجا", True, ORANGE, 36), (" یاد گرفتید؟", True, WHITE, 36)],
          align=PP_ALIGN.CENTER, line=1.3)

    # سه گزینه
    opts = [
        ("از یک کلاس یا کتاب", BLUE_LIGHT),
        ("از انجام دادن و اشتباه کردن", BLUE_LIGHT),
        ("از یک آدم", ORANGE),
    ]
    gap = Inches(0.34)
    cw = (CONTENT_W - 2 * gap) / 3
    y = Inches(4.12)
    ch = Inches(1.18)
    for i, (label, color) in enumerate(opts):
        x = SW - M - cw - i * (cw + gap)
        hi = (color == ORANGE)
        rect(s, x, y, cw, ch, fill=WHITE, alpha=0.20 if hi else 0.09,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11,
             line=ORANGE if hi else BLUE_LIGHT, line_w=Pt(2 if hi else 1))
        tb = textbox(s, x + Inches(0.20), y, cw - Inches(0.40), ch,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, label, size=19, bold=hi,
              color=ORANGE if hi else WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.25)

    fb = textbox(s, Inches(1.4), Inches(5.72), SW - Inches(2.8), Inches(1.0))
    tf = fb.text_frame
    write(tf, "تقریباً همیشه، جواب یک آدم است.", size=25, bold=True,
          color=WHITE, first=True, align=PP_ALIGN.CENTER, line=1.25)
    write(tf, "منتورینگ گروهی یعنی این اتفاق را از «شانسی» به «برنامه‌ریزی‌شده» تبدیل کنیم.",
          size=17, color=BLUE_PALE, align=PP_ALIGN.CENTER, line=1.25)

    notes(s, """
    واقعاً ۶۰ ثانیه سکوت بدهید و از ۳ تا ۴ نفر بپرسید.
    تقریباً همه یک اسم می‌گویند، نه یک دوره. همین، کل استدلال ارائه است.
    اگر کسی «کلاس» گفت، بپرسید: «چه چیزی از آن کلاس هنوز یادتان است؟»
    """)
    return s


# ════════════════════════════ ۴ — چیستی: تعریف در یک جمله ════════════════
def s04(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "منتورینگ گروهی، در یک جمله", chapter=0, kicker="تعریف")

    # جمله تعریف بزرگ
    box = card(s, M, Inches(1.72), CONTENT_W, Inches(1.62), fill=WHITE,
               line=LINE)
    rect(s, SW - M - Inches(0.10), Inches(1.72), Inches(0.10), Inches(1.62),
         fill=ORANGE)
    tb = textbox(s, M + Inches(0.36), Inches(1.72), CONTENT_W - Inches(0.72),
                 Inches(1.62), anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame,
          [("گروهی کوچک از همکاران که ", False, INK, 26),
           ("مرتب", True, ORANGE_DEEP, 26),
           (" دور هم جمع می‌شوند تا روی ", False, INK, 26),
           ("چالش‌های واقعی کاری", True, NAVY, 26),
           (" یکدیگر فکر کنند — با ", False, INK, 26),
           ("پرسش", True, ORANGE_DEEP, 26), ("، نه با ", False, INK, 26),
           ("نسخه‌پیچی", True, INK_SOFT, 26), (".", False, INK, 26)],
          first=True, align=PP_ALIGN.CENTER, line=1.4)

    # چهار مشخصه کلیدی
    specs = [
        ("۸ تا ۱۲", "نفر در هر گروه", NAVY),
        ("۹۰", "دقیقه در هر جلسه", BLUE_MID),
        ("هر ۲ هفته", "یک بار", BLUE_LIGHT),
        ("۶ ماه", "طول یک دوره", ORANGE),
    ]
    gap = Inches(0.30)
    cw = (CONTENT_W - 3 * gap) / 4
    y = Inches(3.65)
    ch = Inches(1.32)
    for i, (big, small, color) in enumerate(specs):
        x = SW - M - cw - i * (cw + gap)
        rect(s, x, y, cw, ch, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.11, line=LINE)
        rect(s, x, y, cw, Inches(0.07), fill=color)
        tb = textbox(s, x + Inches(0.14), y, cw - Inches(0.28), ch,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, big, size=27, bold=True, color=color, first=True,
              align=PP_ALIGN.CENTER, line=1.05)
        write(tb.text_frame, small, size=15, color=INK_SOFT,
              align=PP_ALIGN.CENTER, line=1.15)
        y_txt = y

    # تفاوت با منتورینگ کلاسیک
    dy = Inches(5.25)
    gap2 = Inches(0.34)
    hw = (CONTENT_W - gap2) / 2
    pairs = [
        (SW - M - hw, "منتورینگ کلاسیک", "یک نفر می‌داند، یک نفر یاد می‌گیرد",
         "۱ ← ۱", INK_SOFT, T.RGBColor(0xEC, 0xF0, 0xF3)),
        (M, "منتورینگ گروهی", "همه می‌دانند، همه یاد می‌گیرند",
         "همه ⇄ همه", ORANGE_DEEP, ORANGE_PALE),
    ]
    for x, title, desc, badge, color, pale in pairs:
        rect(s, x, dy, hw, Inches(1.28), fill=pale,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11)
        bd = rect(s, x + hw - Inches(1.65), dy + Inches(0.30), Inches(1.42),
                  Inches(0.66), fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                  radius=0.26)
        btf = bd.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(btf, badge, size=15, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, x + Inches(0.28), dy, hw - Inches(2.05), Inches(1.28),
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=19, bold=True, color=color, first=True,
              line=1.05)
        write(tb.text_frame, desc, size=15, color=INK, line=1.2)

    footer(s, 4, TOTAL, label="با هم بلد می‌شویم")
    notes(s, """
    کلمه کلیدی «پرسش، نه نسخه‌پیچی» را دو بار تکرار کنید — کل تفاوت روش همین است.
    تفاوت با منتورینگ کلاسیک را سریع رد شوید؛ اسلاید بعد مرزها را دقیق می‌کند.
    """)
    return s


# ═══════════════════════════ ۵ — چیستی: هست / نیست ═══════════════════════
def s05(prs, cfg):
    s = base_slide(prs, WHITE)
    head(s, "چه چیزی هست، چه چیزی نیست", chapter=0,
         kicker="مرزبندی — تا انتظار اشتباه شکل نگیرد")

    gap = Inches(0.42)
    cw = (CONTENT_W - gap) / 2
    top = Inches(1.72)
    ch = Inches(0.62)
    rowh = Inches(0.74)
    rows = 5
    boxh = ch + rows * rowh + Inches(0.22)

    cols = [
        (M + cw + gap, "هست", GREEN, GREEN_PALE, "✓", [
            ("فضای فکر کردن با هم", "روی مسئله‌ای که واقعاً درگیرش هستید"),
            ("یادگیری از تجربه همکار", "کسی که مسیر شما را رفته است"),
            ("تمرین پرسیدن", "مهارتی که در جلسات کاری هم به کار می‌آید"),
            ("تعهد به یک اقدام کوچک", "هر جلسه، یک قدم مشخص"),
            ("محرمانه و داوطلبانه", "هیچ‌کس اجباری نیست"),
        ]),
        (M, "نیست", RED, RED_PALE, "✕", [
            ("کلاس آموزشی", "کسی قرار نیست تدریس کند"),
            ("جلسه گزارش‌دهی", "هیچ صورت‌جلسه‌ای به بالا نمی‌رود"),
            ("ارزیابی عملکرد", "روی پرونده پرسنلی شما اثری ندارد"),
            ("جلسه شکایت", "درباره مسئله حرف می‌زنیم، نه درباره آدم‌ها"),
            ("جای نصیحت کردن", "قدیمی‌ترها هم اینجا فقط می‌پرسند"),
        ]),
    ]

    for x, title, cmain, cpale, mark, items in cols:
        card(s, x, top, cw, boxh, fill=WHITE, line=LINE)
        rect(s, x, top, cw, ch, fill=cmain, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.12)
        rect(s, x, top + ch - Inches(0.16), cw, Inches(0.16), fill=cmain)
        tb = textbox(s, x, top, cw, ch, anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, [(mark + "   ", True, WHITE, 19),
                              (title, True, WHITE, 23)],
              first=True, align=PP_ALIGN.CENTER, line=1.0)
        for i, (main, sub) in enumerate(items):
            y = top + ch + Inches(0.11) + i * rowh
            if i % 2 == 0:
                rect(s, x + Inches(0.13), y, cw - Inches(0.26),
                     rowh - Inches(0.06), fill=cpale,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
            tb = textbox(s, x + Inches(0.30), y, cw - Inches(0.60),
                         rowh - Inches(0.06), anchor=MSO_ANCHOR.MIDDLE)
            write(tb.text_frame, [(mark + "  ", True, cmain, 14),
                                  (main, True, INK, 18.5)],
                  first=True, line=1.05)
            write(tb.text_frame, sub, size=13.5, color=INK_SOFT, line=1.15)

    footer(s, 5, TOTAL, label="با هم بلد می‌شویم")
    notes(s, """
    ستون «نیست» مهم‌تر از ستون «هست» است. بیشترین مقاومت از همین سوءتفاهم‌ها می‌آید.
    روی «هیچ صورت‌جلسه‌ای به بالا نمی‌رود» مکث کنید و در چشم افراد نگاه کنید.
    """)
    return s


# ══════════════════════════ ۶ — چیستی: سه نقش در اتاق ════════════════════
def s06(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "در اتاق، سه نقش وجود دارد", chapter=0, kicker="نقش‌ها")

    roles = [
        ("۱", "صاحب چالش", "یک نفر، داوطلب",
         ["مسئله واقعی‌اش را روایت می‌کند",
          "فقط گوش می‌دهد و یادداشت برمی‌دارد",
          "در پایان تصمیم می‌گیرد چه کند"], ORANGE, ORANGE_PALE),
        ("۲", "گروه همتا", "۷ تا ۱۱ نفر",
         ["سؤال می‌پرسند، راه‌حل نمی‌دهند",
          "تجربه مشابه خود را روایت می‌کنند",
          "قضاوت نمی‌کنند"], NAVY, BLUE_PALE),
        ("۳", "تسهیل‌گر", "چرخشی، هر جلسه یک نفر",
         ["زمان را نگه می‌دارد",
          "قواعد را یادآوری می‌کند",
          "خودش کمترین حرف را می‌زند"], GREEN, GREEN_PALE),
    ]

    gap = Inches(0.34)
    cw = (CONTENT_W - 2 * gap) / 3
    top = Inches(1.72)
    ch = Inches(3.62)

    for i, (num, title, sub, points, color, pale) in enumerate(roles):
        x = SW - M - cw - i * (cw + gap)
        card(s, x, top, cw, ch, fill=WHITE, line=LINE)
        rect(s, x, top, cw, Inches(1.08), fill=color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
        rect(s, x, top + Inches(0.92), cw, Inches(0.16), fill=color)
        bd = rect(s, x + cw - Inches(0.72), top + Inches(0.30), Inches(0.42),
                  Inches(0.42), fill=WHITE, shape=MSO_SHAPE.OVAL)
        btf = bd.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(btf, num, size=16, bold=True, color=color, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, x + Inches(0.22), top + Inches(0.12),
                     cw - Inches(1.05), Inches(0.84), anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=21, bold=True, color=WHITE, first=True,
              line=1.05)
        write(tb.text_frame, sub, size=13.5, color=WHITE, line=1.15)

        y = top + Inches(1.28)
        for pt_txt in points:
            bx = textbox(s, x + Inches(0.26), y, cw - Inches(0.52),
                         Inches(0.72))
            write(bx.text_frame,
                  [("◂  ", True, color, 13), (pt_txt, False, INK, 16)],
                  first=True, line=1.25)
            y += Inches(0.74)

    quote_band(s, Inches(5.62),
               [("نکته‌ای که همه‌چیز را عوض می‌کند:   ", True, ORANGE, 17),
                ("هفته بعد، صاحب چالش یک نفر دیگر است. امروز کمک می‌گیرید، دفعه بعد کمک می‌کنید.",
                 False, WHITE, 20)],
               h=Inches(1.05), align=PP_ALIGN.RIGHT)

    footer(s, 6, TOTAL, label="با هم بلد می‌شویم")
    notes(s, """
    تأکید: تسهیل‌گری چرخشی است — این طرح مالِ یک نفر نیست، مالِ گروه است.
    اگر پرسیدند «مدیر هم در گروه هست؟» بگویید: بله، ولی به‌عنوان عضو گروه همتا.
    """)
    return s


# ═════════════════════════ ۷ — چرایی: مسئله ما ═══════════════════════════
def s07(prs, cfg):
    s = base_slide(prs, WHITE)
    head(s, "مسئله ما کمبود دانش نیست", chapter=1, kicker="چرا اصلاً؟")

    pw = Inches(5.05)
    p, sz = img("06-silo.jpg")
    pic = picture_fill(s, p, M, Inches(1.72), pw, Inches(3.05), sz)
    round_picture(pic)

    cap = card(s, M, Inches(4.92), pw, Inches(1.85), fill=NAVY, line=None)
    rect(s, M + pw - Inches(0.09), Inches(4.92), Inches(0.09), Inches(1.85),
         fill=ORANGE)
    tb = textbox(s, M + Inches(0.26), Inches(5.14), pw - Inches(0.55),
                 Inches(1.5))
    tf = tb.text_frame
    write(tf, "مسئله، گردش دانش است", size=15, bold=True, color=ORANGE,
          first=True, line=1.0)
    write(tf, "همه‌چیز را بلدیم — ولی هرکس تکه‌ای از آن را می‌داند و کسی نمی‌داند چه کسی چه می‌داند.",
          size=17, color=WHITE, space_before=9, line=1.3)

    tx = M + pw + Inches(0.48)
    tw = SW - M - tx

    symptoms = [
        ("یک اشتباه، چند بار", "همان خطا در واحد دیگر دوباره تکرار می‌شود"),
        ("راه‌حل‌های موازی", "دو نفر جدا از هم، ماه‌ها روی یک مسئله کار می‌کنند"),
        ("دانش با آدم‌ها می‌رود", "با هر خداحافظی، سال‌ها تجربه از در بیرون می‌رود"),
        ("تازه‌واردهای سرگردان", "شش ماه طول می‌کشد تا بفهمند از چه کسی بپرسند"),
    ]
    y = Inches(1.72)
    rh = Inches(1.02)
    for i, (title, desc) in enumerate(symptoms):
        rect(s, tx, y, tw, rh, fill=PAPER, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.11)
        rect(s, tx + tw - Inches(0.08), y, Inches(0.08), rh, fill=RED)
        nb = rect(s, tx + tw - Inches(0.84), y + Inches(0.28), Inches(0.44),
                  Inches(0.44), fill=RED, shape=MSO_SHAPE.OVAL)
        ntf = nb.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(ntf, fa(i + 1), size=15, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, tx + Inches(0.28), y, tw - Inches(1.22), rh,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=19, bold=True, color=NAVY, first=True,
              line=1.05)
        write(tb.text_frame, desc, size=15, color=INK_SOFT, line=1.2)
        y += rh + Inches(0.11)

    footer(s, 7, TOTAL, label="با هم بلد می‌شویم")
    notes(s, """
    اینجا حتماً یک مثال واقعی از سازمان خودتان بزنید (بدون بردن نام افراد).
    یک مثال واقعی، از هر چهار مورد این فهرست مؤثرتر است.
    """)
    return s


# ══════════════════════ ۸ — چرایی: علم یادگیری ۷۰-۲۰-۱۰ ══════════════════
def s08(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "بودجه آموزش، جای اشتباهی خرج می‌شود", chapter=1,
         kicker="آنچه پژوهش می‌گوید")

    cx = Inches(9.90)
    cy = Inches(3.72)
    R = Inches(1.62)
    R_IN = Inches(0.97)
    GAP = 2.2
    segs = [(0.0, 252.0, GREEN, 0.0),
            (252.0, 324.0, ORANGE, 0.14),
            (324.0, 360.0, BLUE_MID, 0.0)]
    for a0, a1, color, pop in segs:
        T.donut_segment(s, cx, cy, R + Inches(pop), R_IN + Inches(pop * 0.5),
                        a0 + GAP, a1 - GAP, color)

    ctb = textbox(s, cx - Inches(1.0), cy - Inches(0.52), Inches(2.0),
                  Inches(1.05), anchor=MSO_ANCHOR.MIDDLE)
    write(ctb.text_frame, "۷۰-۲۰-۱۰", size=23, bold=True, color=NAVY,
          first=True, align=PP_ALIGN.CENTER, line=1.0)
    write(ctb.text_frame, "مدل یادگیری", size=12.5, color=INK_SOFT,
          align=PP_ALIGN.CENTER, line=1.1)

    src = textbox(s, cx - Inches(2.0), cy + R + Inches(0.42), Inches(4.0),
                  Inches(0.36))
    write(src.text_frame, "مرجع: مؤسسه CCL آمریکا", size=12.5, color=INK_SOFT,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    tx = M
    tw = Inches(7.45)
    rows = [
        ("۷۰٪", "از انجام کار واقعی", GREEN, GREEN_PALE, False),
        ("۲۰٪", "از تعامل با آدم‌های دیگر", ORANGE, ORANGE_PALE, True),
        ("۱۰٪", "از کلاس و کتاب", BLUE_MID, BLUE_PALE, False),
    ]
    y = Inches(1.74)
    for pct, label, color, pale, hi in rows:
        h = Inches(1.00) if hi else Inches(0.84)
        rect(s, tx, y, tw, h, fill=pale, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.12, line=color if hi else None, line_w=Pt(2))
        rect(s, tx + tw - Inches(0.09), y, Inches(0.09), h, fill=color)
        pb = textbox(s, tx + tw - Inches(1.62), y, Inches(1.42), h,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(pb.text_frame, pct, size=30 if hi else 26, bold=True, color=color,
              first=True, align=PP_ALIGN.CENTER, line=1.0)
        lb = textbox(s, tx + Inches(0.28), y, tw - Inches(2.05), h,
                     anchor=MSO_ANCHOR.MIDDLE)
        chunks = [(label, hi, ORANGE_DEEP if hi else INK, 22 if hi else 20)]
        if hi:
            chunks.append(("     ← اینجا", True, ORANGE, 16))
        write(lb.text_frame, chunks, first=True, line=1.1)
        y += h + Inches(0.15)

    # مقایسه بودجه
    by = Inches(5.02)
    rect(s, tx, by, tw, Inches(1.74), fill=WHITE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09, line=LINE)
    tb = textbox(s, tx + Inches(0.28), by + Inches(0.16), tw - Inches(0.56),
                 Inches(0.4))
    write(tb.text_frame, "و بودجه آموزشی ما کجا خرج می‌شود؟", size=16,
          bold=True, color=NAVY, first=True, line=1.0)

    bar_y = by + Inches(0.66)
    bar_w = tw - Inches(0.56)
    bar_h = Inches(0.42)
    rect(s, tx + Inches(0.28), bar_y, bar_w, bar_h, fill=BLUE_PALE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.30)
    rect(s, tx + Inches(0.28) + bar_w * 0.10, bar_y, bar_w * 0.90, bar_h,
         fill=BLUE_MID, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.30)
    lb = textbox(s, tx + Inches(0.28) + bar_w * 0.10, bar_y, bar_w * 0.90,
                 bar_h, anchor=MSO_ANCHOR.MIDDLE)
    write(lb.text_frame, "۹۰٪ روی کلاس و دوره", size=14.5, bold=True,
          color=WHITE, first=True, align=PP_ALIGN.CENTER, line=1.0)

    nb = textbox(s, tx + Inches(0.28), by + Inches(1.20), bar_w, Inches(0.42))
    write(nb.text_frame,
          [("یعنی ", False, INK_SOFT, 16), ("۹۰٪", True, RED, 16),
           (" بودجه، روی ", False, INK_SOFT, 16), ("۱۰٪", True, RED, 16),
           (" یادگیری. منتورینگ گروهی سراغ آن ", False, INK_SOFT, 16),
           ("۲۰٪ رهاشده", True, ORANGE_DEEP, 16), (" می‌رود.", False, INK_SOFT, 16)],
          first=True, line=1.15)

    footer(s, 8, TOTAL, label="با هم بلد می‌شویم")
    notes(s, """
    این اسلاید پشتوانه علمی طرح است و مخصوصاً برای مدیران نوشته شده.
    اگر عدد دقیق بودجه سازمان را دارید، جایگزین کنید — تأثیرش چند برابر می‌شود.
    """)
    return s


# ═══════════════════ ۹ — چرایی: چرا گروه، نه یک منتور؟ ═══════════════════
def s09(prs, cfg):
    s = base_slide(prs, WHITE)
    head(s, "چرا گروه؟ چرا یک منتور کافی نیست؟", chapter=1,
         kicker="ریاضیاتِ ساده")

    # مقایسه ۱←۱ و گروهی
    gap = Inches(0.40)
    cw = (CONTENT_W - gap) / 2
    top = Inches(1.72)
    ch = Inches(2.28)

    left = [
        (SW - M - cw, "یک منتور", "۱ تجربه، ۱ سبک فکر، ۱ نقطه کور",
         INK_SOFT, T.RGBColor(0xEC, 0xF0, 0xF3), 1),
        (M, "یک گروه ۱۰ نفره", "۱۰ تجربه، ۱۰ سبک فکر، پوشش نقاط کور",
         ORANGE_DEEP, ORANGE_PALE, 10),
    ]
    for x, title, desc, color, pale, dots in left:
        rect(s, x, top, cw, ch, fill=pale, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.09)
        tb = textbox(s, x + Inches(0.30), top + Inches(0.22), cw - Inches(0.60),
                     Inches(0.86))
        write(tb.text_frame, title, size=22, bold=True, color=color, first=True,
              line=1.05)
        write(tb.text_frame, desc, size=15.5, color=INK, line=1.2)
        # نقطه‌های نماینده افراد
        dy = top + Inches(1.32)
        dsz = Inches(0.30)
        dgap = Inches(0.13)
        row_w = min(dots, 5) * (dsz + dgap) - dgap
        for k in range(dots):
            r_i = k // 5
            c_i = k % 5
            cnt = min(dots - r_i * 5, 5)
            rw = cnt * (dsz + dgap) - dgap
            dx = x + (cw - rw) / 2 + c_i * (dsz + dgap)
            rect(s, dx, dy + r_i * (dsz + Inches(0.10)), dsz, dsz,
                 fill=color if dots > 1 else INK_SOFT, shape=MSO_SHAPE.OVAL)

    # سه دلیل
    reasons = [
        ("تنوع دیدگاه", "کسی که مسئله شما را از زاویه‌ای می‌بیند که به ذهنتان نرسیده"),
        ("مقیاس‌پذیری", "منتور باتجربه کم داریم؛ همتا زیاد داریم"),
        ("دوطرفه بودن", "در یک جلسه هم کمک می‌گیرید، هم کمک می‌کنید"),
    ]
    ry = Inches(4.22)
    rgap = Inches(0.30)
    rw = (CONTENT_W - 2 * rgap) / 3
    for i, (title, desc) in enumerate(reasons):
        x = SW - M - rw - i * (rw + rgap)
        rect(s, x, ry, rw, Inches(1.28), fill=WHITE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11, line=LINE)
        rect(s, x + rw - Inches(0.08), ry, Inches(0.08), Inches(1.28),
             fill=NAVY)
        tb = textbox(s, x + Inches(0.24), ry, rw - Inches(0.50), Inches(1.28),
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=18.5, bold=True, color=NAVY,
              first=True, line=1.05)
        write(tb.text_frame, desc, size=14.5, color=INK_SOFT, line=1.2)

    quote_band(s, Inches(5.78),
               [("یک منتور به شما می‌گوید ", False, WHITE, 20),
                ("او چه کرد", True, BLUE_LIGHT, 20),
                (". یک گروه به شما نشان می‌دهد ", False, WHITE, 20),
                ("چند راه وجود دارد", True, ORANGE, 20), (".", False, WHITE, 20)],
               h=Inches(0.96))

    footer(s, 9, TOTAL, label="با هم بلد می‌شویم")
    notes(s, """
    استدلال «مقیاس‌پذیری» را برای مدیران پررنگ کنید: منتور باتجربه کمیاب است،
    ولی هر ۱۰ نفر می‌توانند برای هم گروه همتا باشند.
    """)
    return s


# ═════════════════════ ۱۰ — چرایی: دستاورد برای من و سازمان ══════════════
def s10(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "چه چیزی عاید چه کسی می‌شود؟", chapter=1, kicker="دستاورد")

    gap = Inches(0.40)
    cw = (CONTENT_W - gap) / 2
    top = Inches(1.72)
    ch = Inches(0.64)
    rowh = Inches(0.66)
    rows = 4
    boxh = ch + rows * rowh + Inches(0.22)

    cols = [
        (M + cw + gap, "برای شما", NAVY, BLUE_PALE, [
            ("مسئله‌تان زودتر حل می‌شود", "به‌جای هفته‌ها کلنجار، ۹۰ دقیقه با ۱۰ ذهن"),
            ("مهارت پرسیدن یاد می‌گیرید", "در جلسات و مذاکرات کاری هم به کارتان می‌آید"),
            ("شبکه‌تان گسترده می‌شود", "می‌دانید برای هر مسئله سراغ چه کسی بروید"),
            ("کمتر احساس تنهایی می‌کنید", "می‌فهمید دیگران هم همین را تجربه کرده‌اند"),
        ]),
        (M, "برای سازمان", ORANGE_DEEP, ORANGE_PALE, [
            ("تکرار اشتباهات کم می‌شود", "تجربه شکست یک واحد، درسِ واحد دیگر می‌شود"),
            ("دانش ماندگار می‌شود", "تجربه از سرِ افراد به حافظه جمعی منتقل می‌شود"),
            ("تازه‌واردها زودتر جا می‌افتند", "مسیر شش‌ماهه، چندهفته‌ای طی می‌شود"),
            ("افراد می‌مانند", "آدم‌ها جایی که رشد می‌کنند را ترک نمی‌کنند"),
        ]),
    ]

    for x, title, cmain, cpale, items in cols:
        card(s, x, top, cw, boxh, fill=WHITE, line=LINE)
        rect(s, x, top, cw, ch, fill=cmain, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.11)
        rect(s, x, top + ch - Inches(0.16), cw, Inches(0.16), fill=cmain)
        tb = textbox(s, x, top, cw, ch, anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=22, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        for i, (main, sub) in enumerate(items):
            y = top + ch + Inches(0.11) + i * rowh
            if i % 2 == 0:
                rect(s, x + Inches(0.13), y, cw - Inches(0.26),
                     rowh - Inches(0.05), fill=cpale,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
            tb = textbox(s, x + Inches(0.30), y, cw - Inches(0.60),
                         rowh - Inches(0.05), anchor=MSO_ANCHOR.MIDDLE)
            write(tb.text_frame, [("✓  ", True, cmain, 14),
                                  (main, True, INK, 18)], first=True, line=1.05)
            write(tb.text_frame, sub, size=13, color=INK_SOFT, line=1.15)

    quote_band(s, Inches(5.72),
               [("هزینه‌اش؟  ", True, ORANGE, 19),
                ("۹۰ دقیقه، هر دو هفته یک بار — یعنی حدود ", False, WHITE, 21),
                ("یک درصد", True, ORANGE, 21), (" وقت کاری شما.", False, WHITE, 21)],
               h=Inches(1.0))

    footer(s, 10, TOTAL, label="با هم بلد می‌شویم")
    notes(s, """
    اگر مخاطب بیشتر کارشناس است، ستون راست را کامل بخوانید و چپ را سریع رد شوید.
    اگر مدیران حاضرند، برعکس.
    """)
    return s


# ═══════════════════ ۱۱ — اجرا: آناتومی یک جلسه ۹۰ دقیقه‌ای ══════════════
def s11(prs, cfg):
    s = base_slide(prs, WHITE)
    head(s, "یک جلسه واقعی چطور پیش می‌رود؟", chapter=2,
         kicker="آناتومی ۹۰ دقیقه")

    steps = [
        (10, "دایره باز", "هر نفر یک جمله: این هفته چطور بود", BLUE_LIGHT),
        (10, "مرور تعهدها", "جلسه قبل چه قول دادیم و چه شد", BLUE_LIGHT),
        (10, "روایت چالش", "داوطلب مسئله‌اش را می‌گوید، بدون قطع شدن", BLUE_MID),
        (20, "فقط سؤال", "گروه فقط می‌پرسد — هیچ راه‌حلی ممنوع", ORANGE),
        (20, "تجربه‌های مشابه", "«من هم یک بار…» — روایت، نه دستور", ORANGE),
        (15, "انتخاب و تعهد", "صاحب چالش می‌گوید چه می‌کند", NAVY),
        (5, "دایره بسته", "هر نفر یک کلمه: چه چیزی با خود می‌برد", NAVY),
    ]

    bar_y = Inches(1.72)
    bar_h = Inches(0.40)
    total_min = sum(x[0] for x in steps)
    xcur = SW - M
    for mins, title, desc, color in steps:
        w = CONTENT_W * mins / total_min
        rect(s, xcur - w, bar_y, w - Inches(0.035), bar_h, fill=color)
        tb = textbox(s, xcur - w, bar_y, w, bar_h, anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, fa(mins), size=13.5, bold=True, color=WHITE,
              first=True, align=PP_ALIGN.CENTER, line=1.0)
        xcur -= w

    y = Inches(2.32)
    rh = Inches(0.535)
    for i, (mins, title, desc, color) in enumerate(steps):
        if i % 2 == 0:
            rect(s, M, y, CONTENT_W, rh - Inches(0.04), fill=PAPER,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.20)
        chip = rect(s, SW - M - Inches(1.24), y + Inches(0.05), Inches(1.08),
                    rh - Inches(0.14), fill=color,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.26)
        ctf = chip.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(ctf, fa("%d′" % mins), size=14, bold=True, color=WHITE,
              first=True, align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, SW - M - Inches(4.95), y, Inches(3.5),
                     rh - Inches(0.04), anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=20, bold=True, color=NAVY, first=True,
              line=1.0)
        db = textbox(s, M + Inches(0.2), y, Inches(7.6), rh - Inches(0.04),
                     anchor=MSO_ANCHOR.MIDDLE)
        write(db.text_frame, desc, size=17, color=INK_SOFT, first=True,
              line=1.0)
        y += rh

    gy = Inches(6.14)
    rect(s, M, gy, CONTENT_W, Inches(0.66), fill=ORANGE_PALE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.13, line=ORANGE,
         line_w=Pt(1.75))
    tb = textbox(s, M, gy, CONTENT_W, Inches(0.66), anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame,
          [("۴۰ دقیقه از ۹۰ دقیقه — یعنی نیمی از جلسه — ", False, INK, 18),
           ("فقط سؤال و روایت است", True, ORANGE_DEEP, 19),
           ("، نه راه‌حل.", False, INK, 18)],
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    footer(s, 11, TOTAL, label="با هم بلد می‌شویم")
    notes(s, """
    بخش «فقط سؤال» عجیب‌ترین قسمت است و بیشترین مقاومت را دارد.
    بگویید: «می‌دانم سخت است. اسلاید بعد نشان می‌دهم چرا ارزشش را دارد.»
    """)
    return s


# ══════════════════ ۱۲ — اجرا: نصیحت ← پرسش (قلب روش) ════════════════════
def s12(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "قلب روش: به‌جای نصیحت، بپرس", chapter=2, kicker="مهارت کلیدی")

    pw = Inches(3.75)
    p, sz = img("12-listen.jpg")
    pic = picture_fill(s, p, M, Inches(1.72), pw, Inches(2.60), sz)
    round_picture(pic)

    nb = card(s, M, Inches(4.46), pw, Inches(2.30), fill=NAVY, line=None)
    rect(s, M + pw - Inches(0.09), Inches(4.46), Inches(0.09), Inches(2.30),
         fill=ORANGE)
    tb = textbox(s, M + Inches(0.24), Inches(4.68), pw - Inches(0.52),
                 Inches(1.9))
    tf = tb.text_frame
    write(tf, "چرا نصیحت جواب نمی‌دهد؟", size=15, bold=True, color=ORANGE,
          first=True, line=1.0)
    write(tf, "چون شما تمام جزئیات مسئله او را نمی‌دانید.",
          size=16.5, color=WHITE, space_before=9, line=1.28)
    write(tf, "ولی یک سؤال خوب، خودِ او را به جوابی می‌رساند که خودش اجرایش می‌کند.",
          size=16.5, color=BLUE_PALE, space_before=7, line=1.28)

    tx = M + pw + Inches(0.45)
    tw = SW - M - tx

    # سرستون‌ها
    gap = Inches(0.32)
    hw = (tw - gap) / 2
    hy = Inches(1.72)
    heads = [(tx + hw + gap, "به‌جای این…", RED), (tx, "این را بگویید", GREEN)]
    for x, label, color in heads:
        rect(s, x, hy, hw, Inches(0.54), fill=color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
        tb = textbox(s, x, hy, hw, Inches(0.54), anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, label, size=17, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)

    pairs = [
        ("«جای تو بودم فلان کار را می‌کردم.»", "«خودت چه گزینه‌هایی را سنجیدی؟»"),
        ("«این روش اشتباه است.»", "«اگر جواب ندهد، بدترین اتفاق چیست؟»"),
        ("«ما همیشه اینجوری کردیم.»", "«چه چیزی این مسئله را برایت سخت کرده؟»"),
        ("«نگران نباش، درست می‌شود.»", "«موفقیت در این مسئله برایت یعنی چه؟»"),
    ]
    y = Inches(2.42)
    rh = Inches(0.96)
    for bad, good in pairs:
        for x, txt, color, pale, italic in [
            (tx + hw + gap, bad, RED, RED_PALE, True),
            (tx, good, GREEN, GREEN_PALE, False),
        ]:
            rect(s, x, y, hw, rh - Inches(0.10), fill=pale,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.13)
            tb = textbox(s, x + Inches(0.20), y, hw - Inches(0.40),
                         rh - Inches(0.10), anchor=MSO_ANCHOR.MIDDLE)
            write(tb.text_frame, txt, size=16.5, bold=not italic,
                  color=INK_SOFT if italic else INK, first=True,
                  align=PP_ALIGN.CENTER, line=1.25)
        y += rh

    footer(s, 12, TOTAL, label="با هم بلد می‌شویم")
    notes(s, """
    این مهم‌ترین اسلاید ارائه است. آرام بخوانید و بعد از هر جفت، مکث کوتاه کنید.
    از گروه بخواهید یک نمونه دیگر بسازند — اگر توانستند، یعنی روش را گرفته‌اند.
    """)
    return s


# ══════════════════════ ۱۳ — اجرا: قواعد فضای امن ════════════════════════
def s13(prs, cfg):
    s = base_slide(prs, WHITE)
    head(s, "پنج قاعده‌ای که فضا را امن می‌کند", chapter=2,
         kicker="پیمان گروه")

    rules = [
        ("۱", "اینجا سِمَت نداریم", "در این ۹۰ دقیقه هیچ‌کس رئیس کسی نیست", ORANGE, ORANGE_PALE),
        ("۲", "هرچه گفته شد، همین‌جا می‌ماند", "بدون استثنا، بدون صورت‌جلسه", NAVY, BLUE_PALE),
        ("۳", "سؤال بپرس، نسخه نپیچ", "کنجکاوی به‌جای نصیحت", GREEN, GREEN_PALE),
        ("۴", "هیچ سؤالی احمقانه نیست", "و هیچ تجربه‌ای بی‌ارزش نیست", PURPLE,
         T.RGBColor(0xF6, 0xEF, 0xFA)),
        ("۵", "گوش بده تا بفهمی", "نه اینکه منتظر نوبت حرف زدنت باشی", RED, RED_PALE),
    ]

    tw = Inches(8.15)
    tx = SW - M - tw
    y = Inches(1.72)
    rh = Inches(0.92)
    for num, title, desc, color, pale in rules:
        rect(s, tx, y, tw, rh, fill=pale, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.11)
        rect(s, tx + tw - Inches(0.10), y, Inches(0.10), rh, fill=color)
        bd = rect(s, tx + tw - Inches(0.94), y + Inches(0.22), Inches(0.48),
                  Inches(0.48), fill=color, shape=MSO_SHAPE.OVAL)
        btf = bd.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(btf, num, size=17, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, tx + Inches(0.28), y, tw - Inches(1.32), rh,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=22, bold=True, color=color, first=True,
              line=1.05)
        write(tb.text_frame, desc, size=16, color=INK_SOFT, line=1.15)
        y += rh + Inches(0.10)

    px = M
    pw = tx - M - Inches(0.42)
    p, sz = img("13-safe.jpg")
    pic = picture_fill(s, p, px, Inches(1.72), pw, Inches(3.10), sz)
    round_picture(pic)

    card(s, px, Inches(4.98), pw, Inches(1.80), fill=NAVY, line=None)
    rect(s, px + pw - Inches(0.09), Inches(4.98), Inches(0.09), Inches(1.80),
         fill=ORANGE)
    tb = textbox(s, px + Inches(0.24), Inches(5.18), pw - Inches(0.52),
                 Inches(1.45))
    tf = tb.text_frame
    write(tf, "چرا این پنج تا؟", size=14.5, bold=True, color=ORANGE, first=True,
          line=1.0)
    write(tf, "تا وقتی کسی نترسد از گفتنِ «بلد نیستم»، هیچ یادگیری‌ای شروع نمی‌شود.",
          size=16, color=WHITE, space_before=8, line=1.26)

    footer(s, 13, TOTAL, label="با هم بلد می‌شویم")
    notes(s, """
    این قواعد را نخوانید — از گروه بخواهید بلند بخوانند و تأیید کنند.
    پیشنهاد: روی یک برگه بزرگ بنویسید و هر جلسه روی دیوار اتاق بزنید.
    """)
    return s


# ═════════════════════════ ۱۴ — اجرا: چرخه شش‌ماهه ═══════════════════════
def s14(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "شش ماه، دوازده جلسه، یک گروه", chapter=2, kicker="مسیر پیشِ رو")

    phases = [
        ("جلسه ۱", "شکل‌گیری", "آشنایی، بستن پیمان گروه، انتخاب اولین داوطلب", ORANGE),
        ("جلسه ۲ تا ۴", "تمرین", "قواعد هنوز سخت است؛ تسهیل‌گر بیشتر دخالت می‌کند", BLUE_MID),
        ("جلسه ۵ تا ۱۰", "اوج‌گیری", "گروه روان می‌شود؛ تسهیل‌گری چرخشی آغاز می‌شود", NAVY),
        ("جلسه ۱۱ و ۱۲", "جمع‌بندی", "مرور دستاوردها و تصمیم برای دوره بعد", GREEN),
    ]

    tw = Inches(7.95)
    tx = SW - M - tw
    y = Inches(1.72)
    rh = Inches(1.06)
    for i, (when, title, desc, color) in enumerate(phases):
        rect(s, tx, y, tw, rh, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.11, line=LINE)
        rect(s, tx + tw - Inches(0.09), y, Inches(0.09), rh, fill=color)
        wb = rect(s, tx + tw - Inches(1.95), y + Inches(0.26), Inches(1.66),
                  rh - Inches(0.52), fill=color,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.26)
        wtf = wb.text_frame
        wtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(wtf, when, size=14, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, tx + Inches(0.28), y, tw - Inches(2.35), rh,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=20, bold=True, color=color, first=True,
              line=1.05)
        write(tb.text_frame, desc, size=15, color=INK_SOFT, line=1.2)
        if i < len(phases) - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                    tx + tw - Inches(1.20), y + rh + Inches(0.01),
                                    Inches(0.22), Inches(0.16))
            ar.fill.solid(); ar.fill.fore_color.rgb = LINE
            no_line(ar); shadow_off(ar)
        y += rh + Inches(0.18)

    # ستون «چه چیزی از شما می‌خواهیم»
    px = M
    pw = tx - M - Inches(0.42)
    card(s, px, Inches(1.72), pw, Inches(3.06), fill=WHITE, line=LINE)
    rect(s, px, Inches(1.72), pw, Inches(0.60), fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.13)
    rect(s, px, Inches(2.17), pw, Inches(0.15), fill=NAVY)
    tb = textbox(s, px, Inches(1.72), pw, Inches(0.60), anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame, "تعهد شما", size=19, bold=True, color=WHITE,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    asks = ["حضور منظم در ۱۲ جلسه",
            "آوردن یک چالش واقعی",
            "رعایت محرمانگی",
            "یک اقدام کوچک بین دو جلسه"]
    yy = Inches(2.48)
    for a in asks:
        bx = textbox(s, px + Inches(0.24), yy, pw - Inches(0.48), Inches(0.56))
        write(bx.text_frame, [("◂  ", True, ORANGE, 13), (a, False, INK, 16)],
              first=True, line=1.2)
        yy += Inches(0.58)

    card(s, px, Inches(4.94), pw, Inches(1.84), fill=ORANGE_PALE, line=ORANGE)
    tb = textbox(s, px + Inches(0.24), Inches(5.14), pw - Inches(0.48),
                 Inches(1.5))
    tf = tb.text_frame
    write(tf, "و اگر نشد؟", size=14.5, bold=True, color=ORANGE_DEEP, first=True,
          line=1.0)
    write(tf, "بعد از جلسه سوم، هر کس خواست می‌تواند کنار بکشد — بدون توضیح.",
          size=16, color=NAVY, space_before=8, line=1.26)

    footer(s, 14, TOTAL, label="با هم بلد می‌شویم")
    notes(s, """
    جمله «بدون توضیح می‌توانید کنار بکشید» را حتماً بگویید.
    داوطلبانه بودن، شرط کارکردنِ این روش است — نه یک تعارف.
    """)
    return s


# ══════════════════════ ۱۵ — تجربه زنده (اینجا و اکنون) ══════════════════
def s15(prs, cfg):
    s = base_slide(prs, WHITE)
    gradient(s, 0, 0, SW, SH, NAVY_DEEP, NAVY, angle=315)

    kb = textbox(s, M, Inches(0.52), CONTENT_W, Inches(0.4))
    write(kb.text_frame, "حالا خودمان انجامش می‌دهیم", size=16, bold=True,
          color=ORANGE, first=True, align=PP_ALIGN.CENTER, line=1.0)

    tb = textbox(s, M, Inches(0.96), CONTENT_W, Inches(0.68))
    write(tb.text_frame, "یک دور کوتاه — ۳۰ دقیقه", size=34, bold=True,
          color=WHITE, first=True, align=PP_ALIGN.CENTER, line=1.0)

    steps = [
        ("۵ دقیقه", "یک داوطلب", "چالشش را در سه جمله می‌گوید"),
        ("۱۰ دقیقه", "فقط سؤال", "هر کس یک سؤال — راه‌حل ممنوع"),
        ("۱۰ دقیقه", "تجربه مشابه", "«من هم یک بار…» در یک دقیقه"),
        ("۵ دقیقه", "یک تصمیم", "داوطلب می‌گوید چه می‌کند"),
    ]
    gap = Inches(0.30)
    cw = (CONTENT_W - 3 * gap) / 4
    y = Inches(2.02)
    ch = Inches(2.05)
    for i, (mins, title, desc) in enumerate(steps):
        x = SW - M - cw - i * (cw + gap)
        rect(s, x, y, cw, ch, fill=WHITE, alpha=0.10,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09,
             line=BLUE_LIGHT, line_w=Pt(1))
        chip = rect(s, x + (cw - Inches(1.32)) / 2, y + Inches(0.24),
                    Inches(1.32), Inches(0.50), fill=ORANGE,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.26)
        ctf = chip.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(ctf, mins, size=14.5, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, x + Inches(0.18), y + Inches(0.88), cw - Inches(0.36),
                     Inches(1.0))
        write(tb.text_frame, title, size=19, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.1)
        write(tb.text_frame, desc, size=14, color=BLUE_PALE,
              align=PP_ALIGN.CENTER, line=1.2)

    # قاعده امروز
    ry = Inches(4.42)
    rect(s, M, ry, CONTENT_W, Inches(0.86), fill=ORANGE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11)
    tb = textbox(s, M, ry, CONTENT_W, Inches(0.86), anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame,
          [("تنها قاعده امروز:  ", True, WHITE, 18),
           ("در ده دقیقه دوم، هیچ‌کس حق ندارد جمله‌ای بگوید که علامت سؤال ندارد.",
            True, WHITE, 21)],
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    p, sz = img("15-start.jpg")
    pw = Inches(3.55)
    pic = picture_fill(s, p, M, Inches(5.52), pw, Inches(1.32), sz)
    round_picture(pic)

    bx = M + pw + Inches(0.32)
    bw = SW - M - bx
    rect(s, bx, Inches(5.52), bw, Inches(1.32), fill=WHITE, alpha=0.10,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09, line=BLUE_LIGHT)
    tb = textbox(s, bx + Inches(0.30), Inches(5.52), bw - Inches(0.60),
                 Inches(1.32), anchor=MSO_ANCHOR.MIDDLE)
    tf = tb.text_frame
    write(tf, "چه کسی چالشش را می‌آورد؟", size=21, bold=True, color=ORANGE,
          first=True, line=1.1)
    write(tf, "لازم نیست مسئله بزرگی باشد — همان چیزی که اول جلسه در ذهن نگه داشتید کافی است.",
          size=14, color=BLUE_PALE, line=1.2)

    notes(s, """
    این مهم‌ترین ۳۰ دقیقه امروز است. تایمر بگذارید و سفت نگه دارید.
    اگر کسی در بخش «فقط سؤال» راه‌حل داد، با لبخند وسط حرفش بروید:
    «این را نگه دار برای ده دقیقه بعد.»
    اگر داوطلب پیدا نشد، خودتان یک چالش واقعی از کار خودتان بیاورید.
    """)
    return s


# ════════════════════════════ ۱۶ — گام بعدی ══════════════════════════════
def s16(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "قدم بعدی چیست؟", kicker="جمع‌بندی")

    # سه ستون اقدام
    actions = [
        ("امروز", "همین‌جا", "برگه ثبت‌نام را امضا کنید", ORANGE, ORANGE_PALE),
        ("این هفته", cfg.first_session, "اولین جلسه رسمی گروه", NAVY, BLUE_PALE),
        ("شش ماه بعد", "مرور با هم", "می‌بینیم چه چیزی عوض شده", GREEN, GREEN_PALE),
    ]
    gap = Inches(0.34)
    cw = (CONTENT_W - 2 * gap) / 3
    y = Inches(1.68)
    ch = Inches(1.88)
    for i, (when, what, desc, color, pale) in enumerate(actions):
        x = SW - M - cw - i * (cw + gap)
        rect(s, x, y, cw, ch, fill=pale, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.10)
        rect(s, x, y, cw, Inches(0.07), fill=color)
        tb = textbox(s, x + Inches(0.22), y + Inches(0.26), cw - Inches(0.44),
                     Inches(1.45))
        tf = tb.text_frame
        write(tf, when, size=14, bold=True, color=color, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        write(tf, what, size=23, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
              space_before=5, line=1.15)
        write(tf, desc, size=15, color=INK_SOFT, align=PP_ALIGN.CENTER,
              space_before=4, line=1.2)

    # جزئیات جلسه
    dy = Inches(3.78)
    infos = [("ساعت", cfg.time), ("مکان", cfg.place),
             ("ظرفیت", "۱۲ نفر"), ("تماس", cfg.contact)]
    igap = Inches(0.28)
    iw = (CONTENT_W - 3 * igap) / 4
    for i, (label, value) in enumerate(infos):
        x = SW - M - iw - i * (iw + igap)
        rect(s, x, dy, iw, Inches(1.05), fill=WHITE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11, line=LINE)
        tb = textbox(s, x + Inches(0.14), dy, iw - Inches(0.28), Inches(1.05),
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, label, size=13, bold=True, color=ORANGE,
              first=True, align=PP_ALIGN.CENTER, line=1.1)
        write(tb.text_frame, value, size=18, bold=True, color=NAVY,
              align=PP_ALIGN.CENTER, line=1.2)

    # جمله پایانی
    fy = Inches(5.20)
    rect(s, M, fy, CONTENT_W, Inches(1.58), fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    rect(s, SW - M - Inches(0.10), fy, Inches(0.10), Inches(1.58), fill=ORANGE)
    tb = textbox(s, M + Inches(0.5), fy, CONTENT_W - Inches(1.0), Inches(1.58),
                 anchor=MSO_ANCHOR.MIDDLE)
    tf = tb.text_frame
    write(tf, "ما آدم‌های باتجربه‌ای را استخدام کردیم و کنار هم نشاندیم.",
          size=21, color=WHITE, first=True, align=PP_ALIGN.CENTER, line=1.3)
    write(tf, "فقط یادمان رفت فرصتی بسازیم که با هم حرف بزنند.",
          size=21, color=WHITE, align=PP_ALIGN.CENTER, line=1.3)
    write(tf, "این طرح، همان فرصت است.", size=24, bold=True, color=ORANGE,
          align=PP_ALIGN.CENTER, space_before=8, line=1.3)

    footer(s, 16, TOTAL, label="با هم بلد می‌شویم")
    notes(s, """
    نگویید «سؤالی نیست؟». بگویید: «چه کسی می‌خواهد نفر اول باشد؟»
    برگه ثبت‌نام را همان لحظه دست به دست کنید — نه ایمیل کنید، نه فردا.
    اگر کسی گفت «فکر می‌کنم و خبر می‌دهم»، احترام بگذارید و اصرار نکنید.
    """)
    return s


# ═══════════════════════════════════════════════════════ اجرا ════════════
BUILDERS = [s01, s02, s03, s04, s05, s06, s07, s08,
            s09, s10, s11, s12, s13, s14, s15, s16]


def build(cfg):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    T.set_font(cfg.font)
    for fn in BUILDERS:
        fn(prs, cfg)
    prs.save(cfg.out)
    return cfg.out


def main():
    ap = argparse.ArgumentParser(description="ساخت ارائه افتتاحیه منتورینگ گروهی")
    ap.add_argument("--out", default=os.path.join(HERE, "با-هم-بلد-می‌شویم.pptx"))
    ap.add_argument("--font", default="IRANSans")
    ap.add_argument("--date", default="[تاریخ]")
    ap.add_argument("--presenter", default="[نام تسهیل‌گر]")
    ap.add_argument("--org", default="[نام واحد / سازمان]")
    ap.add_argument("--first-session", dest="first_session", default="[تاریخ]")
    ap.add_argument("--time", default="[ساعت]")
    ap.add_argument("--place", default="[مکان]")
    ap.add_argument("--contact", default="[داخلی]")
    cfg = ap.parse_args()
    print("ساخته شد:", build(cfg))


if __name__ == "__main__":
    main()
