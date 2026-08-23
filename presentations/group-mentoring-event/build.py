# -*- coding: utf-8 -*-
"""
ارائه جامع رویداد منتورینگ گروهی — «جایی که با هم بلد می‌شویم»

نسخه تصویرمحور: هر بخش با یک اسلاید تمام‌قاب باز می‌شود و همه تصاویر
تولیدشده برای این رویداد در روایت به کار رفته‌اند.

اجرا:  python3 build.py [--presenter ...] [--font IRANSans] [--out out.pptx]
"""
import argparse
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                "..", "_lib"))
import theme as T
from theme import (
    NAVY, NAVY_DEEP, BLUE_MID, BLUE_LIGHT, BLUE_PALE, ORANGE, ORANGE_DEEP,
    ORANGE_PALE, WHITE, INK, INK_SOFT, PAPER, GREEN, GREEN_PALE, RED, RED_PALE,
    PURPLE, AMBER, LINE, SW, SH, M, CONTENT_W,
    fa, rect, card, gradient, hline, textbox, write, picture_fill,
    round_picture, base_slide, blank, footer, notes, no_line, shadow_off,
    scrim, scrim_gradient, full_bleed,
)

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")
TOTAL = 20
LABEL = "منتورینگ گروهی"

# پنج بخش ارائه
PARTS = ["چرا", "چیست", "چطور", "چه می‌شود", "شروع"]


def img(name):
    p = os.path.join(IMG, name)
    return p, Image.open(p).size


# ═══════════════════════════════════════════════ کمک‌کارهای اسلاید ═══════
def head(slide, title, part=None, kicker=None, on_dark=False):
    """سربرگ با نشانگر بخش."""
    if part is not None:
        cx = M
        for i, name in enumerate(PARTS):
            active = (i == part)
            w = Inches(0.95)
            x = cx + i * (w + Inches(0.08))
            rect(slide, x, Inches(0.46), w, Inches(0.32),
                 fill=NAVY if active else T.RGBColor(0xE4, 0xEC, 0xF2),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.44)
            tb = textbox(slide, x, Inches(0.46), w, Inches(0.32),
                         anchor=MSO_ANCHOR.MIDDLE)
            write(tb.text_frame, name, size=10.5, bold=True,
                  color=WHITE if active else INK_SOFT, first=True,
                  align=PP_ALIGN.CENTER, line=1.0)

    top = Inches(0.40)
    if kicker:
        kb = textbox(slide, M, Inches(0.34), CONTENT_W, Inches(0.34))
        write(kb.text_frame, kicker, size=15.5, bold=True, color=ORANGE,
              first=True, line=1.0)
        top = Inches(0.68)

    tb = textbox(slide, M, top, CONTENT_W, Inches(0.70))
    write(tb.text_frame, title, size=31, bold=True,
          color=WHITE if on_dark else NAVY, first=True, line=1.05)
    hline(slide, SW - M - Inches(1.2), Inches(1.36), Inches(1.2), ORANGE,
          Pt(4.5))
    return slide


def section(prs, num, title, subtitle, image, part_idx, side="right"):
    """
    اسلاید جداکننده بخش: تصویر تمام‌قاب + پرده گرادیانتی + عنوان بزرگ.
    side مشخص می‌کند ستون متن کدام سمت باشد تا اسلایدها یکنواخت نشوند.
    """
    s = blank(prs)
    p, sz = img(image)
    full_bleed(s, p, sz)

    if side == "right":
        # متن سمت راست، تصویر از چپ دیده می‌شود
        scrim_gradient(s, 0, 0, SW, SH, angle=0, a_from=0.10, a_to=0.94)
        tx = SW - M - Inches(6.3)
    else:
        scrim_gradient(s, 0, 0, SW, SH, angle=0, a_from=0.94, a_to=0.10)
        tx = M
    scrim(s, 0, 0, SW, SH, alpha=0.22)

    tw = Inches(6.3)
    nb = textbox(s, tx, Inches(2.18), tw, Inches(0.62))
    write(nb.text_frame, "بخش " + fa(num), size=17, bold=True, color=ORANGE,
          first=True, line=1.0)

    hline(s, tx + tw - Inches(1.3), Inches(2.86), Inches(1.3), ORANGE, Pt(4.5))

    tb = textbox(s, tx, Inches(3.12), tw, Inches(1.9))
    tf = tb.text_frame
    for i, ln in enumerate(title.split("\n")):
        write(tf, ln, size=44, bold=True, color=WHITE, first=(i == 0),
              line=1.18)
    write(tf, subtitle, size=19, color=BLUE_PALE, space_before=14, line=1.3)

    # نشانگر پیشرفت بخش‌ها
    py = SH - Inches(0.90)
    for i, name in enumerate(PARTS):
        w = Inches(0.95)
        x = tx + tw - w - i * (w + Inches(0.10))
        active = (i == part_idx)
        rect(s, x, py, w, Inches(0.32), fill=ORANGE if active else WHITE,
             alpha=None if active else 0.22,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.44)
        lb = textbox(s, x, py, w, Inches(0.32), anchor=MSO_ANCHOR.MIDDLE)
        write(lb.text_frame, name, size=10.5, bold=True,
              color=WHITE if active else BLUE_PALE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
    return s


def band(slide, y, chunks, h=Inches(0.94), fill=NAVY, accent=ORANGE,
         size=21, align=PP_ALIGN.CENTER, x=M, w=None):
    w = w or CONTENT_W
    rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.10)
    rect(slide, x + w - Inches(0.09), y, Inches(0.09), h, fill=accent)
    tb = textbox(slide, x + Inches(0.34), y, w - Inches(0.68), h,
                 anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame, chunks, size=size, color=WHITE, first=True,
          align=align, line=1.2)
    return slide


def photo_card(slide, x, y, w, h, image, caption=None, cap_h=Inches(0.62)):
    """تصویر گوشه‌گرد با نوار عنوان اختیاری روی لبه پایین."""
    p, sz = img(image)
    pic = picture_fill(slide, p, x, y, w, h, sz)
    round_picture(pic)
    if caption:
        cy = y + h - cap_h
        scrim(slide, x, cy, w, cap_h, alpha=0.72)
        tb = textbox(slide, x + Inches(0.20), cy, w - Inches(0.40), cap_h,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, caption, size=14, bold=True, color=WHITE,
              first=True, align=PP_ALIGN.CENTER, line=1.15)
    return pic


# ═════════════════════════════════════════════════ ۱ — جلد ═══════════════
def s01(prs, cfg):
    s = blank(prs)
    p, sz = img("01-hands-together.jpg")
    full_bleed(s, p, sz)
    scrim_gradient(s, 0, 0, SW, SH, angle=0, a_from=0.08, a_to=0.96)
    scrim(s, 0, 0, SW, SH, alpha=0.15)

    tw = Inches(6.6)
    tx = SW - M - tw

    hline(s, tx + tw - Inches(1.4), Inches(1.72), Inches(1.4), ORANGE, Pt(5))

    tb = textbox(s, tx, Inches(2.04), tw, Inches(3.4))
    tf = tb.text_frame
    write(tf, "رویداد معرفی منتورینگ گروهی", size=18, bold=True,
          color=BLUE_LIGHT, first=True, line=1.0)
    write(tf, "جایی که", size=48, bold=True, color=WHITE, space_before=16,
          line=1.22)
    write(tf, [("با هم ", True, WHITE, 48), ("بلد", True, ORANGE, 48),
               (" می‌شویم", True, WHITE, 48)], line=1.22)
    write(tf, "چرا این کار را می‌کنیم، چیست، چطور اجرا می‌شود",
          size=19, color=BLUE_PALE, space_before=20, line=1.3)

    hline(s, tx + tw - Inches(3.2), Inches(5.52), Inches(3.2), BLUE_LIGHT,
          Pt(1.25))
    mb = textbox(s, tx, Inches(5.72), tw, Inches(0.62))
    write(mb.text_frame,
          [("تهیه و تنظیم:  ", False, BLUE_LIGHT, 15.5),
           (cfg.presenter, True, WHITE, 19)], first=True, line=1.15)

    db = textbox(s, tx, Inches(6.30), tw, Inches(0.40))
    write(db.text_frame,
          [("تاریخ:  ", False, BLUE_LIGHT, 14), (cfg.date, True, WHITE, 14),
           ("        ", False, BLUE_LIGHT, 14),
           (cfg.org, False, BLUE_LIGHT, 14)], first=True, line=1.15)

    # آدرس سایت، در نوار نارنجی گوشه پایین
    site_w = Inches(2.62)
    site_h = Inches(0.46)
    rect(s, tx + tw - site_w, Inches(6.76), site_w, site_h, fill=ORANGE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.22)
    sb = textbox(s, tx + tw - site_w, Inches(6.76), site_w, site_h,
                 anchor=MSO_ANCHOR.MIDDLE)
    write(sb.text_frame, "www.coachroom.ir", size=15, bold=True, color=WHITE,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    notes(s, """
    صندلی‌ها را دایره‌ای بچینید و خودتان هم روی یکی از آن‌ها بنشینید.
    جمله شروع: «این آخرین جلسه‌ای است که من جلوی شما می‌ایستم و حرف می‌زنم.»
    """)
    return s


# ══════════════════════════════════════════ ۲ — نقشه امروز ═══════════════
def s02(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "نقشه امروز", kicker="۹۰ دقیقه پیشِ رو")

    parts = [
        ("۱", "چرا", "مسئله‌ای که داریم و کسی درباره‌اش حرف نمی‌زند", RED),
        ("۲", "چیست", "این روش دقیقاً چیست و چه چیزی نیست", NAVY),
        ("۳", "چطور", "یک جلسه واقعی، دقیقه به دقیقه", BLUE_MID),
        ("۴", "چه می‌شود", "چه چیزی عاید شما و سازمان می‌شود", GREEN),
        ("۵", "شروع", "یک دور واقعی، همین امروز", ORANGE),
    ]

    gap = Inches(0.24)
    cw = (CONTENT_W - 4 * gap) / 5
    y = Inches(1.68)
    ch = Inches(2.42)
    for i, (num, title, desc, color) in enumerate(parts):
        x = SW - M - cw - i * (cw + gap)
        card(s, x, y, cw, ch, fill=WHITE, line=LINE)
        rect(s, x, y, cw, Inches(0.07), fill=color)
        nb = rect(s, x + (cw - Inches(0.54)) / 2, y + Inches(0.30),
                  Inches(0.54), Inches(0.54), fill=color, shape=MSO_SHAPE.OVAL)
        ntf = nb.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(ntf, num, size=19, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, x + Inches(0.16), y + Inches(1.02), cw - Inches(0.32),
                     Inches(1.2))
        write(tb.text_frame, title, size=21, bold=True, color=color, first=True,
              align=PP_ALIGN.CENTER, line=1.1)
        write(tb.text_frame, desc, size=13.5, color=INK_SOFT,
              align=PP_ALIGN.CENTER, space_before=5, line=1.25)

    # نوار تعهد + درخواست
    gap2 = Inches(0.36)
    hw = (CONTENT_W - gap2) / 2
    by = Inches(4.42)
    bh = Inches(1.42)

    rect(s, SW - M - hw, by, hw, bh, fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    rect(s, SW - M - Inches(0.09), by, Inches(0.09), bh, fill=ORANGE)
    tb = textbox(s, SW - M - hw + Inches(0.28), by + Inches(0.18),
                 hw - Inches(0.56), Inches(1.1))
    tf = tb.text_frame
    write(tf, "قولِ امروز", size=14, bold=True, color=ORANGE, first=True,
          line=1.0)
    write(tf, "یک‌سوم این جلسه، خودِ تجربه است — نه توضیحِ تجربه.",
          size=17.5, color=WHITE, space_before=8, line=1.28)

    rect(s, M, by, hw, bh, fill=ORANGE_PALE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09, line=ORANGE,
         line_w=Pt(1.75))
    tb = textbox(s, M + Inches(0.28), by + Inches(0.18), hw - Inches(0.56),
                 Inches(1.1))
    tf = tb.text_frame
    write(tf, "از شما یک چیز می‌خواهم", size=14, bold=True, color=ORANGE_DEEP,
          first=True, line=1.0)
    write(tf, "یک چالش کاری واقعی و حل‌نشده را همین حالا در ذهنتان نگه دارید.",
          size=17.5, bold=True, color=NAVY, space_before=8, line=1.28)

    band(s, Inches(6.06),
         [("در پایان جلسه، به همان چالش برمی‌گردیم.", False, WHITE, 20)],
         h=Inches(0.72), size=20)

    footer(s, 2, TOTAL, label=LABEL)
    notes(s, """
    روی بخش ۵ مکث کنید: مخاطب باید بداند این سخنرانی نیست.
    درخواست «یک چالش واقعی در ذهن نگه دارید» را آرام و جدی بگویید.
    """)
    return s


# ════════════════════════════════════ ۳ — بخش ۱: چرا ═════════════════════
def s03(prs, cfg):
    s = section(prs, 1, "چرا این کار را\nمی‌کنیم؟",
                "مسئله‌ای که همه می‌بینند و کسی درباره‌اش حرف نمی‌زند",
                "02-thinking.jpg", 0, side="right")
    notes(s, """
    این اسلاید را سریع رد نکنید. سه ثانیه سکوت، بعد بروید به اسلاید بعد.
    """)
    return s


# ═════════════════════════════════ ۴ — سؤالی که خودشان جواب می‌دهند ══════
def s04(prs, cfg):
    s = blank(prs)
    p, sz = img("06-silo.jpg")
    full_bleed(s, p, sz)
    gradient(s, 0, 0, SW, SH, NAVY_DEEP, NAVY, angle=315,
             alpha1=0.94, alpha2=0.86)

    kb = textbox(s, M, Inches(0.92), CONTENT_W, Inches(0.44))
    write(kb.text_frame, "۶۰ ثانیه فکر کنید", size=16, bold=True, color=ORANGE,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    qb = textbox(s, Inches(1.2), Inches(1.52), SW - Inches(2.4), Inches(1.9))
    tf = qb.text_frame
    write(tf, "آخرین چیزی که در کارتان یاد گرفتید،", size=34, bold=True,
          color=WHITE, first=True, align=PP_ALIGN.CENTER, line=1.3)
    write(tf, [("از کجا", True, ORANGE, 34), (" یاد گرفتید؟", True, WHITE, 34)],
          align=PP_ALIGN.CENTER, line=1.3)

    opts = [("از یک کلاس یا کتاب", False),
            ("از انجام دادن و اشتباه کردن", False),
            ("از یک آدم", True)]
    gap = Inches(0.32)
    cw = (CONTENT_W - 2 * gap) / 3
    y = Inches(3.72)
    ch = Inches(1.14)
    for i, (label, hi) in enumerate(opts):
        x = SW - M - cw - i * (cw + gap)
        rect(s, x, y, cw, ch, fill=WHITE, alpha=0.20 if hi else 0.09,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11,
             line=ORANGE if hi else BLUE_LIGHT, line_w=Pt(2 if hi else 1))
        tb = textbox(s, x + Inches(0.18), y, cw - Inches(0.36), ch,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, label, size=18.5, bold=hi,
              color=ORANGE if hi else WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.25)

    fb = textbox(s, Inches(1.2), Inches(5.30), SW - Inches(2.4), Inches(1.5))
    tf = fb.text_frame
    write(tf, "تقریباً همیشه، جواب یک آدم است.", size=25, bold=True,
          color=WHITE, first=True, align=PP_ALIGN.CENTER, line=1.25)
    write(tf, "نه یک دوره، نه یک کتاب. یک آدم.", size=17, color=BLUE_PALE,
          align=PP_ALIGN.CENTER, space_before=6, line=1.25)
    write(tf, "منتورینگ گروهی یعنی این اتفاق را از «شانسی» به «برنامه‌ریزی‌شده» تبدیل کنیم.",
          size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER,
          space_before=8, line=1.25)

    notes(s, """
    واقعاً ۶۰ ثانیه سکوت بدهید و از سه چهار نفر بپرسید.
    وقتی همه یک اسم گفتند: «توجه کردید؟ هیچ‌کس اسم یک دوره را نبرد.»
    این جواب، کل استدلال ارائه را می‌سازد.
    """)
    return s


# ══════════════════════════════ ۵ — مسئله: دانش در سیلو ══════════════════
def s05(prs, cfg):
    s = base_slide(prs, WHITE)
    head(s, "مسئله ما کمبود دانش نیست", part=0, kicker="تشخیص")

    pw = Inches(5.15)
    photo_card(s, M, Inches(1.66), pw, Inches(2.98), "03-silo.jpg",
               caption="هرکس تکه‌ای می‌داند، کسی همه را نمی‌داند")

    cap = card(s, M, Inches(4.80), pw, Inches(1.98), fill=NAVY, line=None)
    rect(s, M + pw - Inches(0.09), Inches(4.80), Inches(0.09), Inches(1.98),
         fill=ORANGE)
    tb = textbox(s, M + Inches(0.26), Inches(5.02), pw - Inches(0.55),
                 Inches(1.6))
    tf = tb.text_frame
    write(tf, "مسئله، گردش دانش است", size=14.5, bold=True, color=ORANGE,
          first=True, line=1.0)
    write(tf, "همه‌چیز را بلدیم — ولی هرکس تکه‌ای از آن را می‌داند و کسی نمی‌داند چه کسی چه می‌داند.",
          size=16.5, color=WHITE, space_before=9, line=1.3)

    tx = M + pw + Inches(0.46)
    tw = SW - M - tx
    symptoms = [
        ("یک اشتباه، چند بار", "همان خطا در واحد دیگر دوباره تکرار می‌شود"),
        ("راه‌حل‌های موازی", "دو نفر جدا از هم، ماه‌ها روی یک مسئله کار می‌کنند"),
        ("دانش با آدم‌ها می‌رود", "با هر خداحافظی، سال‌ها تجربه از در بیرون می‌رود"),
        ("تازه‌واردهای سرگردان", "شش ماه طول می‌کشد تا بفهمند از چه کسی بپرسند"),
    ]
    y = Inches(1.66)
    rh = Inches(1.20)
    for i, (title, desc) in enumerate(symptoms):
        rect(s, tx, y, tw, rh - Inches(0.10), fill=PAPER,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11)
        rect(s, tx + tw - Inches(0.08), y, Inches(0.08), rh - Inches(0.10),
             fill=RED)
        nb = rect(s, tx + tw - Inches(0.84), y + Inches(0.31), Inches(0.44),
                  Inches(0.44), fill=RED, shape=MSO_SHAPE.OVAL)
        ntf = nb.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(ntf, fa(i + 1), size=15, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, tx + Inches(0.28), y, tw - Inches(1.22),
                     rh - Inches(0.10), anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=19.5, bold=True, color=NAVY,
              first=True, line=1.05)
        write(tb.text_frame, desc, size=15, color=INK_SOFT, line=1.2)
        y += rh

    footer(s, 5, TOTAL, label=LABEL)
    notes(s, """
    اینجا حتماً یک مثال واقعی از سازمان خودتان بزنید، بدون بردن نام کسی.
    یک مثال واقعی از هر چهار مورد این فهرست مؤثرتر است.
    """)
    return s


# ═════════════════════════════ ۶ — علم: ۷۰-۲۰-۱۰ ═════════════════════════
def s06(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "بودجه آموزش، جای اشتباهی خرج می‌شود", part=0,
         kicker="آنچه پژوهش می‌گوید")

    cx = Inches(9.85)
    cy = Inches(3.68)
    R = Inches(1.60)
    R_IN = Inches(0.96)
    GAP = 2.2
    for a0, a1, color, pop in [(0.0, 252.0, GREEN, 0.0),
                               (252.0, 324.0, ORANGE, 0.14),
                               (324.0, 360.0, BLUE_MID, 0.0)]:
        T.donut_segment(s, cx, cy, R + Inches(pop), R_IN + Inches(pop * 0.5),
                        a0 + GAP, a1 - GAP, color)

    ctb = textbox(s, cx - Inches(1.0), cy - Inches(0.50), Inches(2.0),
                  Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
    write(ctb.text_frame, "۷۰-۲۰-۱۰", size=23, bold=True, color=NAVY,
          first=True, align=PP_ALIGN.CENTER, line=1.0)
    write(ctb.text_frame, "مدل یادگیری", size=12.5, color=INK_SOFT,
          align=PP_ALIGN.CENTER, line=1.1)

    src = textbox(s, cx - Inches(2.0), cy + R + Inches(0.40), Inches(4.0),
                  Inches(0.34))
    write(src.text_frame, "مرجع: مؤسسه CCL آمریکا", size=12.5, color=INK_SOFT,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    tx = M
    tw = Inches(7.45)
    rows = [("۷۰٪", "از انجام کار واقعی", GREEN, GREEN_PALE, False),
            ("۲۰٪", "از تعامل با آدم‌های دیگر", ORANGE, ORANGE_PALE, True),
            ("۱۰٪", "از کلاس و کتاب", BLUE_MID, BLUE_PALE, False)]
    y = Inches(1.68)
    for pct, label, color, pale, hi in rows:
        h = Inches(0.98) if hi else Inches(0.82)
        rect(s, tx, y, tw, h, fill=pale, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.12, line=color if hi else None, line_w=Pt(2))
        rect(s, tx + tw - Inches(0.09), y, Inches(0.09), h, fill=color)
        pb = textbox(s, tx + tw - Inches(1.60), y, Inches(1.40), h,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(pb.text_frame, pct, size=29 if hi else 25, bold=True, color=color,
              first=True, align=PP_ALIGN.CENTER, line=1.0)
        lb = textbox(s, tx + Inches(0.28), y, tw - Inches(2.02), h,
                     anchor=MSO_ANCHOR.MIDDLE)
        ch = [(label, hi, ORANGE_DEEP if hi else INK, 21.5 if hi else 19.5)]
        if hi:
            ch.append(("     ← اینجا", True, ORANGE, 15.5))
        write(lb.text_frame, ch, first=True, line=1.1)
        y += h + Inches(0.14)

    by = Inches(4.86)
    rect(s, tx, by, tw, Inches(1.72), fill=WHITE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09, line=LINE)
    tb = textbox(s, tx + Inches(0.28), by + Inches(0.16), tw - Inches(0.56),
                 Inches(0.38))
    write(tb.text_frame, "و بودجه آموزشی ما کجا خرج می‌شود؟", size=15.5,
          bold=True, color=NAVY, first=True, line=1.0)

    bar_y = by + Inches(0.64)
    bar_w = tw - Inches(0.56)
    bar_h = Inches(0.40)
    rect(s, tx + Inches(0.28), bar_y, bar_w, bar_h, fill=BLUE_PALE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.30)
    rect(s, tx + Inches(0.28) + bar_w * 0.10, bar_y, bar_w * 0.90, bar_h,
         fill=BLUE_MID, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.30)
    lb = textbox(s, tx + Inches(0.28) + bar_w * 0.10, bar_y, bar_w * 0.90,
                 bar_h, anchor=MSO_ANCHOR.MIDDLE)
    write(lb.text_frame, "عمدتاً روی کلاس و دوره", size=14, bold=True, color=WHITE,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    nb = textbox(s, tx + Inches(0.28), by + Inches(1.16), bar_w, Inches(0.42))
    write(nb.text_frame,
          [("یعنی ", False, INK_SOFT, 15.5),
           ("بخش عمده بودجه", True, RED, 15.5),
           ("، روی ", False, INK_SOFT, 15.5), ("۱۰٪", True, RED, 15.5),
           (" یادگیری. منتورینگ گروهی سراغ آن ", False, INK_SOFT, 15.5),
           ("۲۰٪ رهاشده", True, ORANGE_DEEP, 15.5), (" می‌رود.", False, INK_SOFT, 15.5)],
          first=True, line=1.15)

    footer(s, 6, TOTAL, label=LABEL)
    notes(s, """
    این اسلاید پشتوانه علمی طرح است، مخصوصاً برای مدیران.
    اگر عدد واقعی بودجه سازمان را دارید، جایگزین کنید.
    """)
    return s


# ═══════════════════════════════════ ۷ — بخش ۲: چیست ═════════════════════
def s07(prs, cfg):
    s = section(prs, 2, "این روش\nچیست؟",
                "تعریف، مرزها، و نقش‌هایی که در اتاق وجود دارد",
                "05-circle.jpg", 1, side="right")
    notes(s, "از اینجا لحن عوض می‌شود: از مسئله به راه‌حل.")
    return s


# ════════════════════════════ ۸ — تعریف در یک جمله ═══════════════════════
def s08(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "منتورینگ گروهی، در یک جمله", part=1, kicker="تعریف")

    card(s, M, Inches(1.64), CONTENT_W, Inches(1.56), fill=WHITE, line=LINE)
    rect(s, SW - M - Inches(0.10), Inches(1.64), Inches(0.10), Inches(1.56),
         fill=ORANGE)
    tb = textbox(s, M + Inches(0.36), Inches(1.64), CONTENT_W - Inches(0.72),
                 Inches(1.56), anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame,
          [("گروهی کوچک از همکاران که ", False, INK, 25),
           ("مرتب", True, ORANGE_DEEP, 25),
           (" دور هم جمع می‌شوند تا روی ", False, INK, 25),
           ("چالش‌های واقعی کاری", True, NAVY, 25),
           (" یکدیگر فکر کنند — با ", False, INK, 25),
           ("پرسش", True, ORANGE_DEEP, 25), ("، نه با ", False, INK, 25),
           ("نسخه‌پیچی", True, INK_SOFT, 25), (".", False, INK, 25)],
          first=True, align=PP_ALIGN.CENTER, line=1.4)

    specs = [("۸ تا ۱۲", "نفر در هر گروه", NAVY),
             ("۹۰", "دقیقه در هر جلسه", BLUE_MID),
             ("هر ۲ هفته", "یک بار", BLUE_LIGHT),
             ("۶ ماه", "طول یک دوره", ORANGE)]
    gap = Inches(0.30)
    cw = (CONTENT_W - 3 * gap) / 4
    y = Inches(3.50)
    ch = Inches(1.28)
    for i, (big, small, color) in enumerate(specs):
        x = SW - M - cw - i * (cw + gap)
        rect(s, x, y, cw, ch, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.11, line=LINE)
        rect(s, x, y, cw, Inches(0.07), fill=color)
        tb = textbox(s, x + Inches(0.14), y, cw - Inches(0.28), ch,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, big, size=26, bold=True, color=color, first=True,
              align=PP_ALIGN.CENTER, line=1.05)
        write(tb.text_frame, small, size=14.5, color=INK_SOFT,
              align=PP_ALIGN.CENTER, line=1.15)

    dy = Inches(5.06)
    gap2 = Inches(0.34)
    hw = (CONTENT_W - gap2) / 2
    pairs = [(SW - M - hw, "منتورینگ کلاسیک", "یک نفر می‌داند، یک نفر یاد می‌گیرد",
              "۱ ← ۱", INK_SOFT, T.RGBColor(0xEC, 0xF0, 0xF3)),
             (M, "منتورینگ گروهی", "همه می‌دانند، همه یاد می‌گیرند",
              "همه ⇄ همه", ORANGE_DEEP, ORANGE_PALE)]
    for x, title, desc, badge, color, pale in pairs:
        rect(s, x, dy, hw, Inches(1.24), fill=pale,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11)
        bd = rect(s, x + hw - Inches(1.62), dy + Inches(0.29), Inches(1.40),
                  Inches(0.64), fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                  radius=0.26)
        btf = bd.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(btf, badge, size=14.5, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, x + Inches(0.28), dy, hw - Inches(2.02), Inches(1.24),
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=18.5, bold=True, color=color,
              first=True, line=1.05)
        write(tb.text_frame, desc, size=14.5, color=INK, line=1.2)

    footer(s, 8, TOTAL, label=LABEL)
    notes(s, """
    «پرسش، نه نسخه‌پیچی» را دو بار تکرار کنید — کل تفاوت روش همین است.
    """)
    return s


# ══════════════════════════════ ۹ — هست / نیست ═══════════════════════════
def s09(prs, cfg):
    s = base_slide(prs, WHITE)
    head(s, "چه چیزی هست، چه چیزی نیست", part=1,
         kicker="مرزبندی — تا انتظار اشتباه شکل نگیرد")

    gap = Inches(0.40)
    cw = (CONTENT_W - gap) / 2
    top = Inches(1.64)
    ch = Inches(0.60)
    rowh = Inches(0.74)
    boxh = ch + 5 * rowh + Inches(0.20)

    cols = [
        (M + cw + gap, "هست", GREEN, GREEN_PALE, "✓", [
            ("فضای فکر کردن با هم", "روی مسئله‌ای که واقعاً درگیرش هستید"),
            ("یادگیری از تجربه همکار", "کسی که مسیر شما را رفته است"),
            ("تمرین پرسیدن", "مهارتی که در جلسات کاری هم به کار می‌آید"),
            ("تعهد به یک اقدام کوچک", "هر جلسه، یک قدم مشخص"),
            ("محرمانه و داوطلبانه", "هیچ‌کس اجباری نیست"),
        ]),
        (M, "نیست", RED, RED_PALE, "✕", [
            ("کلاس آموزشی", "کسی قرار نیست تدریس کند"),
            ("جلسه گزارش‌دهی", "هیچ صورت‌جلسه‌ای به بالا نمی‌رود"),
            ("ارزیابی عملکرد", "روی پرونده پرسنلی شما اثری ندارد"),
            ("جلسه شکایت", "درباره مسئله حرف می‌زنیم، نه درباره آدم‌ها"),
            ("جای نصیحت کردن", "قدیمی‌ترها هم اینجا فقط می‌پرسند"),
        ]),
    ]
    for x, title, cmain, cpale, mark, items in cols:
        card(s, x, top, cw, boxh, fill=WHITE, line=LINE)
        rect(s, x, top, cw, ch, fill=cmain, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.12)
        rect(s, x, top + ch - Inches(0.16), cw, Inches(0.16), fill=cmain)
        tb = textbox(s, x, top, cw, ch, anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, [(mark + "   ", True, WHITE, 18),
                              (title, True, WHITE, 22)],
              first=True, align=PP_ALIGN.CENTER, line=1.0)
        for i, (main, sub) in enumerate(items):
            y = top + ch + Inches(0.10) + i * rowh
            if i % 2 == 0:
                rect(s, x + Inches(0.13), y, cw - Inches(0.26),
                     rowh - Inches(0.06), fill=cpale,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
            tb = textbox(s, x + Inches(0.30), y, cw - Inches(0.60),
                         rowh - Inches(0.06), anchor=MSO_ANCHOR.MIDDLE)
            write(tb.text_frame, [(mark + "  ", True, cmain, 13.5),
                                  (main, True, INK, 18)], first=True, line=1.05)
            write(tb.text_frame, sub, size=13, color=INK_SOFT, line=1.15)

    footer(s, 9, TOTAL, label=LABEL)
    notes(s, """
    ستون «نیست» مهم‌تر از «هست» است. بیشترین مقاومت از همین سوءتفاهم‌ها می‌آید.
    روی «هیچ صورت‌جلسه‌ای به بالا نمی‌رود» مکث کنید و در چشم افراد نگاه کنید.
    """)
    return s


# ═══════════════════════════════ ۱۰ — سه نقش در اتاق ═════════════════════
def s10(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "در اتاق، سه نقش وجود دارد", part=1, kicker="نقش‌ها")

    roles = [
        ("۱", "صاحب چالش", "یک نفر، داوطلب",
         ["مسئله واقعی‌اش را روایت می‌کند",
          "فقط گوش می‌دهد و یادداشت برمی‌دارد",
          "در پایان تصمیم می‌گیرد چه کند"], ORANGE),
        ("۲", "گروه همتا", "۷ تا ۱۱ نفر",
         ["سؤال می‌پرسند، راه‌حل نمی‌دهند",
          "تجربه مشابه خود را روایت می‌کنند",
          "قضاوت نمی‌کنند"], NAVY),
        ("۳", "تسهیل‌گر", "چرخشی، هر جلسه یک نفر",
         ["زمان را نگه می‌دارد",
          "قواعد را یادآوری می‌کند",
          "خودش کمترین حرف را می‌زند"], GREEN),
    ]
    gap = Inches(0.32)
    cw = (CONTENT_W - 2 * gap) / 3
    top = Inches(1.64)
    ch = Inches(3.48)
    for i, (num, title, sub, points, color) in enumerate(roles):
        x = SW - M - cw - i * (cw + gap)
        card(s, x, top, cw, ch, fill=WHITE, line=LINE)
        rect(s, x, top, cw, Inches(1.04), fill=color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
        rect(s, x, top + Inches(0.88), cw, Inches(0.16), fill=color)
        bd = rect(s, x + cw - Inches(0.70), top + Inches(0.29), Inches(0.42),
                  Inches(0.42), fill=WHITE, shape=MSO_SHAPE.OVAL)
        btf = bd.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(btf, num, size=16, bold=True, color=color, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, x + Inches(0.22), top + Inches(0.12),
                     cw - Inches(1.02), Inches(0.80), anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=20.5, bold=True, color=WHITE,
              first=True, line=1.05)
        write(tb.text_frame, sub, size=13, color=WHITE, line=1.15)
        y = top + Inches(1.22)
        for pt_txt in points:
            bx = textbox(s, x + Inches(0.24), y, cw - Inches(0.48),
                         Inches(0.70))
            write(bx.text_frame, [("◂  ", True, color, 12.5),
                                  (pt_txt, False, INK, 15.5)],
                  first=True, line=1.25)
            y += Inches(0.72)

    band(s, Inches(5.42),
         [("نکته‌ای که همه‌چیز را عوض می‌کند:   ", True, ORANGE, 16.5),
          ("هفته بعد، صاحب چالش یک نفر دیگر است. امروز کمک می‌گیرید، دفعه بعد کمک می‌کنید.",
           False, WHITE, 19.5)],
         h=Inches(1.02), align=PP_ALIGN.RIGHT, size=19.5)

    footer(s, 10, TOTAL, label=LABEL)
    notes(s, """
    تسهیل‌گری چرخشی است — این طرح مالِ یک نفر نیست، مالِ گروه است.
    اگر پرسیدند «مدیر هم هست؟» بگویید بله، ولی به‌عنوان عضو گروه همتا.
    """)
    return s


# ═══════════════════════════════════ ۱۱ — بخش ۳: چطور ════════════════════
def s11(prs, cfg):
    s = section(prs, 3, "چطور اجرا\nمی‌شود؟",
                "یک جلسه واقعی، دقیقه به دقیقه — و مهارتی که باید یاد بگیریم",
                "01-kickoff.jpg", 2, side="right")
    notes(s, "از اینجا وارد جزئیات عملی می‌شویم.")
    return s


# ══════════════════════════ ۱۲ — آناتومی جلسه ۹۰ دقیقه‌ای ════════════════
def s12(prs, cfg):
    s = base_slide(prs, WHITE)
    head(s, "یک جلسه واقعی چطور پیش می‌رود؟", part=2,
         kicker="آناتومی ۹۰ دقیقه")

    steps = [
        (10, "دایره باز", "هر نفر یک جمله: این هفته چطور بود", BLUE_LIGHT),
        (10, "مرور تعهدها", "جلسه قبل چه قول دادیم و چه شد", BLUE_LIGHT),
        (10, "روایت چالش", "داوطلب مسئله‌اش را می‌گوید، بدون قطع شدن", BLUE_MID),
        (20, "فقط سؤال", "گروه فقط می‌پرسد — هیچ راه‌حلی ممنوع", ORANGE),
        (20, "تجربه‌های مشابه", "«من هم یک بار…» — روایت، نه دستور", ORANGE),
        (15, "انتخاب و تعهد", "صاحب چالش می‌گوید چه می‌کند", NAVY),
        (5, "دایره بسته", "هر نفر یک کلمه: چه چیزی با خود می‌برد", NAVY),
    ]
    bar_y = Inches(1.64)
    bar_h = Inches(0.38)
    total_min = sum(x[0] for x in steps)
    xcur = SW - M
    for mins, title, desc, color in steps:
        w = CONTENT_W * mins / total_min
        rect(s, xcur - w, bar_y, w - Inches(0.035), bar_h, fill=color)
        tb = textbox(s, xcur - w, bar_y, w, bar_h, anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, fa(mins), size=13, bold=True, color=WHITE,
              first=True, align=PP_ALIGN.CENTER, line=1.0)
        xcur -= w

    y = Inches(2.22)
    rh = Inches(0.535)
    for i, (mins, title, desc, color) in enumerate(steps):
        if i % 2 == 0:
            rect(s, M, y, CONTENT_W, rh - Inches(0.04), fill=PAPER,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.20)
        chip = rect(s, SW - M - Inches(1.24), y + Inches(0.05), Inches(1.08),
                    rh - Inches(0.14), fill=color,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.26)
        ctf = chip.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(ctf, fa("%d′" % mins), size=13.5, bold=True, color=WHITE,
              first=True, align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, SW - M - Inches(4.95), y, Inches(3.5),
                     rh - Inches(0.04), anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=19.5, bold=True, color=NAVY,
              first=True, line=1.0)
        db = textbox(s, M + Inches(0.2), y, Inches(7.6), rh - Inches(0.04),
                     anchor=MSO_ANCHOR.MIDDLE)
        write(db.text_frame, desc, size=16.5, color=INK_SOFT, first=True,
              line=1.0)
        y += rh

    gy = Inches(6.06)
    rect(s, M, gy, CONTENT_W, Inches(0.66), fill=ORANGE_PALE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.13, line=ORANGE,
         line_w=Pt(1.75))
    tb = textbox(s, M, gy, CONTENT_W, Inches(0.66), anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame,
          [("۴۰ دقیقه از ۹۰ دقیقه — یعنی نیمی از جلسه — ", False, INK, 17.5),
           ("فقط سؤال و روایت است", True, ORANGE_DEEP, 18.5),
           ("، نه راه‌حل.", False, INK, 17.5)],
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    footer(s, 12, TOTAL, label=LABEL)
    notes(s, """
    بخش «فقط سؤال» بیشترین مقاومت را دارد.
    بگویید: «می‌دانم سخت است. اسلاید بعد نشان می‌دهم چرا ارزشش را دارد.»
    """)
    return s


# ═════════════════════ ۱۳ — قلب روش: نصیحت ← پرسش ════════════════════════
def s13(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "قلب روش: به‌جای نصیحت، بپرس", part=2, kicker="مهارت کلیدی")

    pw = Inches(3.70)
    photo_card(s, M, Inches(1.64), pw, Inches(2.52), "12-listen.jpg")

    nb = card(s, M, Inches(4.32), pw, Inches(2.46), fill=NAVY, line=None)
    rect(s, M + pw - Inches(0.09), Inches(4.32), Inches(0.09), Inches(2.46),
         fill=ORANGE)
    tb = textbox(s, M + Inches(0.24), Inches(4.54), pw - Inches(0.52),
                 Inches(2.05))
    tf = tb.text_frame
    write(tf, "چرا نصیحت جواب نمی‌دهد؟", size=14.5, bold=True, color=ORANGE,
          first=True, line=1.0)
    write(tf, "چون شما تمام جزئیات مسئله او را نمی‌دانید.", size=16,
          color=WHITE, space_before=9, line=1.28)
    write(tf, "ولی یک سؤال خوب، خودِ او را به جوابی می‌رساند که خودش اجرایش می‌کند.",
          size=16, color=BLUE_PALE, space_before=7, line=1.28)

    tx = M + pw + Inches(0.44)
    tw = SW - M - tx
    gap = Inches(0.30)
    hw = (tw - gap) / 2
    hy = Inches(1.64)
    for x, label, color in [(tx + hw + gap, "به‌جای این…", RED),
                            (tx, "این را بگویید", GREEN)]:
        rect(s, x, hy, hw, Inches(0.52), fill=color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
        tb = textbox(s, x, hy, hw, Inches(0.52), anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, label, size=16.5, bold=True, color=WHITE,
              first=True, align=PP_ALIGN.CENTER, line=1.0)

    pairs = [
        ("«جای تو بودم فلان کار را می‌کردم.»", "«خودت چه گزینه‌هایی را سنجیدی؟»"),
        ("«این روش اشتباه است.»", "«اگر جواب ندهد، بدترین اتفاق چیست؟»"),
        ("«ما همیشه اینجوری کردیم.»", "«چه چیزی این مسئله را برایت سخت کرده؟»"),
        ("«نگران نباش، درست می‌شود.»", "«موفقیت در این مسئله برایت یعنی چه؟»"),
    ]
    y = Inches(2.32)
    rh = Inches(1.14)
    for bad, good in pairs:
        for x, txt, pale, muted in [(tx + hw + gap, bad, RED_PALE, True),
                                    (tx, good, GREEN_PALE, False)]:
            rect(s, x, y, hw, rh - Inches(0.14), fill=pale,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.13)
            tb = textbox(s, x + Inches(0.18), y, hw - Inches(0.36),
                         rh - Inches(0.14), anchor=MSO_ANCHOR.MIDDLE)
            write(tb.text_frame, txt, size=16, bold=not muted,
                  color=INK_SOFT if muted else INK, first=True,
                  align=PP_ALIGN.CENTER, line=1.25)
        y += rh

    footer(s, 13, TOTAL, label=LABEL)
    notes(s, """
    مهم‌ترین اسلاید ارائه. آرام بخوانید، بعد از هر جفت مکث کوتاه کنید.
    از گروه بخواهید یک نمونه دیگر بسازند — اگر توانستند، روش را گرفته‌اند.
    """)
    return s


# ═════════════════════════════ ۱۴ — قواعد فضای امن ═══════════════════════
def s14(prs, cfg):
    s = base_slide(prs, WHITE)
    head(s, "پنج قاعده‌ای که فضا را امن می‌کند", part=2, kicker="پیمان گروه")

    rules = [
        ("۱", "اینجا سِمَت نداریم", "در این ۹۰ دقیقه هیچ‌کس رئیس کسی نیست",
         ORANGE, ORANGE_PALE),
        ("۲", "هرچه گفته شد، همین‌جا می‌ماند", "بدون استثنا، بدون صورت‌جلسه",
         NAVY, BLUE_PALE),
        ("۳", "سؤال بپرس، نسخه نپیچ", "کنجکاوی به‌جای نصیحت", GREEN, GREEN_PALE),
        ("۴", "هیچ سؤالی احمقانه نیست", "و هیچ تجربه‌ای بی‌ارزش نیست", PURPLE,
         T.RGBColor(0xF6, 0xEF, 0xFA)),
        ("۵", "گوش بده تا بفهمی", "نه اینکه منتظر نوبت حرف زدنت باشی", RED,
         RED_PALE),
    ]
    tw = Inches(8.05)
    tx = SW - M - tw
    y = Inches(1.64)
    rh = Inches(0.90)
    for num, title, desc, color, pale in rules:
        rect(s, tx, y, tw, rh, fill=pale, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.11)
        rect(s, tx + tw - Inches(0.10), y, Inches(0.10), rh, fill=color)
        bd = rect(s, tx + tw - Inches(0.92), y + Inches(0.22), Inches(0.46),
                  Inches(0.46), fill=color, shape=MSO_SHAPE.OVAL)
        btf = bd.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(btf, num, size=16.5, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, tx + Inches(0.28), y, tw - Inches(1.30), rh,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=21.5, bold=True, color=color,
              first=True, line=1.05)
        write(tb.text_frame, desc, size=15.5, color=INK_SOFT, line=1.15)
        y += rh + Inches(0.09)

    px = M
    pw = tx - M - Inches(0.40)
    photo_card(s, px, Inches(1.64), pw, Inches(3.06), "13-safe.jpg")

    card(s, px, Inches(4.86), pw, Inches(1.92), fill=NAVY, line=None)
    rect(s, px + pw - Inches(0.09), Inches(4.86), Inches(0.09), Inches(1.92),
         fill=ORANGE)
    tb = textbox(s, px + Inches(0.24), Inches(5.06), pw - Inches(0.52),
                 Inches(1.55))
    tf = tb.text_frame
    write(tf, "چرا این پنج تا؟", size=14, bold=True, color=ORANGE, first=True,
          line=1.0)
    write(tf, "تا وقتی کسی نترسد از گفتنِ «بلد نیستم»، هیچ یادگیری‌ای شروع نمی‌شود.",
          size=15.5, color=WHITE, space_before=8, line=1.26)

    footer(s, 14, TOTAL, label=LABEL)
    notes(s, """
    این قواعد را نخوانید — از گروه بخواهید بلند بخوانند و تأیید کنند.
    روی یک برگه بزرگ بنویسید و هر جلسه روی دیوار بزنید.
    """)
    return s


# ══════════════════════════════ ۱۵ — چرخه شش‌ماهه ════════════════════════
def s15(prs, cfg):
    s = base_slide(prs, PAPER)
    head(s, "شش ماه، دوازده جلسه، یک گروه", part=2, kicker="مسیر پیشِ رو")

    phases = [
        ("جلسه ۱", "شکل‌گیری", "آشنایی، بستن پیمان گروه، انتخاب اولین داوطلب",
         ORANGE),
        ("جلسه ۲ تا ۴", "تمرین", "قواعد هنوز سخت است؛ تسهیل‌گر بیشتر دخالت می‌کند",
         BLUE_MID),
        ("جلسه ۵ تا ۱۰", "اوج‌گیری", "گروه روان می‌شود؛ تسهیل‌گری چرخشی آغاز می‌شود",
         NAVY),
        ("جلسه ۱۱ و ۱۲", "جمع‌بندی", "مرور دستاوردها و تصمیم برای دوره بعد",
         GREEN),
    ]
    tw = Inches(7.85)
    tx = SW - M - tw
    y = Inches(1.64)
    rh = Inches(1.04)
    for i, (when, title, desc, color) in enumerate(phases):
        rect(s, tx, y, tw, rh, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.11, line=LINE)
        rect(s, tx + tw - Inches(0.09), y, Inches(0.09), rh, fill=color)
        wb = rect(s, tx + tw - Inches(1.92), y + Inches(0.26), Inches(1.64),
                  rh - Inches(0.52), fill=color,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.26)
        wtf = wb.text_frame
        wtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(wtf, when, size=13.5, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, tx + Inches(0.28), y, tw - Inches(2.30), rh,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=19.5, bold=True, color=color,
              first=True, line=1.05)
        write(tb.text_frame, desc, size=14.5, color=INK_SOFT, line=1.2)
        if i < len(phases) - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                    tx + tw - Inches(1.18),
                                    y + rh + Inches(0.01),
                                    Inches(0.22), Inches(0.15))
            ar.fill.solid(); ar.fill.fore_color.rgb = LINE
            no_line(ar); shadow_off(ar)
        y += rh + Inches(0.17)

    px = M
    pw = tx - M - Inches(0.40)
    photo_card(s, px, Inches(1.64), pw, Inches(2.05), "04-climb.jpg")

    card(s, px, Inches(3.82), pw, Inches(1.60), fill=WHITE, line=LINE)
    rect(s, px, Inches(3.82), pw, Inches(0.54), fill=NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
    rect(s, px, Inches(4.22), pw, Inches(0.14), fill=NAVY)
    tb = textbox(s, px, Inches(3.82), pw, Inches(0.54), anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame, "تعهد شما", size=17, bold=True, color=WHITE,
          first=True, align=PP_ALIGN.CENTER, line=1.0)
    asks = ["حضور منظم در ۱۲ جلسه", "آوردن یک چالش واقعی",
            "رعایت محرمانگی", "یک اقدام کوچک بین دو جلسه"]
    yy = Inches(4.46)
    for a in asks:
        bx = textbox(s, px + Inches(0.20), yy, pw - Inches(0.40), Inches(0.24))
        write(bx.text_frame, [("◂  ", True, ORANGE, 11.5),
                              (a, False, INK, 13.5)], first=True, line=1.0)
        yy += Inches(0.25)

    card(s, px, Inches(5.58), pw, Inches(1.20), fill=ORANGE_PALE, line=ORANGE)
    tb = textbox(s, px + Inches(0.22), Inches(5.58), pw - Inches(0.44),
                 Inches(1.20), anchor=MSO_ANCHOR.MIDDLE)
    tf = tb.text_frame
    write(tf, "و اگر نشد؟", size=13.5, bold=True, color=ORANGE_DEEP, first=True,
          line=1.0)
    write(tf, "بعد از جلسه سوم، هر کس خواست می‌تواند کنار بکشد — بدون توضیح.",
          size=14.5, color=NAVY, space_before=6, line=1.24)

    footer(s, 15, TOTAL, label=LABEL)
    notes(s, """
    جمله «بدون توضیح می‌توانید کنار بکشید» را حتماً بگویید.
    داوطلبانه بودن شرط کارکردن این روش است، نه یک تعارف.
    """)
    return s


# ════════════════════════════ ۱۶ — بخش ۴: چه می‌شود ══════════════════════
def s16(prs, cfg):
    s = section(prs, 4, "چه چیزی\nعوض می‌شود؟",
                "برای شما، برای تیم، و برای سازمان",
                "08-growth.jpg", 3, side="right")
    notes(s, "لحن اینجا امیدوارانه می‌شود. سرعت را کمی بالا ببرید.")
    return s


# ═══════════════════════════════ ۱۷ — دستاوردها ══════════════════════════
def s17(prs, cfg):
    s = base_slide(prs, WHITE)
    head(s, "چه چیزی عاید چه کسی می‌شود؟", part=3, kicker="دستاورد")

    gap = Inches(0.38)
    cw = (CONTENT_W - gap) / 2
    top = Inches(1.64)
    ch = Inches(0.62)
    rowh = Inches(0.68)
    boxh = ch + 4 * rowh + Inches(0.20)

    cols = [
        (M + cw + gap, "برای شما", NAVY, BLUE_PALE, [
            ("مسئله‌تان زودتر حل می‌شود", "به‌جای هفته‌ها کلنجار، ۹۰ دقیقه با ۱۰ ذهن"),
            ("مهارت پرسیدن یاد می‌گیرید", "در جلسات و مذاکرات کاری هم به کارتان می‌آید"),
            ("شبکه‌تان گسترده می‌شود", "می‌دانید برای هر مسئله سراغ چه کسی بروید"),
            ("کمتر احساس تنهایی می‌کنید", "می‌فهمید دیگران هم همین را تجربه کرده‌اند"),
        ]),
        (M, "برای سازمان", ORANGE_DEEP, ORANGE_PALE, [
            ("تکرار اشتباهات کم می‌شود", "تجربه شکست یک واحد، درسِ واحد دیگر می‌شود"),
            ("دانش ماندگار می‌شود", "تجربه از سرِ افراد به حافظه جمعی منتقل می‌شود"),
            ("تازه‌واردها زودتر جا می‌افتند", "مسیر شش‌ماهه، چندهفته‌ای طی می‌شود"),
            ("افراد می‌مانند", "آدم‌ها جایی که رشد می‌کنند را ترک نمی‌کنند"),
        ]),
    ]
    for x, title, cmain, cpale, items in cols:
        card(s, x, top, cw, boxh, fill=WHITE, line=LINE)
        rect(s, x, top, cw, ch, fill=cmain, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             radius=0.11)
        rect(s, x, top + ch - Inches(0.16), cw, Inches(0.16), fill=cmain)
        tb = textbox(s, x, top, cw, ch, anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, title, size=21.5, bold=True, color=WHITE,
              first=True, align=PP_ALIGN.CENTER, line=1.0)
        for i, (main, sub) in enumerate(items):
            y = top + ch + Inches(0.10) + i * rowh
            if i % 2 == 0:
                rect(s, x + Inches(0.13), y, cw - Inches(0.26),
                     rowh - Inches(0.05), fill=cpale,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14)
            tb = textbox(s, x + Inches(0.30), y, cw - Inches(0.60),
                         rowh - Inches(0.05), anchor=MSO_ANCHOR.MIDDLE)
            write(tb.text_frame, [("✓  ", True, cmain, 13.5),
                                  (main, True, INK, 17.5)], first=True,
                  line=1.05)
            write(tb.text_frame, sub, size=13, color=INK_SOFT, line=1.15)

    band(s, Inches(5.20),
         [("هزینه‌اش؟  ", True, ORANGE, 18),
          ("۹۰ دقیقه، هر دو هفته یک بار — یعنی حدود ", False, WHITE, 20),
          ("یک درصد", True, ORANGE, 20), (" وقت کاری شما.", False, WHITE, 20)],
         h=Inches(0.92), size=20)

    # نوار تصویری کوچک: سه تصویر به‌عنوان جمع‌بندی بصری
    sy = Inches(6.24)
    sh_ = Inches(0.62)
    strip = [("11-commit.jpg", "تعهد فردی"), ("10-safe.jpg", "فضای امن"),
             ("05-circle.jpg", "گفتگوی گروهی")]
    sgap = Inches(0.20)
    sw_ = (CONTENT_W - 2 * sgap) / 3
    for i, (fname, cap) in enumerate(strip):
        x = SW - M - sw_ - i * (sw_ + sgap)
        p, sz = img(fname)
        pic = picture_fill(s, p, x, sy, sw_, sh_, sz)
        round_picture(pic)
        scrim(s, x, sy, sw_, sh_, alpha=0.55)
        tb = textbox(s, x, sy, sw_, sh_, anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, cap, size=13, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)

    footer(s, 17, TOTAL, label=LABEL)
    notes(s, """
    اگر مخاطب بیشتر کارشناس است، ستون راست را کامل بخوانید و چپ را سریع رد شوید.
    اگر مدیران حاضرند، برعکس.
    """)
    return s


# ═══════════════════════════════════ ۱۸ — بخش ۵: شروع ════════════════════
def s18(prs, cfg):
    s = section(prs, 5, "حالا خودمان\nانجامش می‌دهیم",
                "یک دور واقعی، همین امروز، ۳۰ دقیقه",
                "15-start.jpg", 4, side="right")
    notes(s, "اینجا انرژی جلسه باید بالا برود. بلند شوید و راه بروید.")
    return s


# ══════════════════════════════ ۱۹ — تجربه زنده ══════════════════════════
def s19(prs, cfg):
    s = blank(prs)
    gradient(s, 0, 0, SW, SH, NAVY_DEEP, NAVY, angle=315)

    kb = textbox(s, M, Inches(0.52), CONTENT_W, Inches(0.38))
    write(kb.text_frame, "تجربه زنده", size=15.5, bold=True, color=ORANGE,
          first=True, align=PP_ALIGN.CENTER, line=1.0)
    tb = textbox(s, M, Inches(0.94), CONTENT_W, Inches(0.66))
    write(tb.text_frame, "یک دور کوتاه — ۳۰ دقیقه", size=33, bold=True,
          color=WHITE, first=True, align=PP_ALIGN.CENTER, line=1.0)

    steps = [("۵ دقیقه", "یک داوطلب", "چالشش را در سه جمله می‌گوید"),
             ("۱۰ دقیقه", "فقط سؤال", "هر کس یک سؤال — راه‌حل ممنوع"),
             ("۱۰ دقیقه", "تجربه مشابه", "«من هم یک بار…» در یک دقیقه"),
             ("۵ دقیقه", "یک تصمیم", "داوطلب می‌گوید چه می‌کند")]
    gap = Inches(0.28)
    cw = (CONTENT_W - 3 * gap) / 4
    y = Inches(1.94)
    ch = Inches(2.00)
    for i, (mins, title, desc) in enumerate(steps):
        x = SW - M - cw - i * (cw + gap)
        rect(s, x, y, cw, ch, fill=WHITE, alpha=0.10,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09,
             line=BLUE_LIGHT, line_w=Pt(1))
        chip = rect(s, x + (cw - Inches(1.30)) / 2, y + Inches(0.22),
                    Inches(1.30), Inches(0.48), fill=ORANGE,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.26)
        ctf = chip.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(ctf, mins, size=14, bold=True, color=WHITE, first=True,
              align=PP_ALIGN.CENTER, line=1.0)
        tb = textbox(s, x + Inches(0.16), y + Inches(0.84), cw - Inches(0.32),
                     Inches(1.0))
        write(tb.text_frame, title, size=18.5, bold=True, color=WHITE,
              first=True, align=PP_ALIGN.CENTER, line=1.1)
        write(tb.text_frame, desc, size=13.5, color=BLUE_PALE,
              align=PP_ALIGN.CENTER, line=1.2)

    ry = Inches(4.28)
    rect(s, M, ry, CONTENT_W, Inches(0.84), fill=ORANGE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11)
    tb = textbox(s, M, ry, CONTENT_W, Inches(0.84), anchor=MSO_ANCHOR.MIDDLE)
    write(tb.text_frame,
          [("تنها قاعده امروز:  ", True, WHITE, 17.5),
           ("در ده دقیقه دوم، هیچ‌کس حق ندارد جمله‌ای بگوید که علامت سؤال ندارد.",
            True, WHITE, 20)],
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    rect(s, M, Inches(5.38), CONTENT_W, Inches(1.42), fill=WHITE, alpha=0.10,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09, line=BLUE_LIGHT)
    tb = textbox(s, M + Inches(0.40), Inches(5.38), CONTENT_W - Inches(0.80),
                 Inches(1.42), anchor=MSO_ANCHOR.MIDDLE)
    tf = tb.text_frame
    write(tf, "چه کسی چالشش را می‌آورد؟", size=23, bold=True, color=ORANGE,
          first=True, align=PP_ALIGN.CENTER, line=1.15)
    write(tf, "لازم نیست مسئله بزرگی باشد — همان چیزی که اول جلسه در ذهن نگه داشتید کافی است.",
          size=16, color=BLUE_PALE, align=PP_ALIGN.CENTER, line=1.25)

    notes(s, """
    مهم‌ترین ۳۰ دقیقه امروز. تایمر بگذارید و سفت نگه دارید.
    اگر کسی در «فقط سؤال» راه‌حل داد، با لبخند وسط حرفش بروید:
    «این را نگه دار برای ده دقیقه بعد.»
    اگر داوطلب پیدا نشد، خودتان یک چالش واقعی بیاورید.
    در پایان از صاحب چالش بپرسید: «الان چیزی می‌دانی که نیم‌ساعت پیش نمی‌دانستی؟»
    """)
    return s


# ══════════════════════════════ ۲۰ — دعوت پایانی ═════════════════════════
def s20(prs, cfg):
    s = blank(prs)
    p, sz = img("12-road.jpg")
    full_bleed(s, p, sz)
    gradient(s, 0, 0, SW, SH, NAVY_DEEP, ORANGE_DEEP, angle=45,
             alpha1=0.88, alpha2=0.34)
    scrim(s, 0, 0, SW, SH, alpha=0.26)

    kb = textbox(s, M, Inches(0.56), CONTENT_W, Inches(0.4))
    write(kb.text_frame, "قدم بعدی", size=16, bold=True, color=ORANGE,
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    hb = textbox(s, M, Inches(1.00), CONTENT_W, Inches(0.76))
    write(hb.text_frame, "چه کسی می‌خواهد نفر اول باشد؟", size=38, bold=True,
          color=WHITE, first=True, align=PP_ALIGN.CENTER, line=1.0)

    qb = textbox(s, Inches(1.7), Inches(1.98), SW - Inches(3.4), Inches(1.5))
    tf = qb.text_frame
    write(tf, "ما آدم‌های باتجربه‌ای را استخدام کردیم و کنار هم نشاندیم.",
          size=21, color=WHITE, first=True, align=PP_ALIGN.CENTER, line=1.3)
    write(tf, "فقط یادمان رفت فرصتی بسازیم که با هم حرف بزنند.",
          size=21, color=WHITE, align=PP_ALIGN.CENTER, line=1.3)
    write(tf, "این طرح، همان فرصت است.", size=25, bold=True, color=ORANGE,
          align=PP_ALIGN.CENTER, space_before=8, line=1.3)

    infos = [("اولین جلسه", cfg.first_session), ("ساعت", cfg.time),
             ("مکان", cfg.place), ("ظرفیت", "۱۲ نفر"), ("تماس", cfg.contact)]
    gap = Inches(0.22)
    cw = (CONTENT_W - 4 * gap) / 5
    y = Inches(3.90)
    ch = Inches(1.16)
    for i, (label, value) in enumerate(infos):
        x = SW - M - cw - i * (cw + gap)
        rect(s, x, y, cw, ch, fill=WHITE, alpha=0.14,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10, line=WHITE,
             line_w=Pt(1))
        tb = textbox(s, x + Inches(0.12), y, cw - Inches(0.24), ch,
                     anchor=MSO_ANCHOR.MIDDLE)
        write(tb.text_frame, label, size=12.5, bold=True, color=ORANGE,
              first=True, align=PP_ALIGN.CENTER, line=1.1)
        write(tb.text_frame, value, size=17, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, line=1.2)

    hline(s, (SW - Inches(3.0)) / 2, Inches(5.42), Inches(3.0), ORANGE, Pt(2.5))

    fb = textbox(s, M, Inches(5.66), CONTENT_W, Inches(1.1))
    tf = fb.text_frame
    write(tf, "برگه ثبت‌نام همین حالا دست به دست می‌شود.", size=23, bold=True,
          color=WHITE, first=True, align=PP_ALIGN.CENTER, line=1.25)
    write(tf, "بهترین زمان برای شروع، الان است.", size=20, color=ORANGE,
          align=PP_ALIGN.CENTER, line=1.25)

    # امضای تهیه‌کننده و آدرس سایت
    sig = textbox(s, M, Inches(6.86), CONTENT_W, Inches(0.40))
    write(sig.text_frame,
          [(cfg.presenter, True, WHITE, 14), ("     |     ", False, ORANGE, 14),
           ("www.coachroom.ir", True, ORANGE, 14)],
          first=True, align=PP_ALIGN.CENTER, line=1.0)

    notes(s, """
    نگویید «سؤالی نیست؟». بگویید «چه کسی می‌خواهد نفر اول باشد؟» و
    خودتان دستتان را بالا ببرید تا اولین داوطلب راحت‌تر بلند شود.
    برگه ثبت‌نام را همان لحظه دست به دست کنید — نه ایمیل، نه فردا.
    """)
    return s


# ═══════════════════════════════════════════════════════ اجرا ════════════
BUILDERS = [s01, s02, s03, s04, s05, s06, s07, s08, s09, s10,
            s11, s12, s13, s14, s15, s16, s17, s18, s19, s20]


def build(cfg):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    T.set_font(cfg.font)
    for fn in BUILDERS:
        fn(prs, cfg)
    prs.save(cfg.out)
    return cfg.out


def main():
    ap = argparse.ArgumentParser(description="ساخت ارائه جامع رویداد منتورینگ گروهی")
    ap.add_argument("--out", default=os.path.join(
        HERE, "رویداد-منتورینگ-گروهی.pptx"))
    ap.add_argument("--font", default="IRANSans")
    ap.add_argument("--date", default="[تاریخ]")
    ap.add_argument("--presenter", default="وحید مجیدی")
    ap.add_argument("--org", default="[نام واحد / سازمان]")
    ap.add_argument("--first-session", dest="first_session", default="[تاریخ]")
    ap.add_argument("--time", default="[ساعت]")
    ap.add_argument("--place", default="[مکان]")
    ap.add_argument("--contact", default="[داخلی]")
    cfg = ap.parse_args()
    print("ساخته شد:", build(cfg))


if __name__ == "__main__":
    main()
